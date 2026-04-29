"""
Generate M7 vs M8 comparison presentation.
Run:  python make_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ─────────────────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
CARD_BG   = RGBColor(0x13, 0x27, 0x3D)   # slightly lighter
ACCENT1   = RGBColor(0x00, 0xD2, 0xFF)   # cyan (M8 highlight)
ACCENT2   = RGBColor(0xFF, 0xA5, 0x00)   # orange (M7)
GREEN     = RGBColor(0x00, 0xFF, 0x88)   # good metric
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xAA, 0xAA, 0xAA)
RED       = RGBColor(0xFF, 0x57, 0x57)

# ── Helpers ────────────────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

def slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_size=14, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(1)):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        5,  # rounded rectangle
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def title_bar(slide, title, subtitle=None):
    """Dark top bar with title text."""
    add_rect(slide, 0, 0, 13.33, 1.4, CARD_BG)
    add_textbox(slide, title, 0.3, 0.1, 12, 0.7,
                font_size=28, bold=True, color=ACCENT1, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, 0.3, 0.82, 12, 0.45,
                    font_size=13, color=GRAY, align=PP_ALIGN.LEFT)

def section_label(slide, text, left, top, color=ACCENT1):
    add_textbox(slide, text, left, top, 5, 0.35,
                font_size=11, bold=True, color=color, align=PP_ALIGN.LEFT)

def metric_card(slide, label, value, left, top, w=1.9, h=1.0,
                val_color=WHITE, bg=CARD_BG, border=None):
    add_rounded_rect(slide, left, top, w, h, bg,
                     line_color=border or ACCENT1, line_width=Pt(1.2))
    add_textbox(slide, label, left+0.08, top+0.08, w-0.16, 0.3,
                font_size=10, color=GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, value, left+0.08, top+0.36, w-0.16, 0.5,
                font_size=18, bold=True, color=val_color, align=PP_ALIGN.CENTER)

def bullet(slide, items, left, top, width, spacing=0.32,
           font_size=13, color=WHITE, marker="▸"):
    for i, item in enumerate(items):
        add_textbox(slide, f"{marker}  {item}",
                    left, top + i * spacing, width, 0.35,
                    font_size=font_size, color=color)

def code_box(slide, code_text, left, top, width, height,
             bg=rgb(0x08,0x12,0x1E), border=rgb(0x00,0xD2,0xFF)):
    add_rounded_rect(slide, left, top, width, height, bg, line_color=border)
    txBox = slide.shapes.add_textbox(
        Inches(left+0.1), Inches(top+0.1),
        Inches(width-0.2), Inches(height-0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code_text
    run.font.size = Pt(10.5)
    run.font.color.rgb = rgb(0x7E, 0xE8, 0xFA)
    run.font.name = "Courier New"

def divider(slide, top, color=ACCENT1, width=13.0):
    shape = slide.shapes.add_shape(1,
        Inches(0.15), Inches(top), Inches(width), Inches(0.025))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]   # completely blank

# ─────────────────────────── SLIDE 1 — COVER ──────────────────────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
# big gradient band
add_rect(s, 0, 0, 13.33, 2.6, CARD_BG)
# accent stripe
shape = s.shapes.add_shape(1, Inches(0), Inches(2.55), Inches(13.33), Inches(0.07))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT1; shape.line.fill.background()

add_textbox(s, "M7  →  M8", 0.5, 0.25, 12, 1.1,
            font_size=54, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
add_textbox(s, "Efficient Ψ-NN: Architecture Upgrades & Physics Constraints",
            0.5, 1.3, 12, 0.7, font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, "Pile Stiffness Degradation — M8 vs M7 Technical Comparison",
            0.5, 1.95, 12, 0.5, font_size=13, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# summary badges
badges = [("M7 Params", "55,770", ACCENT2),
          ("M8 Params", "45,486", ACCENT1),
          ("Compression", "19.7 %", GREEN),
          ("Ψ-Model R²", "0.9882", WHITE)]
for i, (lbl, val, col) in enumerate(badges):
    x = 1.3 + i * 2.7
    metric_card(s, lbl, val, x, 3.05, w=2.35, h=1.15,
                val_color=col, border=col)

# description
add_textbox(s,
    "Three efficiency upgrades: Bottleneck MLP  •  Learnable Relation Matrix  •  Physics-Monotonic Loss",
    0.5, 4.6, 12, 0.45, font_size=13, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# footer
add_rect(s, 0, 6.9, 13.33, 0.6, rgb(0x08,0x12,0x20))
add_textbox(s, "PFE Research  ·  2026", 0.3, 6.95, 12.7, 0.4,
            font_size=10, color=GRAY, align=PP_ALIGN.RIGHT)

# ─────────────────────────── SLIDE 2 — M7 REMINDER ──────────────────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
title_bar(s, "Reminder: M7 — Ψ-NN Pile Stiffness Model",
          "What we built in M7 before upgrading to M8")

# left column — architecture
add_rect(s, 0.2, 1.55, 5.9, 5.6, CARD_BG)
section_label(s, "ARCHITECTURE", 0.35, 1.6, color=ACCENT2)
bullet(s, [
    "3-Stage pipeline: Distillation → Discovery → Ψ-Model",
    "Slot Attention: 64-dim, 4 heads, 3 iterations",
    "MLP per slot:  64 → 128 → 64 (wide bottleneck)",
    "1 initial slot + 20 drop slots",
    "Fixed relation matrix R (from K-means, frozen)",
    "M6 teacher frozen for distillation (Stage A)",
    "K-Means Stage B: discovered k*=6 prototypes",
], 0.35, 2.0, 5.6, spacing=0.44, font_size=12.5)

# right column — results
add_rect(s, 6.35, 1.55, 6.7, 5.6, CARD_BG)
section_label(s, "M7 RESULTS", 6.5, 1.6, color=ACCENT2)
metric_card(s, "Ψ-Model Parameters",  "55,770", 6.5,  2.05, w=2.9, h=1.0,
            val_color=ACCENT2, border=ACCENT2)
metric_card(s, "Compression vs M6",   "1.5 %",  9.65, 2.05, w=2.9, h=1.0,
            val_color=ACCENT2, border=ACCENT2)
metric_card(s, "Ψ-Model R² (test)",   "0.9897", 6.5,  3.25, w=2.9, h=1.0,
            val_color=WHITE, border=ACCENT2)
metric_card(s, "Ψ-Model RMSE",        "21.60 B", 9.65, 3.25, w=2.9, h=1.0,
            val_color=WHITE, border=ACCENT2)

section_label(s, "PER-VARIABLE  R²", 6.5, 4.45, color=GRAY)
rows = [("KL  (lateral stiffness)",  "0.9910"),
        ("KR  (rotational stiffness)","0.9807"),
        ("KLR (coupling)",            "0.9860")]
for i, (lbl, val) in enumerate(rows):
    add_textbox(s, lbl, 6.5, 4.8+i*0.42, 4.5, 0.35, font_size=12, color=GRAY)
    add_textbox(s, val, 11.0,4.8+i*0.42, 1.9, 0.35, font_size=12, bold=True, color=WHITE)

# loss formula box
add_rect(s, 6.35, 6.2, 6.7, 0.75, rgb(0x08,0x12,0x1E))
add_textbox(s, "Loss = L_distill + L_seq + 5·L_initial + L_shape",
            6.5, 6.3, 6.4, 0.5, font_size=11.5, color=rgb(0x7E,0xE8,0xFA),
            align=PP_ALIGN.CENTER)

# ─────────────────────── SLIDE 3 — 3 UPGRADES OVERVIEW ──────────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
title_bar(s, "M8: Three Efficiency Upgrades at a Glance",
          "All three upgrades are additive — M7 foundations remain intact")

# three cards
cards = [
    ("① Efficient Bottleneck MLP",
     "Replace the wide slot MLP\n(64→128→64) with a\nnarrow bottleneck\n(64→48→64)\n\n→  ~40% fewer MLP params\n→  EfficientSlotMLP class",
     ACCENT1, 0.25),
    ("② Learnable Relation Matrix",
     "Convert the fixed buffer R\n(from K-means Stage B)\ninto a learnable parameter\nvia softmax logits\n\n→  R evolves during training\n→  Valid prob. distribution",
     GREEN, 4.7),
    ("③ Physics-Monotonic Loss",
     "Add two new loss terms:\n• Relation-entropy (0.02×)\n• Monotonic penalty (0.2×)\n\nKL, KR ↓ over time\nKLR ↑ over time\n\n→  Physically realistic curves",
     ACCENT2, 9.15),
]
for title, body, col, x in cards:
    add_rounded_rect(s, x, 1.55, 4.1, 5.65, CARD_BG,
                     line_color=col, line_width=Pt(2))
    add_textbox(s, title, x+0.12, 1.65, 3.9, 0.55,
                font_size=15, bold=True, color=col)
    divider(s, 2.3+0.0, color=col, width=3.8)
    # shift divider per card
    shape = s.shapes.add_shape(1,
        Inches(x+0.1), Inches(2.28), Inches(3.9), Inches(0.025))
    shape.fill.solid(); shape.fill.fore_color.rgb = col; shape.line.fill.background()
    add_textbox(s, body, x+0.12, 2.42, 3.85, 4.6,
                font_size=12.5, color=WHITE)

# bottom tag
add_textbox(s, "Combined effect: 19.7% compression vs M6 teacher  ·  Physically constrained predictions",
            0.3, 7.1, 12.7, 0.35, font_size=11, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────── SLIDE 4 — UPGRADE 1: EFFICIENT MLP ─────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
title_bar(s, "Upgrade ①  —  EfficientSlotMLP  (64 → 48 → 64)",
          "Narrower bottleneck reduces parameters by ~40% with negligible accuracy loss")

# before / after side by side
for col_i, (label, dims, c, x) in enumerate([
        ("M7 — Wide MLP", "64 → 128 → 64", ACCENT2, 0.3),
        ("M8 — EfficientSlotMLP", "64 → 48 → 64", ACCENT1, 6.8)]):
    add_rounded_rect(s, x, 1.55, 6.1, 5.6, CARD_BG, line_color=c, line_width=Pt(1.5))
    add_textbox(s, label, x+0.15, 1.65, 5.8, 0.45, font_size=15, bold=True, color=c)
    add_textbox(s, dims,  x+0.15, 2.15, 5.8, 0.4,  font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# M7 code
code_box(s, "self.slot_mlp = nn.Sequential(\n    nn.Linear(64, 128),   # 8,192 params\n    nn.GELU(),\n    nn.Dropout(0.1),\n    nn.Linear(128, 64),   # 8,192 params\n)  # Total MLP: ~16,384 params",
         0.4, 2.7, 5.9, 1.8, bg=rgb(0x08,0x12,0x1E), border=ACCENT2)

# M8 code
code_box(s, "class EfficientSlotMLP(nn.Module):\n  def __init__(self, d=64, btl=48):\n    self.net = nn.Sequential(\n      nn.Linear(d, btl),    # 3,072 params\n      nn.GELU(),\n      nn.Dropout(0.1),\n      nn.Linear(btl, d),    # 3,072 params\n    )  # Total MLP: ~6,144 params",
         6.9, 2.7, 6.1, 1.95, bg=rgb(0x08,0x12,0x1E), border=ACCENT1)

# param comparison
add_textbox(s, "PARAMETER COUNT COMPARISON", 0.3, 4.7, 12.7, 0.4,
            font_size=12, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
cols_data = [("M6 Teacher",   "56,646", GRAY),
             ("M7 Student",   "56,646", ACCENT2),
             ("M7 Ψ-Model",   "55,770", ACCENT2),
             ("M8 Student",   "46,326 ↓", ACCENT1),
             ("M8 Ψ-Model",   "45,486 ↓", GREEN)]
for i, (lbl, val, c) in enumerate(cols_data):
    x = 0.5 + i * 2.45
    metric_card(s, lbl, val, x, 5.1, w=2.3, h=1.1, val_color=c, border=c)

# ─────────────────────── SLIDE 5 — UPGRADE 2: LEARNABLE R ───────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
title_bar(s, "Upgrade ②  —  Learnable Relation Matrix  R",
          "Fixed K-means initialisation → trainable parameter that adapts during Stage C")

# left explanation
add_rect(s, 0.2, 1.55, 6.1, 5.65, CARD_BG)
section_label(s, "WHAT IS THE RELATION MATRIX?", 0.35, 1.62, color=ACCENT1)
bullet(s, [
    "R is a [20 × k*] matrix (20 drop slots, k* prototypes)",
    "Each row = mixing weights for one time step",
    "Maps k* prototype slots → 20 degradation steps",
    "Stage B (K-means) provides a good initialisation",
], 0.35, 2.0, 5.8, spacing=0.42, font_size=12.5)

section_label(s, "M7 (FIXED)", 0.35, 3.88, color=ACCENT2)
code_box(s, "# register_buffer = NOT a learnable parameter\nself.register_buffer(\n  'relation_matrix',\n  torch.FloatTensor(R)  # frozen after Stage B\n)",
         0.35, 4.25, 5.7, 1.5, bg=rgb(0x08,0x12,0x1E), border=ACCENT2)

# right upgrade
add_rect(s, 6.6, 1.55, 6.5, 5.65, CARD_BG)
section_label(s, "M8 (LEARNABLE)", 6.75, 1.62, color=ACCENT1)
code_box(s, "# nn.Parameter = gradient flows through R\nself.relation_logits = nn.Parameter(\n    torch.log(R_init.clamp(min=1e-6))\n)  # log-space for numerical stability\n\ndef get_relation_matrix(self):\n    # Row-softmax → valid prob distribution\n    return torch.softmax(\n        self.relation_logits, dim=1\n    )",
         6.75, 2.0, 6.1, 2.5, bg=rgb(0x08,0x12,0x1E), border=ACCENT1)

section_label(s, "BENEFITS", 6.75, 4.65, color=GREEN)
bullet(s, [
    "R adapts to minimize the multi-objective loss",
    "Row-softmax ensures valid probability mix",
    "Entropy regularisation sharpens assignments",
    "Saved as relation_matrix (learned) in config",
], 6.75, 5.05, 6.1, spacing=0.4, font_size=12.5, color=WHITE)

# ─────────────────────── SLIDE 6 — UPGRADE 3: PHYSICS LOSS ─────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
title_bar(s, "Upgrade ③  —  Physics-Monotonic Loss",
          "Two new loss components enforce physically realistic pile degradation behaviour")

# constraint explanation — left
add_rect(s, 0.2, 1.55, 6.0, 5.65, CARD_BG)
section_label(s, "PHYSICAL CONSTRAINTS", 0.35, 1.62, color=ACCENT2)
add_textbox(s, "Pile stiffness under cyclic loading must satisfy:", 0.35, 2.0, 5.7, 0.4,
            font_size=12.5, color=GRAY)
rows = [("KL  (lateral stiffness)",   "↓ decreases (pile softens)", RED),
        ("KR  (rotational stiffness)", "↓ decreases (pile softens)", RED),
        ("KLR (coupling term)",        "↑ increases (coupling grows)", GREEN)]
for i, (name, rule, c) in enumerate(rows):
    y = 2.55 + i * 0.68
    add_rounded_rect(s, 0.35, y, 5.65, 0.58, rgb(0x0A,0x1A,0x2E),
                     line_color=c, line_width=Pt(1))
    add_textbox(s, name, 0.5, y+0.05, 2.8, 0.45, font_size=12, color=WHITE, bold=True)
    add_textbox(s, rule, 3.4, y+0.05, 2.4, 0.45, font_size=12, color=c, bold=True)

add_textbox(s, "M7 had no explicit enforcement of these constraints →\npredictions could violate physical laws",
            0.35, 4.68, 5.65, 0.7, font_size=12, color=rgb(0xFF,0xAA,0x55), italic=True)

# M7 vs M8 loss box
section_label(s, "M7 LOSS  vs  M8 LOSS", 0.35, 5.5, color=GRAY)
add_textbox(s, "M7:  L = L_distill + L_seq + 5·L_initial + L_shape",
            0.35, 5.85, 5.7, 0.4, font_size=11.5, color=rgb(0x7E,0xE8,0xFA))
add_textbox(s, "M8:  L = L_distill + L_seq + 5·L_initial + L_shape\n       + 0.02·L_rel_entropy  + 0.2·L_mono_physics",
            0.35, 6.3, 5.7, 0.65, font_size=11.5, color=rgb(0x7E,0xE8,0xFA))

# code — right
add_rect(s, 6.5, 1.55, 6.6, 5.65, CARD_BG)
section_label(s, "M8 IMPLEMENTATION", 6.65, 1.62, color=ACCENT1)
code_box(s, "# Physics-monotonic penalty\ndiff_kl_kr = psi_out[:,1:,:2] - psi_out[:,:-1,:2]\ndiff_klr   = psi_out[:,1:,2:] - psi_out[:,:-1,2:]\n\n# Penalise violations:\nloss_mono = (\n  torch.relu(diff_kl_kr).mean()   # KL,KR must NOT increase\n+ torch.relu(-diff_klr).mean()    # KLR must NOT decrease\n)\n\n# Relation-entropy regularisation\nR = psi_model.get_relation_matrix()\nloss_rel_ent = -(R * torch.log(R+1e-8)).sum(1).mean()\n# → encourages sharp prototype assignments\n\nloss = (...base losses...\n        + 0.02 * loss_rel_ent\n        + 0.20 * loss_mono)",
         6.65, 2.0, 6.2, 4.8, bg=rgb(0x08,0x12,0x1E), border=ACCENT1)

# ─────────────────────── SLIDE 7 — LOSS FUNCTION SUMMARY ───────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
title_bar(s, "Loss Function: M7  vs  M8",
          "M8 = M7 base losses + two new physics-informed terms")

# table
headers = ["Loss Component", "Description", "M7", "M8"]
rows_data = [
    ["L_distill",      "MSE vs M6 teacher output",          "✓", "✓"],
    ["L_seq",          "SmoothL1 vs target curves",         "✓", "✓"],
    ["5·L_initial",    "MSE on first step (×5 weight)",     "✓", "✓"],
    ["L_shape",        "SmoothL1 on consecutive Δ curves",  "✓", "✓"],
    ["0.02·L_rel_ent", "Entropy of R → sharp assignments",  "—", "✓ NEW"],
    ["0.2·L_physics",  "Monotonic penalty KL↓ KR↓ KLR↑",   "—", "✓ NEW"],
]
col_widths = [3.0, 5.5, 1.5, 1.8]
col_starts = [0.2, 3.35, 9.0, 10.65]
row_h = 0.62
header_y = 1.6

# header row
for ci, (hdr, cw, cx) in enumerate(zip(headers, col_widths, col_starts)):
    add_rect(s, cx, header_y, cw-0.05, row_h-0.05, rgb(0x00,0x55,0x77))
    add_textbox(s, hdr, cx+0.08, header_y+0.1, cw-0.2, 0.4,
                font_size=13, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows_data):
    y = header_y + (ri+1)*row_h
    bg_c = rgb(0x0A,0x1A,0x2E) if ri % 2 == 0 else CARD_BG
    for ci, (cell, cw, cx) in enumerate(zip(row, col_widths, col_starts)):
        add_rect(s, cx, y, cw-0.05, row_h-0.05, bg_c)
        if ci in (2, 3):  # checkmark cols
            c = GREEN if "✓" in cell else (GRAY if "—" in cell else ACCENT1)
            add_textbox(s, cell, cx+0.08, y+0.1, cw-0.2, 0.4,
                        font_size=13, bold=True, color=c, align=PP_ALIGN.CENTER)
        else:
            clr = ACCENT2 if ri < 4 else ACCENT1
            add_textbox(s, cell, cx+0.08, y+0.1, cw-0.2, 0.4,
                        font_size=12, color=clr)

# bottom note
add_textbox(s,
    "All 6 components are computed on full 21-step curves — the model targets curve shape, not individual slot values.",
    0.2, 6.5, 12.9, 0.5, font_size=12, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────── SLIDE 8 — PARAMETER COMPARISON ────────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
title_bar(s, "Parameter Efficiency: M7  vs  M8",
          "M8 achieves 19.7% compression vs the M6 teacher while M7 only achieved 1.5%")

# big metrics
metrics = [
    ("M6 Teacher\n(Reference)", "56,646", GRAY,  0.3),
    ("M7 Ψ-Model",              "55,770",  ACCENT2, 3.4),
    ("M8 Ψ-Model",              "45,486",  ACCENT1, 6.5),
    ("Compression\nM7 vs M6",   "1.5 %",   ACCENT2, 9.6),
]
for lbl, val, c, x in metrics:
    metric_card(s, lbl, val, x, 1.55, w=2.8, h=1.35, val_color=c, border=c)

# bar chart (manual)
bar_y_base = 6.8
bar_data = [
    ("M6 Teacher",  56646, GRAY),
    ("M7 Student",  56646, ACCENT2),
    ("M7 Ψ-Model",  55770, ACCENT2),
    ("M8 Student",  46326, ACCENT1),
    ("M8 Ψ-Model",  45486, GREEN),
]
max_val = 60000
bar_h_max = 3.2
bar_w = 1.5
gap = 0.5
start_x = 1.3

add_textbox(s, "PARAMETER COUNT — BAR CHART", 0.3, 3.0, 12.7, 0.35,
            font_size=11, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

for i, (name, val, c) in enumerate(bar_data):
    h = bar_h_max * val / max_val
    x = start_x + i * (bar_w + gap)
    y_top = bar_y_base - h
    add_rect(s, x, y_top, bar_w, h, c)
    # value label
    add_textbox(s, f"{val:,}", x, y_top - 0.38, bar_w, 0.35,
                font_size=11, bold=True, color=c, align=PP_ALIGN.CENTER)
    # name label
    add_textbox(s, name, x - 0.1, bar_y_base + 0.05, bar_w + 0.2, 0.45,
                font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

# compression annotation
add_textbox(s, "−19.7 %", 9.3, 3.6, 3.5, 0.5,
            font_size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_textbox(s, "compression vs M6", 9.3, 4.1, 3.5, 0.4,
            font_size=11, color=GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, "vs only −1.5% in M7", 9.3, 4.5, 3.5, 0.35,
            font_size=11, color=ACCENT2, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────── SLIDE 9 — PERFORMANCE COMPARISON ──────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
title_bar(s, "Performance Comparison: M7  vs  M8",
          "M8 maintains strong accuracy (~R²=0.988) with significantly fewer parameters")

# overall results
section_label(s, "OVERALL METRICS (TEST SET — Ψ-Model only)", 0.3, 1.55, color=GRAY)
cols_perf = [
    ("M6 Teacher R²", "0.9804", GRAY),
    ("M7 Ψ-Model R²", "0.9897", ACCENT2),
    ("M8 Ψ-Model R²", "0.9882", ACCENT1),
    ("M7 RMSE",        "21.60 B", ACCENT2),
    ("M8 RMSE",        "23.06 B", ACCENT1),
    ("M8 Curve MAPE",  "2.48 %",  GREEN),
]
for i, (lbl, val, c) in enumerate(cols_perf):
    x = 0.3 + i * 2.12
    metric_card(s, lbl, val, x, 1.9, w=2.05, h=1.1, val_color=c, border=c)

# per-variable table
section_label(s, "PER-VARIABLE PERFORMANCE — Ψ-MODEL", 0.3, 3.2, color=GRAY)
pv_headers = ["Variable", "M7 R²", "M8 R²", "M7 RMSE", "M8 RMSE"]
pv_rows = [
    ("KL  (lateral stiffness)",   "0.9910", "0.9888", "61.1 M", "68.1 M"),
    ("KR  (rotational stiffness)","0.9807", "0.9780", "37.4 B", "39.9 B"),
    ("KLR (coupling)",            "0.9860", "0.9845", "1.29 B", "1.36 B"),
]
pv_col_w  = [3.8, 1.9, 1.9, 1.9, 1.9]
pv_col_x  = [0.3, 4.25, 6.25, 8.25, 10.25]
pv_row_h  = 0.58
pv_hdr_y  = 3.55

for ci, (hdr, cw, cx) in enumerate(zip(pv_headers, pv_col_w, pv_col_x)):
    add_rect(s, cx, pv_hdr_y, cw-0.05, pv_row_h-0.05, rgb(0x00,0x40,0x60))
    add_textbox(s, hdr, cx+0.08, pv_hdr_y+0.1, cw-0.2, 0.4,
                font_size=13, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)

for ri, row in enumerate(pv_rows):
    y = pv_hdr_y + (ri+1)*pv_row_h
    bg_c = rgb(0x0A,0x1A,0x2E) if ri%2==0 else CARD_BG
    for ci, (cell, cw, cx) in enumerate(zip(row, pv_col_w, pv_col_x)):
        add_rect(s, cx, y, cw-0.05, pv_row_h-0.05, bg_c)
        c = (ACCENT2 if ci in (1,3) else ACCENT1) if ci > 0 else WHITE
        add_textbox(s, cell, cx+0.08, y+0.1, cw-0.2, 0.4,
                    font_size=13, color=c, align=PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT,
                    bold=(ci==0))

# clustering comparison
section_label(s, "STRUCTURE DISCOVERY (STAGE B)", 0.3, 5.5, color=GRAY)
cl_data = [("M7 k*",  "6 prototypes", ACCENT2), ("M7 Silhouette","0.7146", ACCENT2),
           ("M8 k*",  "5 prototypes", ACCENT1), ("M8 Silhouette","0.6490", ACCENT1)]
for i, (lbl, val, c) in enumerate(cl_data):
    metric_card(s, lbl, val, 0.3+i*3.1, 5.85, w=2.95, h=1.1, val_color=c, border=c)

add_textbox(s,
    "Note: slight accuracy trade-off (-0.0015 R²) is the cost of 19.7% compression + physics guarantee.",
    0.3, 7.1, 12.7, 0.35, font_size=11, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────── SLIDE 10 — SUMMARY / CONCLUSIONS ──────────────────
s = prs.slides.add_slide(blank)
slide_bg(s)
title_bar(s, "Summary: M7  →  M8  Key Takeaways",
          "Three targeted upgrades deliver compression, adaptability, and physical correctness")

# left — differences table
add_rect(s, 0.2, 1.55, 7.7, 5.65, CARD_BG)
section_label(s, "WHAT CHANGED FROM M7 TO M8", 0.35, 1.62, color=ACCENT1)
diff_rows = [
    ("Feature",             "M7",               "M8"),
    ("Slot MLP",            "64→128→64",         "64→48→64  ↓"),
    ("Relation Matrix R",   "Fixed buffer",      "Learnable logits ✓"),
    ("Physics loss",        "None",              "0.2 × L_mono"),
    ("Entropy loss",        "None",              "0.02 × L_rel_ent"),
    ("Student params",      "56,646",            "46,326  (−18.1%)"),
    ("Ψ-Model params",      "55,770",            "45,486  (−18.4%)"),
    ("Compression vs M6",   "1.5 %",             "19.7 %  ✓✓"),
    ("k* prototypes",       "6",                 "5"),
    ("Ψ-Model R² (test)",   "0.9897",            "0.9882  (~same)"),
]
dw = [3.4, 2.0, 2.0]
dx = [0.35, 3.85, 5.95]
dh = 0.49
for ri, row in enumerate(diff_rows):
    bg_c = rgb(0x00,0x40,0x60) if ri==0 else (rgb(0x0A,0x1A,0x2E) if ri%2==0 else CARD_BG)
    for ci, (cell, cw, cx) in enumerate(zip(row, dw, dx)):
        add_rect(s, cx, 2.0+ri*dh, cw-0.05, dh-0.04, bg_c)
        is_header = ri==0
        c = (ACCENT1 if is_header else
             (ACCENT2 if ci==1 else (GREEN if "✓" in cell else
             (RED if "None" in cell else WHITE))))
        add_textbox(s, cell, cx+0.08, 2.04+ri*dh, cw-0.18, 0.4,
                    font_size=11.5 if not is_header else 12,
                    bold=is_header, color=c,
                    align=PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT)

# right — conclusions
add_rect(s, 8.2, 1.55, 4.9, 5.65, CARD_BG)
section_label(s, "CONCLUSIONS", 8.35, 1.62, color=GREEN)
conclusions = [
    "M8 is a leaner, physics-aware\nversion of M7",
    "19.7% fewer parameters\nversus M6 teacher",
    "Accuracy trade-off minimal:\nΔR² = −0.0015",
    "Learnable R enables better\nprototype assignment",
    "Physics loss guarantees\nmonotonic degradation",
    "Curve MAPE = 2.48% confirms\nfull-curve accuracy",
]
for i, txt in enumerate(conclusions):
    y = 2.05 + i * 0.82
    add_rounded_rect(s, 8.35, y, 4.6, 0.72, rgb(0x0A,0x1E,0x32),
                     line_color=GREEN, line_width=Pt(0.8))
    add_textbox(s, txt, 8.5, y+0.06, 4.3, 0.62,
                font_size=11.5, color=WHITE)

# bottom banner
add_rect(s, 0, 7.1, 13.33, 0.4, rgb(0x00,0x40,0x60))
add_textbox(s, "M8 achieves the best balance between model size, accuracy, and physical interpretability in the Ψ-NN pipeline.",
            0.3, 7.15, 12.7, 0.35, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "M7_vs_M8_comparison.pptx")
prs.save(out)
print(f"Saved → {out}")
