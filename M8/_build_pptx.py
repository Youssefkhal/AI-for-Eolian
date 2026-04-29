"""
Build new.pptx — M6 architecture presentation
Run: python M6/_build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree
import copy
from datetime import datetime

# ── colour palette ──────────────────────────────────────────────────────────
BLUE    = RGBColor(0x1A, 0x73, 0xE8)   # header blue
DKBLUE  = RGBColor(0x0D, 0x47, 0xA1)   # darker blue
GREEN   = RGBColor(0x2E, 0xCC, 0x71)   # slot / attention green
PURPLE  = RGBColor(0x8E, 0x44, 0xAD)   # embedding purple
ORANGE  = RGBColor(0xE6, 0x7E, 0x22)   # physics orange
RED     = RGBColor(0xE7, 0x4C, 0x3C)   # output red
YELLOW  = RGBColor(0xF3, 0x9C, 0x12)   # LSTM / drop yellow
TEAL    = RGBColor(0x1A, 0xBC, 0x9C)   # cross-attn teal
LGREY   = RGBColor(0xEC, 0xF0, 0xF1)   # light background
DGREY   = RGBColor(0x2C, 0x3E, 0x50)   # dark text
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# ── helpers ──────────────────────────────────────────────────────────────────

def add_slide(prs, layout_idx=6):
    """Blank layout slide."""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(1.5), radius=None):
    """Add a rounded-rect shape, return shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    sf = shape.fill
    if fill:
        sf.solid()
        sf.fore_color.rgb = fill
    else:
        sf.background()
    sl = shape.line
    if line:
        sl.color.rgb = line
        sl.width = line_w
    else:
        sl.fill.background()
    if radius:
        # set rounded corners via XML
        sp = shape._element
        spPr = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
        if spPr is not None:
            prstGeom = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
            if prstGeom is not None:
                prstGeom.set('prst', 'roundRect')
                avLst = prstGeom.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                if avLst is None:
                    avLst = etree.SubElement(prstGeom, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                gd = etree.SubElement(avLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
                gd.set('name', 'adj')
                gd.set('fmla', f'val {radius}')
    return shape


def txbox(slide, text, l, t, w, h,
          size=Pt(11), bold=False, italic=False,
          color=DGREY, align=PP_ALIGN.LEFT,
          wrap=True):
    """Add a plain text box."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def header_bar(slide, title, subtitle=None):
    """Dark blue header bar across the top."""
    rect(slide, 0, 0, 13.33, 1.1, fill=DKBLUE)
    txbox(slide, title, 0.25, 0.08, 12, 0.55,
          size=Pt(24), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txbox(slide, subtitle, 0.25, 0.62, 12, 0.38,
              size=Pt(13), italic=True, color=RGBColor(0xBB, 0xDE, 0xFB),
              align=PP_ALIGN.LEFT)


def block(slide, text, l, t, w, h,
          fill=LGREY, line=BLUE, txt_size=Pt(10), bold=False,
          txt_color=DGREY, align=PP_ALIGN.CENTER):
    """Filled rectangle with centred label."""
    r = rect(slide, l, t, w, h, fill=fill, line=line, radius=15000)
    tf = r.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = txt_size
    run.font.bold = bold
    run.font.color.rgb = txt_color
    return r


def arrow(slide, x1, y1, x2, y2, color=DGREY, width=Pt(1.5)):
    """Simple connector arrow."""
    from pptx.util import Pt
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = width
    return conn


def bullet_box(slide, items, l, t, w, h,
               header=None, header_color=BLUE,
               item_size=Pt(11), item_color=DGREY,
               fill=LGREY, line=BLUE):
    """Box with optional header + bullet list."""
    rect(slide, l, t, w, h, fill=fill, line=line, radius=10000)
    tb = slide.shapes.add_textbox(Inches(l+0.12), Inches(t+0.1), Inches(w-0.24), Inches(h-0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    if header:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = header
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = header_color
    for item in items:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"  \u2022  {item}"
        run.font.size = item_size
        run.font.color.rgb = item_color
    return tb


def table_slide(slide, headers, rows,
                l, t, w, h,
                header_fill=DKBLUE, row_fills=None):
    """Add a table. headers = list of str, rows = list of list of str."""
    cols = len(headers)
    n_rows = len(rows) + 1
    col_w = w / cols
    row_h = h / n_rows

    # header
    for ci, hdr in enumerate(headers):
        r = rect(slide, l + ci*col_w, t, col_w, row_h, fill=header_fill)
        txbox(slide, hdr, l + ci*col_w + 0.05, t + 0.04,
              col_w - 0.1, row_h - 0.08,
              size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        fill = row_fills[ri] if row_fills else (LGREY if ri % 2 == 0 else WHITE)
        for ci, cell in enumerate(row):
            rect(slide, l + ci*col_w, t + (ri+1)*row_h, col_w, row_h,
                 fill=fill, line=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(0.5))
            txbox(slide, cell,
                  l + ci*col_w + 0.05, t + (ri+1)*row_h + 0.03,
                  col_w - 0.1, row_h - 0.06,
                  size=Pt(9.5), color=DGREY, align=PP_ALIGN.CENTER)


# ============================================================
# BUILD PRESENTATION
# ============================================================

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── SLIDE 1 — Title ─────────────────────────────────────────────────────────
sl = add_slide(prs)
rect(sl, 0, 0, 13.33, 7.5, fill=DKBLUE)
rect(sl, 0, 2.8, 13.33, 2.0, fill=BLUE)

txbox(sl, "M6 — Pile Stiffness Degradation Prediction",
      0.5, 1.0, 12.3, 1.4,
      size=Pt(34), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "Slot-Attention Transformer · Physics-Informed · LSTM-Free",
      0.5, 2.5, 12.3, 0.8,
      size=Pt(18), italic=True, color=RGBColor(0xBB, 0xDE, 0xFB),
      align=PP_ALIGN.CENTER)
txbox(sl, "Architecture  ·  Fine-tuning Strategy  ·  MLP vs LSTM  ·  Efficiency  ·  Error Metrics",
      0.5, 3.25, 12.3, 0.55,
      size=Pt(13), color=WHITE, align=PP_ALIGN.CENTER)

for i, label in enumerate(["Predict KL, KR, KLR over 21 steps",
                            "64.1% fewer parameters than M5",
                            "Hard-coded monotonicity guarantees"]):
    txbox(sl, f"✓  {label}",
          1.2 + i*3.8, 4.2, 3.5, 0.45,
          size=Pt(12), color=RGBColor(0xA8, 0xD8, 0xA8), align=PP_ALIGN.LEFT)

txbox(sl, "PFE Project  |  April 2026",
      0.5, 6.9, 12.3, 0.4,
      size=Pt(10), italic=True, color=RGBColor(0x90, 0xCA, 0xF9),
      align=PP_ALIGN.CENTER)


# ── SLIDE 2 — Problem Statement & Overview ──────────────────────────────────
sl = add_slide(prs)
header_bar(sl, "Problem & Goal", "Predicting pile-soil stiffness degradation under cyclic loading")

txbox(sl, "Given 8 geotechnical input features, predict 3 stiffness quantities over 21 cyclic loading steps",
      0.3, 1.2, 12.7, 0.45, size=Pt(13), color=DGREY)

# Input box
block(sl, "INPUT  (8 features)\nPI, Gmax, ν, Dp, Tp, Lp, Ip, Dp/Lp",
      0.3, 1.75, 4.0, 1.25, fill=RGBColor(0xD6, 0xEA, 0xF8), line=BLUE,
      txt_size=Pt(11), txt_color=DKBLUE)

arrow(sl, 4.3, 2.38, 5.0, 2.38)

# Model box
block(sl, "M6 MODEL\n(Slot-Attention Transformer)",
      5.0, 1.75, 3.3, 1.25, fill=RGBColor(0xD5, 0xF5, 0xE3), line=GREEN,
      txt_size=Pt(11), txt_color=RGBColor(0x1A, 0x7A, 0x3C))

arrow(sl, 8.3, 2.38, 9.0, 2.38)

# Output box
block(sl, "OUTPUT  (21 × 3)\nKL(t), KR(t), KLR(t)\nfor t = 1 … 21",
      9.0, 1.75, 4.0, 1.25, fill=RGBColor(0xFD, 0xBD, 0xBD), line=RED,
      txt_size=Pt(11), txt_color=RGBColor(0x8B, 0x00, 0x00))

# Key features row
items = [
    ("21 Learnable Slots", GREEN, "One slot per degradation step — temporal structure without RNN"),
    ("Physics Constraints", ORANGE, "ΔKL, ΔKR ≤ 0  |  ΔKLR ≥ 0  (hard-coded monotonicity)"),
    ("Cumulative Sum", TEAL, "K(t) = K⁰ + Σ Δₙ  — error never resets"),
    ("Signed Log Scale", PURPLE, "Handles mixed-sign targets (KLR < 0, KL/KR > 0)"),
]
for i, (title, col, desc) in enumerate(items):
    x = 0.3 + i * 3.25
    rect(sl, x, 3.25, 3.0, 0.38, fill=col, radius=12000)
    txbox(sl, title, x+0.05, 3.27, 2.9, 0.34,
          size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, desc, x, 3.68, 3.05, 0.55,
          size=Pt(9), color=DGREY, align=PP_ALIGN.LEFT)

# Stiffness meaning
txbox(sl, "Stiffness quantities:  KL = lateral  |  KR = rotational  |  KLR = cross-coupling (always < 0, increases toward 0 under loading)",
      0.3, 4.35, 12.7, 0.38, size=Pt(10), italic=True, color=RGBColor(0x55, 0x55, 0x55))

# Data pipeline
bullet_box(sl,
    ["44 raw time steps → subsampled to 21 (equal spacing)",
     "Signed log transform: sign(y)·log(1+|y|)  →  preserves negative KLR",
     "RobustScaler per output variable (median/IQR, robust to outliers)",
     "Additive drop model: network learns Δ, not absolute K — physically motivated"],
    0.3, 4.8, 12.7, 2.45,
    header="Data Pipeline", header_color=BLUE,
    fill=RGBColor(0xF4, 0xF6, 0xF7), line=RGBColor(0xAA, 0xBB, 0xCC))


# ── SLIDE 3 — Full Architecture Diagram ─────────────────────────────────────
sl = add_slide(prs)
header_bar(sl, "Architecture Overview", "M6 — SlotAttentionDegradation  (LSTM-free, physics-constrained)")

# Draw the architecture as blocks + arrows
# Col centres: input=0.7, embed=2.5, slots=5.0, refine=8.5, split=11/12
Y0 = 1.3   # top row y

# INPUT
block(sl, "INPUT\n[B, 8]", 0.15, Y0, 1.4, 0.85,
      fill=RGBColor(0xD6, 0xEA, 0xF8), line=BLUE,
      txt_size=Pt(9), txt_color=DKBLUE)
txbox(sl, "8 features", 0.15, Y0+0.87, 1.4, 0.28, size=Pt(8), italic=True,
      color=DGREY, align=PP_ALIGN.CENTER)

arrow(sl, 1.55, Y0+0.42, 2.1, Y0+0.42, color=BLUE)

# EMBEDDING
block(sl, "EMBEDDING\nLinear 8→64\nLayerNorm\nGELU\n[B,1,64]",
      2.1, Y0, 1.85, 1.25,
      fill=RGBColor(0xE8, 0xD5, 0xF5), line=PURPLE,
      txt_size=Pt(9), txt_color=PURPLE)

arrow(sl, 3.95, Y0+0.62, 4.5, Y0+0.62, color=PURPLE)

# SLOTS INIT
block(sl, "21 LEARNABLE\nSLOTS\n[B, 21, 64]",
      4.5, Y0, 2.1, 1.1,
      fill=RGBColor(0xD5, 0xF5, 0xE3), line=GREEN,
      txt_size=Pt(9), txt_color=RGBColor(0x1A, 0x7A, 0x3C))
txbox(sl, "Slot 1 (initial) + Slots 2–21 (drops)", 4.5, Y0+1.12, 2.1, 0.3,
      size=Pt(7.5), italic=True, color=DGREY, align=PP_ALIGN.CENTER)

arrow(sl, 6.6, Y0+0.55, 7.05, Y0+0.55, color=GREEN)

# REFINEMENT BOX
rect(sl, 7.05, Y0-0.05, 4.55, 2.7,
     fill=RGBColor(0xF0, 0xFB, 0xF4), line=GREEN, line_w=Pt(2), radius=15000)
txbox(sl, "ITERATIVE SLOT REFINEMENT  ×3",
      7.1, Y0-0.02, 4.4, 0.3,
      size=Pt(9.5), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# inside refinement
block(sl, "Cross-Attention\nQ=slots\nK=V=input",
      7.15, Y0+0.32, 1.28, 0.95,
      fill=RGBColor(0xD1, 0xF2, 0xEB), line=TEAL, txt_size=Pt(8))
txbox(sl, "Add & LN", 7.15, Y0+1.3, 1.28, 0.25,
      size=Pt(7.5), italic=True, color=DGREY, align=PP_ALIGN.CENTER)

block(sl, "Self-Attention\nQ=K=V\n=slots",
      8.6, Y0+0.32, 1.28, 0.95,
      fill=RGBColor(0xD5, 0xF5, 0xE3), line=GREEN, txt_size=Pt(8))
txbox(sl, "Add & LN", 8.6, Y0+1.3, 1.28, 0.25,
      size=Pt(7.5), italic=True, color=DGREY, align=PP_ALIGN.CENTER)

block(sl, "Slot MLP\n64→128→64",
      10.07, Y0+0.32, 1.42, 0.95,
      fill=RGBColor(0xE8, 0xD5, 0xF5), line=PURPLE, txt_size=Pt(8))
txbox(sl, "Add & LN", 10.07, Y0+1.3, 1.42, 0.25,
      size=Pt(7.5), italic=True, color=DGREY, align=PP_ALIGN.CENTER)

arrow(sl, 8.43, Y0+0.82, 8.6, Y0+0.82, color=TEAL)
arrow(sl, 9.88, Y0+0.82, 10.07, Y0+0.82, color=GREEN)

# Down arrow from refinement
arrow(sl, 9.28, Y0+2.65, 9.28, Y0+3.15, color=GREEN)

# REFINED SLOTS
block(sl, "Refined Slots  [B, 21, 64]",
      7.35, Y0+3.15, 3.85, 0.65,
      fill=RGBColor(0xD5, 0xF5, 0xE3), line=GREEN, txt_size=Pt(10))

# SPLIT
arrow(sl, 8.5, Y0+3.8, 6.5, Y0+4.5, color=RED)
arrow(sl, 10.1, Y0+3.8, 11.0, Y0+4.5, color=YELLOW)

# Slot 1
block(sl, "Slot 1\n(initial)\n[B,1,64]",
      5.4, Y0+4.5, 2.1, 0.9,
      fill=RGBColor(0xFD, 0xBD, 0xBD), line=RED, txt_size=Pt(9))
arrow(sl, 6.45, Y0+5.4, 6.45, Y0+5.85, color=RED)
block(sl, "Initial MLP\n64→32→3  →  K⁰",
      5.4, Y0+5.85, 2.1, 0.75,
      fill=RGBColor(0xFD, 0xBD, 0xBD), line=RED, txt_size=Pt(8.5))
txbox(sl, "[B, 1, 3]", 5.4, Y0+6.62, 2.1, 0.25,
      size=Pt(8), italic=True, color=DGREY, align=PP_ALIGN.CENTER)

# Slots 2-21
block(sl, "Slots 2–21\n[B,20,64]",
      9.7, Y0+4.5, 2.3, 0.9,
      fill=RGBColor(0xFF, 0xF3, 0xCD), line=YELLOW, txt_size=Pt(9))
arrow(sl, 10.85, Y0+5.4, 10.85, Y0+5.85, color=YELLOW)
block(sl, "Drop MLP (shared)\n64→32→3  → Δ drops",
      9.7, Y0+5.85, 2.3, 0.75,
      fill=RGBColor(0xFF, 0xF3, 0xCD), line=YELLOW, txt_size=Pt(8.5))
arrow(sl, 10.85, Y0+6.6, 10.85, Y0+6.95, color=ORANGE)
block(sl, "Physics Constraint\n−|Δ| for KL,KR  ;  +|Δ| for KLR",
      9.3, Y0+6.95, 3.0, 0.7,
      fill=RGBColor(0xFD, 0xEC, 0xD9), line=ORANGE, txt_size=Pt(8.5))


# ── SLIDE 4 — Fine-Tuning / Training Strategy ───────────────────────────────
sl = add_slide(prs)
header_bar(sl, "Fine-Tuning Strategy", "How M6 was trained from scratch and what was changed from M5")

txbox(sl, "M6 is not a pretrained model fine-tuned from M5 — it is trained from scratch, but the training strategy was carefully redesigned based on M5's failure modes.",
      0.3, 1.18, 12.7, 0.48, size=Pt(11), italic=True, color=DGREY)

# 4 strategy cards
cards = [
    (BLUE, "① Initial Loss Weighting  ×5",
     ["MSE on Slot 1 (initial K⁰) multiplied by 5.0",
      "Any error in K⁰ propagates forward via cumulative sum",
      "Forces the model to anchor the trajectory correctly first",
      "M5 used ×3 — M6 raises to ×5 (LSTM no longer corrects drift)"]),
    (ORANGE, "② Shape Loss (Huber on Δ)",
     ["Penalises changes in consecutive drops: Huber(Δₜ − Δₜ₋₁)",
      "Encourages smooth degradation curves, not zigzag predictions",
      "Huber norm: robust to large isolated errors (vs pure L2)"]),
    (GREEN, "③ Gradient Clipping & Scheduler",
     ["Gradient clip = 2.0 (M5: 1.0) — relaxed to let early gradients flow",
      "ReduceLROnPlateau: patience 80 (M5: 100) — faster adaptation",
      "Early stopping patience 200 epochs (M5: 300)"]),
    (PURPLE, "④ Weight Decay & Normalisation",
     ["Weight decay = 0.005 (M5: 0.01) — reduced to prevent over-regularisation",
      "RobustScaler on targets: median/IQR instead of mean/std",
      "Signed log transform preserves sign of cross-coupling term KLR"]),
]

for i, (color, title, bullets) in enumerate(cards):
    x = 0.25 + (i % 2) * 6.5
    y = 1.8 + (i // 2) * 2.7
    rect(sl, x, y, 6.2, 2.5, fill=RGBColor(0xF8, 0xF9, 0xFA), line=color, line_w=Pt(2), radius=12000)
    txbox(sl, title, x+0.15, y+0.12, 5.9, 0.4,
          size=Pt(12), bold=True, color=color)
    for bi, b in enumerate(bullets):
        txbox(sl, f"• {b}", x+0.15, y+0.58+bi*0.44, 5.9, 0.42,
              size=Pt(10), color=DGREY)

txbox(sl, "Total loss = MSE_overall + 5.0 × MSE_initial + λ_shape × ShapeLoss",
      0.25, 7.1, 12.7, 0.32,
      size=Pt(11), bold=True, color=DKBLUE, align=PP_ALIGN.CENTER)


# ── SLIDE 5 — MLP vs LSTM ────────────────────────────────────────────────────
sl = add_slide(prs)
header_bar(sl, "MLP vs LSTM Decoder", "Why M6 replaced the LSTM with a direct per-slot MLP")

# MLP side
rect(sl, 0.2, 1.15, 5.9, 5.8, fill=RGBColor(0xF0, 0xFB, 0xF4), line=GREEN, line_w=Pt(2), radius=15000)
txbox(sl, "M6 — Per-Slot MLP  ✓", 0.4, 1.2, 5.5, 0.45,
      size=Pt(14), bold=True, color=GREEN)

mlp_items = [
    ("Slots 2–21 as input", "Each slot vector [64-dim] goes independently into the shared MLP"),
    ("Single forward pass", "64 → 32 → 3  (GELU activation)  —  all 20 slots in parallel"),
    ("No hidden state", "No sequential dependency between slots at decoding time"),
    ("Temporal structure from attention", "3 rounds of cross + self-attention encode temporal order"),
    ("Physics applied after MLP", "−|output| for KL, KR  ;  +|output| for KLR"),
    ("Cumsum reconstructs trajectory", "K(t) = K⁰ + Σ Δₙ  →  absolute stiffness values"),
]
for i, (h, d) in enumerate(mlp_items):
    txbox(sl, f"• {h}", 0.4, 1.75+i*0.74, 5.7, 0.28, size=Pt(11), bold=True, color=GREEN)
    txbox(sl, f"  {d}", 0.4, 2.0+i*0.74, 5.7, 0.45, size=Pt(10), color=DGREY)

# LSTM side
rect(sl, 6.5, 1.15, 6.55, 5.8, fill=RGBColor(0xFF, 0xF9, 0xE6), line=YELLOW, line_w=Pt(2), radius=15000)
txbox(sl, "M5 — LSTM Decoder  (removed in M6)", 6.7, 1.2, 6.2, 0.45,
      size=Pt(14), bold=True, color=YELLOW)

lstm_items = [
    ("Input: aggregated slot vector", "Mean of all 21 slots → single [64-dim] context vector"),
    ("2-layer LSTM unrolled 20 steps", "Hidden state hₜ = LSTM(hₜ₋₁, context)  — sequential"),
    ("Decoder cross-attention", "Each LSTM step attends back to all refined slots"),
    ("Positional embeddings", "Learnable pos_embed added to LSTM inputs for step identity"),
    ("h₀, c₀ projections", "Separate MLPs to initialise LSTM hidden & cell states"),
    ("101,312 extra parameters", "seq_decoder + decoder_cross_attn + h0_proj + c0_proj + pos_embed"),
]
for i, (h, d) in enumerate(lstm_items):
    txbox(sl, f"• {h}", 6.7, 1.75+i*0.74, 6.2, 0.28,
          size=Pt(11), bold=True, color=YELLOW)
    txbox(sl, f"  {d}", 6.7, 2.0+i*0.74, 6.2, 0.45, size=Pt(10), color=DGREY)

# Key insight bar
rect(sl, 0.2, 7.07, 12.9, 0.33, fill=DKBLUE, radius=8000)
txbox(sl, "Key insight: self-attention already encodes temporal order through slot position — LSTM adds sequential decoding overhead but the ordering is already captured",
      0.3, 7.09, 12.7, 0.3,
      size=Pt(9.5), color=WHITE, align=PP_ALIGN.CENTER)


# ── SLIDE 6 — Efficiency Gains ───────────────────────────────────────────────
sl = add_slide(prs)
header_bar(sl, "Efficiency Gains", "M6 vs M5 — measured on identical hardware (CPU, batch=32)")

# Big numbers row
metrics = [
    ("157,958  →  56,646", "Parameters", "64.1% FEWER", GREEN),
    ("632.7 KB  →  231.2 KB", "Model File Size", "63.5% SMALLER", BLUE),
    ("11.38 ms  →  9.07 ms", "Inference Time (CPU)", "1.25× FASTER", ORANGE),
]
for i, (val, label, gain, col) in enumerate(metrics):
    x = 0.4 + i * 4.3
    rect(sl, x, 1.2, 4.0, 2.1, fill=RGBColor(0xF8, 0xF9, 0xFA), line=col, line_w=Pt(2.5), radius=18000)
    txbox(sl, val, x+0.1, 1.28, 3.8, 0.7,
          size=Pt(14), bold=True, color=col, align=PP_ALIGN.CENTER)
    txbox(sl, label, x+0.1, 1.95, 3.8, 0.4,
          size=Pt(11), color=DGREY, align=PP_ALIGN.CENTER)
    txbox(sl, gain, x+0.1, 2.45, 3.8, 0.7,
          size=Pt(22), bold=True, color=col, align=PP_ALIGN.CENTER)

# Where did the parameters go?
txbox(sl, "Where did 101,312 parameters go?",
      0.4, 3.45, 12.5, 0.42, size=Pt(14), bold=True, color=DKBLUE)

removed = [
    ("seq_decoder  (LSTM 2-layer)", "66,176", "Entire sequential decoder — largest chunk"),
    ("decoder_cross_attn", "16,640", "Cross-attention between LSTM steps and slots"),
    ("h0_proj + c0_proj", "16,640", "Linear projections to initialise LSTM hidden/cell state"),
    ("pos_embed  (learnable)", "1,280", "One embedding vector per step  (20 × 64)"),
    ("TOTAL REMOVED", "100,736", ""),
    ("TOTAL ADDED (3rd iteration reuses existing layers)", "+0", "Zero new params — only extra compute"),
]
headers = ["Module Removed", "Parameters", "Purpose"]
rows = [[a, b, c] for a, b, c in removed]
rf = [LGREY, WHITE, LGREY, WHITE, RGBColor(0xD0, 0xF0, 0xD0), RGBColor(0xFF, 0xF3, 0xCD)]
table_slide(sl, headers, rows, 0.4, 3.95, 12.5, 3.3,
            header_fill=DKBLUE, row_fills=rf)

txbox(sl, "The 3rd refinement iteration reuses the SAME cross-attention + self-attention + MLP layers — zero parameter cost, only extra forward-pass compute (~10 ms).",
      0.4, 7.08, 12.5, 0.35, size=Pt(10), italic=True, color=DGREY)


# ── SLIDE 7 — Accuracy vs Efficiency Trade-off ───────────────────────────────
sl = add_slide(prs)
header_bar(sl, "Accuracy vs. Efficiency Trade-off", "Full metric comparison M5 ↔ M6")

txbox(sl, "M5 is more accurate on all metrics — M6 trades some accuracy for 64% fewer parameters and stronger physics guarantees.",
      0.3, 1.15, 12.7, 0.38, size=Pt(11), italic=True, color=DGREY)

headers = ["Target", "Metric", "M5", "M6", "Δ (M6 − M5)"]
rows = [
    ["Overall",  "R²",   "0.9913", "0.9804", "−0.0109"],
    ["Overall",  "RMSE", "1.98 × 10¹⁰", "2.98 × 10¹⁰", "+51%"],
    ["Overall",  "MAE",  "2.98 × 10⁹",  "3.33 × 10⁹",  "+12%"],
    ["KL",       "R²",   "0.9907", "0.9768", "−0.0139"],
    ["KL",       "RMSE", "6.21 × 10⁷",  "9.79 × 10⁷",  "+58%"],
    ["KR",       "R²",   "0.9838", "0.9634", "−0.0204"],
    ["KR",       "RMSE", "3.42 × 10¹⁰", "5.15 × 10¹⁰", "+51%"],
    ["KLR",      "R²",   "0.9882", "0.9690", "−0.0192"],
    ["KLR",      "RMSE", "1.18 × 10⁹",  "1.92 × 10⁹",  "+63%"],
    ["Slot 1→21","R² trend", "0.9908→0.9968", "0.9811→0.9809", "M5 improves; M6 flat"],
    ["Slot 1→21","RMSE trend","2.40×10¹⁰→2.69×10⁹","3.45×10¹⁰→6.57×10⁹","Both decrease (magnitudes shrink)"],
]
rf2 = []
for i in range(len(rows)):
    if rows[i][0] in ("Overall",) and rows[i][1]=="R²":
        rf2.append(RGBColor(0xE8, 0xF8, 0xE8))
    elif i % 2 == 0:
        rf2.append(LGREY)
    else:
        rf2.append(WHITE)

table_slide(sl, headers, rows, 0.25, 1.6, 12.8, 5.55,
            header_fill=DKBLUE, row_fills=rf2)

txbox(sl, "KR has the largest RMSE gap — rotational stiffness benefits most from LSTM's sequential context.  KLR R² gap is large because the new +|Δ| constraint changes the output distribution.",
      0.25, 7.1, 12.8, 0.35, size=Pt(9.5), italic=True, color=DGREY)


# ── SLIDE 8 — Error Metrics Explained ────────────────────────────────────────
sl = add_slide(prs)
header_bar(sl, "Error Metrics", "R², RMSE, MAE — definitions, properties, and how to read them")

# Definitions — 3 cards side by side
defs = [
    ("R²  (Coefficient of Determination)", BLUE,
     "R² = 1 − Σ(y−ŷ)² / Σ(y−ȳ)²",
     ["Range: −∞ to 1.0  (1.0 = perfect fit)",
      "Scale-free: comparable across variables",
      "Negative R² → model worse than predicting the mean",
      "Best metric for slot-to-slot comparison"]),
    ("RMSE  (Root Mean Square Error)", ORANGE,
     "RMSE = √[ (1/n) Σ(y−ŷ)² ]",
     ["Same physical units as the target (N, N·m…)",
      "Always ≥ 0  (squares → never negative)",
     "Penalises large errors heavily (squared term)",
      "Large values here reflect large stiffness magnitudes (10¹⁰)"]),
    ("MAE  (Mean Absolute Error)", GREEN,
     "MAE = (1/n) Σ|y−ŷ|",
     ["Same units as target, always ≥ 0",
      "Robust to outliers (no squaring)",
      "For same target: MAE ≤ RMSE always",
      "Easier to interpret physically"]),
]
for i, (title, col, formula, bullets) in enumerate(defs):
    x = 0.2 + i * 4.35
    rect(sl, x, 1.18, 4.1, 4.15, fill=RGBColor(0xF8, 0xF9, 0xFA), line=col, line_w=Pt(2), radius=14000)
    txbox(sl, title, x+0.1, 1.24, 3.9, 0.42, size=Pt(11), bold=True, color=col)
    txbox(sl, formula, x+0.1, 1.68, 3.9, 0.42,
          size=Pt(12), bold=True, color=DGREY, align=PP_ALIGN.CENTER)
    rect(sl, x+0.1, 2.14, 3.9, 0.02, fill=col)
    for bi, b in enumerate(bullets):
        txbox(sl, f"• {b}", x+0.1, 2.22+bi*0.68, 3.9, 0.62,
              size=Pt(10), color=DGREY)

# Granularity levels
txbox(sl, "Granularity Levels", 0.2, 5.4, 12.9, 0.38,
      size=Pt(13), bold=True, color=DKBLUE)

levels = [
    ("Overall", "Flatten all N×21×3 predictions", "Global model health check"),
    ("Per-variable", "Collapse time axis: N×21 per KL/KR/KLR", "Which output is hardest?"),
    ("Per-slot", "Collapse variable axis: N×3 per step 1–21", "Where on trajectory is error worst?"),
    ("Per-scenario", "One prediction of 21×3", "Which input combinations fail?"),
]
for i, (lvl, calc, use) in enumerate(levels):
    x = 0.2 + i * 3.27
    rect(sl, x, 5.82, 3.1, 1.5,
         fill=[RGBColor(0xD6, 0xEA, 0xF8), RGBColor(0xD5, 0xF5, 0xE3),
               RGBColor(0xFD, 0xEC, 0xD9), RGBColor(0xE8, 0xD5, 0xF5)][i],
         line=[BLUE, GREEN, ORANGE, PURPLE][i], radius=12000)
    txbox(sl, lvl, x+0.07, 5.88, 2.95, 0.32, size=Pt(11), bold=True,
          color=[DKBLUE, GREEN, ORANGE, PURPLE][i])
    txbox(sl, calc, x+0.07, 6.23, 2.95, 0.35, size=Pt(9.5), color=DGREY)
    txbox(sl, use, x+0.07, 6.6, 2.95, 0.6, size=Pt(9), italic=True, color=DGREY)


# ── SLIDE 9 — Errors vs Slot Results ─────────────────────────────────────────
sl = add_slide(prs)
header_bar(sl, "Understanding Errors Across 21 Slots", "Why per-slot RMSE/MAE decrease ≠ better accuracy — use R²")

# Pattern explanation
rect(sl, 0.2, 1.15, 8.65, 3.35, fill=RGBColor(0xF0, 0xF8, 0xFF), line=BLUE, line_w=Pt(1.5), radius=12000)
txbox(sl, "Observed Pattern: RMSE and MAE both decrease from Slot 1 → 21", 0.35, 1.2, 8.4, 0.38,
      size=Pt(12), bold=True, color=BLUE)

obs = [
    ("M5", "Slot 1 RMSE = 2.40×10¹⁰", "Slot 21 RMSE = 2.69×10⁹",  "Decrease: ~9×"),
    ("M6", "Slot 1 RMSE = 3.45×10¹⁰", "Slot 21 RMSE = 6.57×10⁹",  "Decrease: ~5×"),
]
for i, (m, s1, s21, dec) in enumerate(obs):
    x = 0.35 + i * 4.2
    txbox(sl, m, x, 1.63, 4.0, 0.3, size=Pt(11), bold=True,
          color=[BLUE, ORANGE][i])
    txbox(sl, f"{s1}  →  {s21}  ({dec})", x, 1.96, 4.0, 0.3,
          size=Pt(10), color=DGREY)

rect(sl, 0.35, 2.35, 8.4, 0.04, fill=RGBColor(0xBB, 0xCC, 0xDD))
txbox(sl, "WHY this happens — two reasons:", 0.35, 2.45, 8.4, 0.32,
      size=Pt(11), bold=True, color=DKBLUE)
reasons = [
    "① Physical magnitudes shrink:  Stiffness degradation is larger in early cycles (big initial drops), so absolute errors in N or N·m are naturally larger at Slot 1 than Slot 21.  RMSE and MAE are scale-dependent — they inherit this pattern.",
    "② Cumulative sum drift:  Early errors propagate forward (K(t) = K⁰ + Σ Δₙ), so Slot 1 carries the highest absolute burden.  Later drops are smaller, so errors contributed per step are also smaller.",
]
for i, r in enumerate(reasons):
    txbox(sl, r, 0.35, 2.82+i*0.6, 8.4, 0.55, size=Pt(10), color=DGREY)

# R² is the right metric box
rect(sl, 9.05, 1.15, 4.1, 3.35, fill=RGBColor(0xF0, 0xFB, 0xF4), line=GREEN, line_w=Pt(2), radius=12000)
txbox(sl, "Use R² for slot comparison", 9.15, 1.22, 3.9, 0.38,
      size=Pt(12), bold=True, color=GREEN)
r2_pts = [
    "R² is scale-normalised — measures how much better the model is than predicting the mean",
    "For M5: Slot R² improves  0.9908 → 0.9968 ↑  (LSTM adds sequential context as trajectory progresses)",
    "For M6: Slot R² stays flat  0.9811 → 0.9809 →  (no sequential memory — each slot predicted independently)",
    "This flat profile is M6's signature trade-off — uniform quality, no improvement over trajectory",
]
for i, p in enumerate(r2_pts):
    txbox(sl, f"• {p}", 9.15, 1.65+i*0.72, 3.9, 0.65, size=Pt(9.5), color=DGREY)

# Webapp percentage
rect(sl, 0.2, 4.6, 12.9, 1.4, fill=RGBColor(0xFF, 0xF8, 0xE3), line=YELLOW, line_w=Pt(1.5), radius=12000)
txbox(sl, "Webapp: Percentage Error per Step — the most interpretable metric", 0.35, 4.67, 12.6, 0.38,
      size=Pt(12), bold=True, color=RGBColor(0x7D, 0x60, 0x08))
txbox(sl, "Err% = |ŷ − y| / |y| × 100",
      0.35, 5.08, 4.5, 0.3, size=Pt(11), bold=True, color=DGREY)
txbox(sl, "Scale-independent: a 10% error at Slot 1 and Slot 21 are equally meaningful — unlike RMSE which is dominated by early-slot magnitudes.",
      4.9, 5.08, 8.1, 0.3, size=Pt(10), italic=True, color=DGREY)
for j, (label, col) in enumerate([("< 5%  Green (excellent)", RGBColor(0x27, 0xAE, 0x60)),
                                    ("5–20%  Orange (acceptable)", RGBColor(0xE6, 0x7E, 0x22)),
                                    ("> 20%  Red (poor)", RGBColor(0xC0, 0x39, 0x2B))]):
    txbox(sl, f"■  {label}", 0.35+j*4.2, 5.45, 4.1, 0.32,
          size=Pt(10), bold=True, color=col)

# Concise comparison table
txbox(sl, "Summary: which metric to use when?", 0.2, 6.1, 12.9, 0.35,
      size=Pt(12), bold=True, color=DKBLUE)
summary = [
    ["R²",    "Scale-free, 0–1",       "Slot-to-slot comparison, variable comparison, headline accuracy"],
    ["RMSE",  "Physical units, ≥ 0",   "Absolue error magnitude; dominated by early slots / large-magnitude variables"],
    ["MAE",   "Physical units, ≥ 0",   "Same as RMSE but more robust; physical interpretation"],
    ["Err%",  "Percentage, ≥ 0",       "Webapp display; best for per-step human-readable quality check"],
]
table_slide(sl, ["Metric", "Properties", "Best used for"],
            summary, 0.2, 6.48, 12.9, 0.97,
            header_fill=DKBLUE,
            row_fills=[LGREY, WHITE, LGREY, RGBColor(0xFD, 0xEC, 0xD9)])


# ── SLIDE 10 — Summary & Conclusion ──────────────────────────────────────────
sl = add_slide(prs)
rect(sl, 0, 0, 13.33, 7.5, fill=DKBLUE)
rect(sl, 0, 1.9, 13.33, 3.7, fill=BLUE)

txbox(sl, "M6 — Summary", 0.5, 0.25, 12.3, 1.0,
      size=Pt(34), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "Simpler · Faster · Physics-Guaranteed · LSTM-Free",
      0.5, 1.2, 12.3, 0.55,
      size=Pt(16), italic=True, color=RGBColor(0xBB, 0xDE, 0xFB),
      align=PP_ALIGN.CENTER)

summary_items = [
    ("Architecture", "21 learnable slots · 3× refinement (cross-attn → self-attn → MLP) · per-slot MLP prediction"),
    ("vs M5", "Removed LSTM decoder (101,312 params) · Added 3rd iteration · Added KLR sign constraint"),
    ("Efficiency", "64.1% fewer parameters · 63.5% smaller file · 1.25× faster inference"),
    ("Accuracy", "R² 0.9804 (vs 0.9913) · All three targets accurately predicted · KR largest gap"),
    ("Physics", "ΔKL ≤ 0, ΔKR ≤ 0, ΔKLR ≥ 0 hard-coded · Monotonicity structurally guaranteed"),
    ("Metrics", "R² best for slot comparison (scale-free) · RMSE/MAE in physical units (always ≥ 0)"),
]
for i, (k, v) in enumerate(summary_items):
    x = 0.4 + (i % 2) * 6.5
    y = 2.05 + (i // 2) * 1.1
    txbox(sl, k, x, y, 1.2, 0.35, size=Pt(10), bold=True,
          color=RGBColor(0xA8, 0xD8, 0xA8))
    txbox(sl, v, x+1.2, y, 5.1, 0.55, size=Pt(10), color=WHITE)

txbox(sl, "M6 = M5 − LSTM + stronger constraints + 3 attention iterations",
      0.5, 6.85, 12.3, 0.45,
      size=Pt(13), bold=True, color=RGBColor(0xA8, 0xD8, 0xA8),
      align=PP_ALIGN.CENTER)


# ── SLIDE 11 — What the Psi-NN Paper Contributes ───────────────────────────
sl = add_slide(prs)
header_bar(sl, "What the Ψ-NN Paper Adds", "Automatic structure discovery for physics-informed models")

txbox(sl, "Paper: Liu et al. (Nature Communications, 2025) — Automatic network structure discovery of physics informed neural networks via knowledge distillation",
        0.25, 1.15, 12.8, 0.4, size=Pt(10.5), italic=True, color=DGREY)

rect(sl, 0.25, 1.7, 4.0, 4.95, fill=RGBColor(0xF0, 0xF8, 0xFF), line=BLUE, line_w=Pt(2), radius=14000)
txbox(sl, "Core innovation", 0.4, 1.8, 3.7, 0.35, size=Pt(13), bold=True, color=BLUE)
for i, text in enumerate([
      "1. Train a teacher network for physics accuracy first.",
      "2. Train a student network with regularization, but without mixing that regularization into the physics loss.",
      "3. Cluster learned weights and detect repeated, opposite, or permuted patterns.",
      "4. Rebuild a smaller structured network with those rules hardwired.",
]):
      txbox(sl, text, 0.4, 2.22 + i * 0.88, 3.65, 0.7, size=Pt(10), color=DGREY)

rect(sl, 4.55, 1.7, 4.2, 4.95, fill=RGBColor(0xF0, 0xFB, 0xF4), line=GREEN, line_w=Pt(2), radius=14000)
txbox(sl, "Why it is physically interesting", 4.7, 1.8, 3.9, 0.35, size=Pt(13), bold=True, color=GREEN)
for i, text in enumerate([
      "The paper argues that PINNs usually enforce physics only through loss terms, not through internal wiring.",
      "Ψ-NN tries to convert physical patterns into architecture rules: symmetry, sign reversal, permutation, parameter sharing.",
      "That means the model can satisfy some constraints structurally at every forward pass, not only on average during training.",
      "This is the same philosophy as your M6 sign constraints and cumulative-sum monotonicity.",
]):
      txbox(sl, text, 4.7, 2.22 + i * 0.88, 3.9, 0.7, size=Pt(10), color=DGREY)

rect(sl, 9.0, 1.7, 4.08, 4.95, fill=RGBColor(0xFF, 0xF8, 0xE3), line=ORANGE, line_w=Pt(2), radius=14000)
txbox(sl, "Paper evidence", 9.15, 1.8, 3.8, 0.35, size=Pt(13), bold=True, color=ORANGE)
paper_rows = [
      ["Mechanism", "Teacher/student staged optimization"],
      ["Discovery", "Clustering + parameter reconstruction"],
      ["Physics encoded", "Symmetry, anti-symmetry, conservation-style structure"],
      ["Reported gain", "~50% fewer iterations to same loss on Laplace"],
      ["Reported gain", "~95% lower final L2 error on Laplace vs PINN"],
      ["Extra value", "Transferability across PDE settings / parameters"],
]
table_slide(sl, ["Claim", "Paper result"], paper_rows, 9.15, 2.22, 3.65, 3.95,
                  header_fill=DKBLUE,
                  row_fills=[LGREY, WHITE, LGREY, WHITE, LGREY, WHITE])

txbox(sl, "Important: those exact gains come from PDE benchmarks, not from your pile dataset, so they support the idea but do not guarantee the same numbers for M6.",
        0.25, 6.85, 12.8, 0.35, size=Pt(9.5), italic=True, color=DGREY)


# ── SLIDE 12 — Translation to Pile Physics ─────────────────────────────────
sl = add_slide(prs)
header_bar(sl, "How Ψ-NN Translates to Pile Physics", "What physical structure could be discovered or hardwired in M6")

txbox(sl, "The paper is about PDE symmetries, but the same idea can be translated to degradation physics: identify repeated physical patterns, then build them directly into the network.",
        0.25, 1.18, 12.8, 0.4, size=Pt(11), italic=True, color=DGREY)

headers = ["Paper idea", "Meaning in paper", "Translation to M6 / pile degradation"]
rows = [
      ["Sign rule", "Positive / negative parameter relations", "Already used: ΔKL, ΔKR ≤ 0 and ΔKLR ≥ 0 are structural physics constraints"],
      ["Shared prototypes", "Many weights are copies of a small set", "Drops across slots may come from a few degradation regimes: early, middle, late cycles"],
      ["Permutation / symmetry", "Different neurons are the same up to reordering", "KL and KR may share partial structure under similar soil-pile response patterns"],
      ["Relation matrix", "Explicit rules between parameters", "Could encode: slot j uses prototype p with scale αj, not a fully free head"],
      ["Transferability", "Structure survives parameter changes", "A discovered slot structure may transfer across different PI, Gmax, Lp or loading scenarios"],
      ["Reconstruction", "Rebuild a smaller architecture", "Replace many free slot-specific patterns by a smaller family of physically meaningful patterns"],
]
table_slide(sl, headers, rows, 0.25, 1.75, 12.8, 3.95,
                  header_fill=DKBLUE,
                  row_fills=[LGREY, WHITE, LGREY, WHITE, LGREY, WHITE])

rect(sl, 0.25, 5.95, 6.2, 1.15, fill=RGBColor(0xF0, 0xFB, 0xF4), line=GREEN, line_w=Pt(2), radius=12000)
txbox(sl, "Most direct physics bridge", 0.4, 6.05, 5.9, 0.28, size=Pt(12), bold=True, color=GREEN)
txbox(sl, "M6 already behaves like a small Ψ-NN idea: instead of only punishing bad physics in the loss, it hardwires monotonic degradation through sign constraints + cumulative sum.",
        0.4, 6.35, 5.8, 0.55, size=Pt(10), color=DGREY)

rect(sl, 6.7, 5.95, 6.35, 1.15, fill=RGBColor(0xFF, 0xF8, 0xE3), line=ORANGE, line_w=Pt(2), radius=12000)
txbox(sl, "Most innovative next step", 6.85, 6.05, 6.0, 0.28, size=Pt(12), bold=True, color=ORANGE)
txbox(sl, "Discover whether the 20 drop slots really need 20 independent behaviors, or whether they collapse into a few reusable physical prototypes with learned scaling rules.",
        6.85, 6.35, 5.95, 0.55, size=Pt(10), color=DGREY)


# ── SLIDE 13 — Can This Make M6 More Efficient? ────────────────────────────
sl = add_slide(prs)
header_bar(sl, "Can the Paper Make M6 More Efficient?", "Yes — if used as structure-guided compression, not as blind pruning")

rect(sl, 0.25, 1.2, 4.05, 2.1, fill=RGBColor(0xF0, 0xFB, 0xF4), line=GREEN, line_w=Pt(2), radius=13000)
txbox(sl, "Short answer: YES", 0.4, 1.32, 3.7, 0.32, size=Pt(16), bold=True, color=GREEN)
txbox(sl, "The paper's strongest idea is not just adding more physics. It is using physics + distillation to compress a model while protecting accuracy.",
        0.4, 1.72, 3.6, 0.95, size=Pt(10.5), color=DGREY)

rect(sl, 4.55, 1.2, 4.15, 2.1, fill=RGBColor(0xF0, 0xF8, 0xFF), line=BLUE, line_w=Pt(2), radius=13000)
txbox(sl, "Why it could help M6", 4.7, 1.32, 3.8, 0.32, size=Pt(14), bold=True, color=BLUE)
for i, text in enumerate([
      "M6 is already compact, so naive pruning may hurt accuracy.",
      "Ψ-NN-style discovery would remove redundancy only after learning the task well.",
      "That is better aligned with your objective: keep accuracy, compress only repeated structure.",
]):
      txbox(sl, text, 4.7, 1.72 + i * 0.42, 3.75, 0.35, size=Pt(9.7), color=DGREY)

rect(sl, 8.95, 1.2, 4.13, 2.1, fill=RGBColor(0xFF, 0xF8, 0xE3), line=ORANGE, line_w=Pt(2), radius=13000)
txbox(sl, "What to try in this project", 9.1, 1.32, 3.8, 0.32, size=Pt(14), bold=True, color=ORANGE)
for i, text in enumerate([
      "Teacher = current best M5 or M6 model.",
      "Student = regularized compact slot model.",
      "Then cluster slot/head weights and rebuild a smaller structured version.",
]):
      txbox(sl, text, 9.1, 1.72 + i * 0.42, 3.75, 0.35, size=Pt(9.7), color=DGREY)

txbox(sl, "Concrete efficiency ideas inspired by the paper", 0.25, 3.55, 12.8, 0.35, size=Pt(13), bold=True, color=DKBLUE)
idea_rows = [
      ["1. Prototype drop heads", "Cluster the 20 slot-drop behaviors into 3–5 physical regimes (early / middle / late) and use one shared prototype per regime", "Fewer free patterns, better interpretability, lower overfitting risk"],
      ["2. Structured KL/KR coupling", "Force partial weight sharing between KL and KR branches, with small correction terms", "May keep accuracy because both are monotonic stiffness channels"],
      ["3. Distilled M6-lite", "Train a tiny student from current M6 predictions, then add the same sign/cumsum physics rules", "Best chance to reduce size without losing much accuracy"],
      ["4. Slot relation matrix", "Represent some slot weights as scaled copies or sign-variants of prototype slots", "Closer to the Ψ-NN reconstruction idea"],
      ["5. Transfer by scenario family", "Discover one structure on one soil/loading regime, reuse it on nearby regimes", "Could reduce retraining cost and improve generalization"],
]
table_slide(sl, ["Idea", "How to apply it", "Why it may preserve accuracy"],
                  idea_rows, 0.25, 4.0, 12.8, 2.55,
                  header_fill=DKBLUE,
                  row_fills=[LGREY, WHITE, LGREY, WHITE, LGREY])

txbox(sl, "Best recommendation: use the paper as a second-stage compression method after training, not as a replacement for your current M6 training pipeline.",
        0.25, 6.75, 12.8, 0.28, size=Pt(10.5), bold=True, color=RGBColor(0x8B, 0x00, 0x00), align=PP_ALIGN.CENTER)
txbox(sl, "Expected outcome: moderate extra efficiency gains with smaller accuracy loss than manual simplification, but this needs experiments on your pile dataset because PDE symmetry results do not transfer one-to-one.",
        0.25, 7.02, 12.8, 0.25, size=Pt(9.4), italic=True, color=DGREY, align=PP_ALIGN.CENTER)


# ── SLIDE 14 — Does It Understand Neuron Layering? ────────────────────────
sl = add_slide(prs)
header_bar(sl, "Does the Method Understand Neuron Layering?", "Yes: it analyzes each layer matrix and finds reusable neuron patterns")

txbox(sl, "Simple answer: the method does not 'understand' like a human; it statistically detects repeated patterns in layer weights, then converts those patterns into architecture rules.",
      0.25, 1.18, 12.8, 0.4, size=Pt(11), italic=True, color=DGREY)

# Left: layer-by-layer process
rect(sl, 0.25, 1.7, 6.35, 5.45, fill=RGBColor(0xF0, 0xF8, 0xFF), line=BLUE, line_w=Pt(2), radius=14000)
txbox(sl, "Layer-by-layer discovery (simple view)", 0.4, 1.8, 6.0, 0.3, size=Pt(12), bold=True, color=BLUE)

steps = [
    "1) Take one trained layer matrix W (for example, from a slot head).",
    "2) Treat each row as one neuron's connection pattern.",
    "3) Measure row similarity: same shape, opposite sign, or permuted order.",
    "4) Cluster rows into groups: one prototype + relation rule.",
    "5) Rebuild layer with fewer independent rows and explicit relations.",
]
for i, text in enumerate(steps):
    txbox(sl, text, 0.45, 2.18 + i * 0.58, 5.95, 0.45, size=Pt(10.2), color=DGREY)

# Small matrix-like visual
rect(sl, 0.5, 5.15, 2.2, 1.7, fill=WHITE, line=RGBColor(0x88, 0xAA, 0xCC), radius=9000)
txbox(sl, "W (before)", 0.55, 5.2, 2.1, 0.22, size=Pt(9), bold=True, color=BLUE, align=PP_ALIGN.CENTER)
for r in range(4):
    rect(sl, 0.65, 5.45 + r * 0.3, 0.35, 0.2, fill=RGBColor(0xD6, 0xEA, 0xF8), line=WHITE)
    rect(sl, 1.05, 5.45 + r * 0.3, 0.35, 0.2, fill=RGBColor(0xD6, 0xEA, 0xF8), line=WHITE)
    rect(sl, 1.45, 5.45 + r * 0.3, 0.35, 0.2, fill=RGBColor(0xD6, 0xEA, 0xF8), line=WHITE)
    rect(sl, 1.85, 5.45 + r * 0.3, 0.35, 0.2, fill=RGBColor(0xD6, 0xEA, 0xF8), line=WHITE)

arrow(sl, 2.75, 5.95, 3.25, 5.95, color=BLUE)

rect(sl, 3.3, 5.15, 2.9, 1.7, fill=WHITE, line=RGBColor(0x88, 0xCC, 0xAA), radius=9000)
txbox(sl, "Prototype + rules", 3.35, 5.2, 2.8, 0.22, size=Pt(9), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txbox(sl, "p1", 3.55, 5.52, 0.6, 0.22, size=Pt(9), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txbox(sl, "p2", 3.55, 5.82, 0.6, 0.22, size=Pt(9), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txbox(sl, "row3 = -p1", 4.2, 5.52, 1.8, 0.22, size=Pt(9), color=DGREY)
txbox(sl, "row4 = perm(p2)", 4.2, 5.82, 1.8, 0.22, size=Pt(9), color=DGREY)

# Right: what that means for M6
rect(sl, 6.75, 1.7, 6.35, 5.45, fill=RGBColor(0xF0, 0xFB, 0xF4), line=GREEN, line_w=Pt(2), radius=14000)
txbox(sl, "What this means for M6", 6.9, 1.8, 6.0, 0.3, size=Pt(12), bold=True, color=GREEN)
for i, text in enumerate([
    "Drop MLP (64→32→3) can be analyzed neuron-by-neuron after training.",
    "If several hidden neurons behave almost the same, keep only one prototype neuron and map others by rules.",
    "If some slot behaviors are sign-flips/scaled copies, encode that explicitly instead of learning each slot independently.",
    "Result: fewer effective free parameters and more interpretable internal structure.",
    "Risk control: always validate on KR/KLR metrics, because these channels are most sensitive in your model.",
]):
    txbox(sl, f"• {text}", 6.95, 2.2 + i * 0.78, 5.95, 0.65, size=Pt(10), color=DGREY)


# ── SLIDE 15 — Graphical Architecture to Build ─────────────────────────────
sl = add_slide(prs)
header_bar(sl, "Graphical Architecture: M6 + Ψ-NN Compression", "Practical pipeline to keep accuracy and improve efficiency")

txbox(sl, "Use this as a two-stage engineering workflow: train for accuracy first, then discover and hardwire structure.",
      0.25, 1.18, 12.8, 0.35, size=Pt(11), italic=True, color=DGREY)

# Main flow boxes
block(sl, "A) Teacher\nStrong baseline\n(M5 or current M6)",
      0.35, 1.75, 2.3, 1.15,
      fill=RGBColor(0xD6, 0xEA, 0xF8), line=BLUE, txt_size=Pt(10), txt_color=DKBLUE)
arrow(sl, 2.68, 2.32, 3.2, 2.32, color=BLUE)

block(sl, "B) Student\nCompact model +\nregularization", 3.2, 1.75, 2.3, 1.15,
      fill=RGBColor(0xD5, 0xF5, 0xE3), line=GREEN, txt_size=Pt(10), txt_color=GREEN)
arrow(sl, 5.53, 2.32, 6.05, 2.32, color=GREEN)

block(sl, "C) Weight Analysis\nCluster rows/neurons\nfind relations", 6.05, 1.75, 2.45, 1.15,
      fill=RGBColor(0xE8, 0xD5, 0xF5), line=PURPLE, txt_size=Pt(10), txt_color=PURPLE)
arrow(sl, 8.55, 2.32, 9.05, 2.32, color=PURPLE)

block(sl, "D) Rebuild\nStructured M6-lite\n(prototypes + rules)", 9.05, 1.75, 2.45, 1.15,
      fill=RGBColor(0xFD, 0xEC, 0xD9), line=ORANGE, txt_size=Pt(10), txt_color=ORANGE)
arrow(sl, 11.53, 2.32, 12.0, 2.32, color=ORANGE)

block(sl, "E) Validate\nR², RMSE, MAE\nper variable + per slot", 12.0, 1.75, 1.1, 1.15,
      fill=RGBColor(0xFDE, 0xBD, 0xBD) if False else RGBColor(0xFD, 0xBD, 0xBD), line=RED, txt_size=Pt(8.7), txt_color=RGBColor(0x8B, 0, 0))

# Architecture inside student/rebuilt
rect(sl, 0.35, 3.2, 12.75, 3.85, fill=RGBColor(0xF8, 0xF9, 0xFA), line=RGBColor(0xAA, 0xBB, 0xCC), line_w=Pt(1.5), radius=12000)
txbox(sl, "Inside the model (what changes and what stays)", 0.5, 3.28, 12.4, 0.3, size=Pt(12), bold=True, color=DKBLUE)

cols = [0.55, 3.05, 5.55, 8.05, 10.55]
labels = [
    ("Input + Embed", "8→64"),
    ("Slot Refinement", "Cross/Self/MLP ×3"),
    ("Slot Heads", "Initial head + Drop head"),
    ("Ψ-Extraction", "Prototype neurons + relation rules"),
    ("Structured Heads", "smaller effective DoF"),
]
fills = [RGBColor(0xD6, 0xEA, 0xF8), RGBColor(0xD5, 0xF5, 0xE3), RGBColor(0xFF, 0xF3, 0xCD), RGBColor(0xE8, 0xD5, 0xF5), RGBColor(0xFD, 0xEC, 0xD9)]
lines = [BLUE, GREEN, YELLOW, PURPLE, ORANGE]
for i, x in enumerate(cols):
    block(sl, labels[i][0] + "\n" + labels[i][1], x, 3.75, 2.1, 1.0,
          fill=fills[i], line=lines[i], txt_size=Pt(9.2), txt_color=DGREY)
    if i < len(cols) - 1:
        arrow(sl, x + 2.1, 4.25, cols[i + 1], 4.25, color=DGREY)

txbox(sl, "Keep physics hard constraints unchanged: ΔKL, ΔKR ≤ 0 and ΔKLR ≥ 0; cumsum reconstruction remains identical.",
      0.55, 5.02, 12.2, 0.3, size=Pt(10), bold=True, color=RGBColor(0x8B, 0, 0))

table_slide(
    sl,
    ["Part", "Before", "After Ψ-NN-style rebuild"],
    [
        ["Refinement block", "Same", "Same (no need to change)"],
        ["Drop head behavior", "Fully free", "Prototype-constrained"],
        ["Interpretability", "Low", "Higher (explicit relations)"],
        ["Expected efficiency", "Current", "Better with controlled accuracy drop"],
    ],
    0.55, 5.4, 12.2, 1.55,
    header_fill=DKBLUE,
    row_fills=[LGREY, WHITE, LGREY, WHITE],
)


# ── SLIDE 16 — Worked Example (Simple) ─────────────────────────────────────
sl = add_slide(prs)
header_bar(sl, "Worked Example: From 20 Drop Slots to 4 Prototypes", "Detailed enough to understand, simple enough to explain")

txbox(sl, "Goal: show how neuron/slot structure discovery can compress the model without changing the physics output logic.",
      0.25, 1.18, 12.8, 0.35, size=Pt(11), italic=True, color=DGREY)

rect(sl, 0.25, 1.65, 6.35, 5.35, fill=RGBColor(0xF0, 0xF8, 0xFF), line=BLUE, line_w=Pt(2), radius=12000)
txbox(sl, "Step-by-step example", 0.4, 1.75, 6.0, 0.3, size=Pt(12), bold=True, color=BLUE)

example_lines = [
    "1) Train baseline M6 normally (same losses and constraints).",
    "2) Collect slot-drop latent vectors for slots 2..21 on validation set.",
    "3) Cluster these 20 slots into 4 groups: early, early-mid, late-mid, late.",
    "4) Learn 4 prototype drop mappings P1..P4 instead of 20 independent behaviors.",
    "5) For each slot j, predict Δ using prototype Pk plus a tiny scalar adapter αj.",
    "6) Keep sign constraints and cumulative sum exactly unchanged.",
    "7) Compare old/new on per-variable per-slot RMSE and MAE.",
]
for i, txt in enumerate(example_lines):
    txbox(sl, txt, 0.45, 2.15 + i * 0.58, 5.95, 0.45, size=Pt(10), color=DGREY)

rect(sl, 0.45, 6.35, 5.95, 0.45, fill=RGBColor(0xFF, 0xF8, 0xE3), line=ORANGE, radius=7000)
txbox(sl, "Presentation sentence: 'We replace many free slot behaviors with a few physically reusable prototypes.'",
      0.52, 6.45, 5.8, 0.25, size=Pt(9.3), bold=True, color=RGBColor(0x7D, 0x60, 0x08))

# Right side visual mapping
rect(sl, 6.85, 1.65, 6.2, 5.35, fill=RGBColor(0xF0, 0xFB, 0xF4), line=GREEN, line_w=Pt(2), radius=12000)
txbox(sl, "Graphical mapping", 7.0, 1.75, 5.9, 0.3, size=Pt(12), bold=True, color=GREEN)

block(sl, "Slots 2..21\n(20 drop slots)", 7.0, 2.15, 2.5, 0.9,
      fill=RGBColor(0xFF, 0xF3, 0xCD), line=YELLOW, txt_size=Pt(9.5))
arrow(sl, 9.52, 2.6, 10.05, 2.6, color=DGREY)
block(sl, "Cluster into\n4 groups", 10.05, 2.15, 1.85, 0.9,
      fill=RGBColor(0xE8, 0xD5, 0xF5), line=PURPLE, txt_size=Pt(9.5))

for i, label in enumerate(["P1", "P2", "P3", "P4"]):
    block(sl, label, 7.05 + i * 1.45, 3.35, 1.2, 0.55,
          fill=RGBColor(0xD5, 0xF5, 0xE3), line=GREEN, txt_size=Pt(10), bold=True, txt_color=GREEN)

txbox(sl, "slot j uses Pk with scale αj", 7.0, 4.02, 5.9, 0.3, size=Pt(10), color=DGREY, align=PP_ALIGN.CENTER)
arrow(sl, 9.95, 4.35, 9.95, 4.75, color=GREEN)
block(sl, "Physics block\n−|ΔKL,ΔKR| , +|ΔKLR|\nthen cumulative sum", 8.1, 4.75, 3.7, 1.0,
      fill=RGBColor(0xFD, 0xEC, 0xD9), line=ORANGE, txt_size=Pt(9))

table_slide(
    sl,
    ["Check", "Accept if"],
    [
        ["Overall accuracy", "R² drop is small (project-defined threshold)"],
        ["KR stability", "KR RMSE/MAE does not degrade sharply"],
        ["KLR physics", "Sign/monotonic behavior remains physically valid"],
        ["Efficiency", "Params and runtime improve enough to justify change"],
    ],
    7.0, 5.9, 5.9, 1.05,
    header_fill=DKBLUE,
    row_fills=[LGREY, WHITE, LGREY, WHITE],
)


# ── SAVE ─────────────────────────────────────────────────────────────────────
out = r"c:\Users\youss\Downloads\PFE\M6\new_out.pptx"
try:
      prs.save(out)
      saved_to = out
except PermissionError:
      try:
            saved_to = r"c:\Users\youss\Downloads\PFE\M6\new_out_updated.pptx"
            prs.save(saved_to)
      except PermissionError:
            stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            saved_to = rf"c:\Users\youss\Downloads\PFE\M6\new_out_{stamp}.pptx"
            prs.save(saved_to)
print(f"Saved → {saved_to}")
print(f"Slides: {len(prs.slides)}")
