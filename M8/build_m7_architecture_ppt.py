from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT_PATH = r"c:\Users\youss\Downloads\PFE\M7\M7_PsiNN_Architecture_Exact_Explanation.pptx"

prs = Presentation()

# Simple visual system
COLOR_BG = RGBColor(245, 248, 252)
COLOR_TITLE = RGBColor(16, 37, 66)
COLOR_TEXT = RGBColor(35, 48, 66)
COLOR_ACCENT = RGBColor(10, 116, 218)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.65), Inches(1.15), Inches(12.0), Inches(0.6))
        stf = sub.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(16)
        sp.font.color.rgb = COLOR_ACCENT


def add_bullets(slide, x, y, w, h, items, font_size=22):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_TEXT


def add_slide_title_and_bullets(title, subtitle, bullets, font_size=20):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, title, subtitle)
    add_bullets(slide, 0.8, 1.9, 12.0, 5.0, bullets, font_size=font_size)
    return slide


# Slide 1: Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(
    slide,
    "M7 Ψ-NN Architecture",
    "Exact explanation of what was rewritten: slot structure vs MLP heads"
)
add_bullets(
    slide,
    0.8,
    2.1,
    12.0,
    3.8,
    [
        "Project: Pile stiffness degradation (KL, KR, KLR over 21 steps)",
        "Pipeline: Stage A Distillation -> Stage B Structure Discovery -> Stage C Structured Reconstruction",
        "Core point: M7 does not only rewrite MLP; it rewrites slot parameterization in Stage C"
    ],
    font_size=20,
)

# Slide 2: Add the exact answer
add_slide_title_and_bullets(
    "Direct Answer (Exact)",
    "What is rewritten in M7?",
    [
        "In M7, it is not only the MLP that is rewritten.",
        "The slot-based model structure is rewritten for the final Ψ model.",
        "Teacher (M6 copy) and Stage-A Student keep full slot-attention architecture.",
        "Stage-C Ψ model replaces 20 independent drop slots with k* prototype slots + relation matrix R.",
        "Attention refinement and projection heads are still present; main rewrite is slot representation/reconstruction."
    ],
)

# Slide 3: M6 baseline
add_slide_title_and_bullets(
    "Baseline (M6)",
    "Original SlotAttentionDegradation",
    [
        "Input features (8) -> embedding (d_model=64)",
        "21 learnable slots: 1 initial slot + 20 drop slots",
        "Iterative refinement: cross-attn + self-attn + slot MLP (3 iterations)",
        "Heads: initial_proj and drop_proj to predict KL, KR, KLR sequence",
        "Physics constraints: KL/KR drops negative, KLR drop positive"
    ],
)

# Slide 4: Stage A
add_slide_title_and_bullets(
    "Stage A: Distillation Student",
    "What changes and what stays",
    [
        "Architecture mostly unchanged from M6 (same slot-attention backbone)",
        "Student adds L1 regularization on slot vectors to encourage sparse structure",
        "Loss = MSE(student, teacher) + 0.5*MSE(student, target) + mu*L1(slots)",
        "Goal: learn teacher behavior while exposing slot redundancy patterns"
    ],
)

# Slide 5: Stage B
add_slide_title_and_bullets(
    "Stage B: Structure Discovery",
    "How k* and R are discovered",
    [
        "Extract refined student slot vectors and keep only 20 drop slots",
        "Average drop-slot vectors across dataset -> 20 x 64 representation",
        "Compute cosine similarity matrix for slot relations",
        "Run KMeans for k=2..10 and choose k* by best silhouette score",
        "Build relation matrix R (20 x k*) via inverse-distance soft assignment"
    ],
)

# Slide 6: Stage C core rewrite
add_slide_title_and_bullets(
    "Stage C: Ψ-Model Core Rewrite",
    "This is the architectural rewrite",
    [
        "Instead of 20 independent drop-slot parameters, learn only k* prototype slots",
        "Reconstruct full 20 drop slots using: drop_slots = R @ prototypes",
        "Apply learned per-slot scales for fine correction",
        "Concatenate with initial slot -> full 21 slots",
        "Then run same iterative attention refinement and output heads"
    ],
)

# Slide 7: Equation view
add_slide_title_and_bullets(
    "Mathematical View",
    "Structured reconstruction logic",
    [
        "Prototype bank: P in R^(k* x d)",
        "Relation matrix: R in R^(20 x k*)",
        "Reconstructed drop slots: S_drop = R P",
        "Scaled drop slots: S_drop_scaled = S_drop * alpha (per-slot scale)",
        "Final slot set: S = [s_init ; S_drop_scaled]"
    ],
)

# Slide 8: What is NOT rewritten
add_slide_title_and_bullets(
    "What Is Not Rewritten",
    "Important clarification",
    [
        "Cross-attention block is not replaced",
        "Self-attention block is not replaced",
        "Slot MLP refinement block is not replaced",
        "Initial and drop projection heads remain conceptually the same",
        "Main new part is slot parameterization (prototypes + R + scaling)"
    ],
)

# Slide 9: Code mapping
add_slide_title_and_bullets(
    "Code Mapping (train.py)",
    "Where each concept lives",
    [
        "SlotAttentionDegradation: M6 teacher copy",
        "SlotAttentionStudent: Stage-A model + l1_regularization",
        "stage_b_structure_discovery: k* selection and relation matrix construction",
        "SlotAttentionPsiModel: prototype slots, relation matrix, slot reconstruction",
        "stage_c_structured_training: structured Ψ training with multi-objective loss"
    ],
    font_size=19,
)

# Slide 10: Training objective in Stage C
add_slide_title_and_bullets(
    "Stage C Objective",
    "Why the Ψ model still matches physics and teacher",
    [
        "Distillation loss: match teacher outputs",
        "Sequence loss: fit true target sequence",
        "Initial condition loss: stronger weight on first step",
        "Shape loss: match step-to-step change profile",
        "Result: compressed slot structure with maintained predictive behavior"
    ],
)

# Slide 11: Parameter interpretation
add_slide_title_and_bullets(
    "Interpretability Gain",
    "From many slots to shared prototypes",
    [
        "Each prototype captures a recurring degradation pattern",
        "Each physical step slot is expressed as combination of prototypes",
        "R matrix tells which prototype influences which drop step",
        "Cluster memberships provide a direct structure explanation",
        "This is the Ψ-NN idea transferred to pile degradation"
    ],
)

# Slide 12: Summary
add_slide_title_and_bullets(
    "Final Summary",
    "Answer to your exact question",
    [
        "M7 is not only MLP rewrite.",
        "Stage C rewrites slot parameterization globally using prototypes + relation matrix.",
        "Attention/refinement blocks and output heads remain in the same family as M6.",
        "So: all-model-with-slots is structurally rewritten at the slot level in Stage C."
    ],
)

prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")
