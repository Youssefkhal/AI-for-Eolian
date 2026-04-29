from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

SRC_PPT_PATH = r"c:\Users\youss\Downloads\PFE\M7\new_out_20260414_075825_graphical_stages.pptx"
OUT_PPT_PATH = r"c:\Users\youss\Downloads\PFE\M7\new_out_20260414_075825_graphical_stages_with_result_interpretation.pptx"

prs = Presentation(SRC_PPT_PATH)

BG = RGBColor(244, 248, 252)
TITLE = RGBColor(21, 45, 74)
TEXT = RGBColor(35, 52, 70)
WHITE = RGBColor(255, 255, 255)
ACCENT_A = RGBColor(0, 112, 192)
ACCENT_B = RGBColor(46, 125, 50)
ACCENT_C = RGBColor(230, 126, 34)
ACCENT_D = RGBColor(120, 60, 170)
LIGHT_GREEN = RGBColor(226, 239, 218)
LIGHT_BLUE = RGBColor(221, 235, 247)
LIGHT_ORANGE = RGBColor(252, 234, 220)
LIGHT_PURPLE = RGBColor(248, 232, 246)
LIGHT_WHITE = RGBColor(255, 255, 255)


def setup_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.8))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TITLE

    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.55), Inches(0.95), Inches(12.0), Inches(0.45))
        sf = s.text_frame
        sp = sf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(15)
        sp.font.color.rgb = ACCENT_D
    return slide


def readable_text_color(color):
    luminance = 0.2126 * color[0] + 0.7152 * color[1] + 0.0722 * color[2]
    return TEXT if luminance > 175 else WHITE


def add_box(slide, x, y, w, h, text, color, fs=16, bold=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = RGBColor(200, 210, 220)

    tf = shp.text_frame
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(fs)
        p.font.bold = bold if i == 0 else False
        p.font.color.rgb = readable_text_color(color)
    return shp


def add_arrow(slide, x, y, w, h, color=RGBColor(130, 140, 155)):
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.color.rgb = color
    return arr


slide = setup_slide(
    "Interpretation of Current M7 Result",
    "What Stage A, B, and C mean for the discovered pile degradation structure"
)
add_box(slide, 0.7, 1.6, 3.8, 1.0, "Stage A succeeded\nStudent learned M6 well", ACCENT_A, fs=17, bold=True)
add_box(slide, 4.8, 1.6, 3.3, 1.0, "Stage B succeeded\nStructure was discovered", ACCENT_B, fs=17, bold=True)
add_box(slide, 8.4, 1.6, 3.8, 1.0, "Stage C started\nStructured Psi model training", ACCENT_C, fs=17, bold=True)
add_box(
    slide, 0.8, 3.0, 11.5, 2.7,
    "Stage A meaning:\n"
    "- Student stopped early at epoch 779 because validation stopped improving enough\n"
    "- Best distillation loss = 0.001106, so student output is very close to teacher output\n"
    "- This means M6 behavior has been transferred successfully into the student",
    LIGHT_WHITE, fs=15
)

slide = setup_slide(
    "Stage B Interpretation",
    "The most important structural result: 20 drop slots collapsed into 6 prototype groups"
)
add_box(slide, 0.7, 1.6, 2.6, 1.0, "Original structure\n20 drop slots", ACCENT_A, fs=18, bold=True)
add_arrow(slide, 3.55, 1.95, 1.0, 0.35)
add_box(slide, 4.7, 1.6, 2.9, 1.0, "Discovered structure\n6 prototypes", ACCENT_B, fs=18, bold=True)
add_arrow(slide, 7.85, 1.95, 1.0, 0.35)
add_box(slide, 9.0, 1.6, 3.2, 1.0, "Meaning\n6 degradation modes", ACCENT_C, fs=18, bold=True)
add_box(
    slide, 0.7, 3.0, 5.7, 3.0,
    "Cluster result:\n"
    "- Prototype 5: slots 2-7\n"
    "- Prototype 0: slots 8-10\n"
    "- Prototype 2: slots 11-14\n"
    "- Prototype 3: slots 15-17\n"
    "- Prototype 1: slots 18-19\n"
    "- Prototype 4: slots 20-21",
    LIGHT_GREEN, fs=15
)
add_box(
    slide, 6.7, 3.0, 5.5, 3.0,
    "Interpretation:\n"
    "- The model does not need 20 fully independent drop-slot patterns\n"
    "- It found only 6 distinct pattern families\n"
    "- Nearby time steps behave similarly and can share prototype structure\n"
    "- Silhouette = 0.715 means the clustering is strong and well separated",
    LIGHT_BLUE, fs=15
)

slide = setup_slide(
    "Physical Meaning of the Prototype Groups",
    "How the discovered slot groups can be read as staged degradation behavior"
)
add_box(
    slide, 0.7, 1.6, 3.8, 1.6,
    "Early degradation mode\nSlots 2-7\nLikely captures the first strong degradation regime",
    LIGHT_ORANGE, fs=16, bold=True
)
add_box(
    slide, 4.75, 1.6, 3.0, 1.6,
    "Mid regime A\nSlots 8-10\nTransition after early drop",
    LIGHT_GREEN, fs=16, bold=True
)
add_box(
    slide, 8.0, 1.6, 4.2, 1.6,
    "Mid regime B/C\nSlots 11-17\nSeveral grouped middle-stage degradation patterns",
    LIGHT_BLUE, fs=16, bold=True
)
add_box(
    slide, 0.7, 3.7, 5.6, 2.0,
    "Late-tail modes\nSlots 18-21\nSmall end-stage adjustment groups\nThese likely represent residual or saturation behavior",
    LIGHT_PURPLE, fs=16, bold=True
)
add_box(
    slide, 6.6, 3.7, 5.6, 2.0,
    "Practical meaning:\n"
    "M7 suggests the 20 degradation steps are organized into about 6 recurring physical patterns,\nnot 20 unrelated behaviors.",
    LIGHT_WHITE, fs=15
)

slide = setup_slide(
    "Why Compression Looks Small",
    "1.5% fewer total parameters does not mean the discovery is weak"
)
add_box(
    slide, 0.7, 1.6, 5.5, 2.3,
    "Teacher params: 56,646\nPsi-model params: 55,770\nCompression: 1.5% fewer parameters",
    ACCENT_C, fs=18, bold=True
)
add_box(
    slide, 6.5, 1.6, 5.7, 2.3,
    "Reason:\nOnly the slot bank was compressed strongly.\nAttention layers, slot MLP, norms, and heads still dominate total size.",
    LIGHT_WHITE, fs=16
)
add_box(
    slide, 0.7, 4.3, 11.5, 1.8,
    "So the main gain is not huge raw parameter reduction.\nThe main gain is structural interpretability, shared degradation modes, and a cleaner internal architecture.",
    LIGHT_GREEN, fs=16, bold=True
)

slide = setup_slide(
    "One-Sentence Conclusion",
    "Final interpretation of the current result"
)
add_box(
    slide, 0.9, 2.1, 11.1, 2.5,
    "M7 has already shown that the pile stiffness degradation process learned by M6 is internally organized into about 6 recurring prototype modes, meaning the 20 drop steps are not independent but belong to a smaller structured family of degradation behaviors.",
    LIGHT_WHITE, fs=19, bold=True
)

prs.save(OUT_PPT_PATH)
print(f"Updated presentation: {OUT_PPT_PATH}")
print(f"Total slides now: {len(prs.slides)}")
