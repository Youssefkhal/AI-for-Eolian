from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

SRC = r"c:\Users\youss\Downloads\PFE\M7\new_out_20260414_075825_graphical_stages_with_result_interpretation.pptx"
OUT = r"c:\Users\youss\Downloads\PFE\M7\new_out_20260414_075825_graphical_stages_with_result_interpretation.pptx"

prs = Presentation(SRC)

BG         = RGBColor(244, 248, 252)
TITLE_C    = RGBColor(21,  45,  74)
TEXT       = RGBColor(35,  52,  70)
WHITE      = RGBColor(255, 255, 255)
ACCENT_A   = RGBColor(0,   112, 192)
ACCENT_B   = RGBColor(46,  125,  50)
ACCENT_C   = RGBColor(230, 126,  34)
ACCENT_D   = RGBColor(120,  60, 170)
ACCENT_RED = RGBColor(192,   0,   0)
LIGHT_GREEN  = RGBColor(226, 239, 218)
LIGHT_BLUE   = RGBColor(221, 235, 247)
LIGHT_ORANGE = RGBColor(252, 234, 220)
LIGHT_PURPLE = RGBColor(248, 232, 246)
LIGHT_GREY   = RGBColor(240, 240, 240)
DARK_GREY    = RGBColor(80,  80,  80)


def readable(color):
    lum = 0.2126 * color[0] + 0.7152 * color[1] + 0.0722 * color[2]
    return TEXT if lum > 175 else WHITE


def setup_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.8))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(29)
    p.font.bold = True
    p.font.color.rgb = TITLE_C

    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.55), Inches(0.93), Inches(12.0), Inches(0.45))
        sp = s.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = ACCENT_D
    return slide


def box(slide, x, y, w, h, text, color, fs=15, bold=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = RGBColor(200, 210, 220)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(fs)
        p.font.bold = bold if i == 0 else False
        p.font.color.rgb = readable(color)
    return shp


def arrow(slide, x, y, w, h, color=RGBColor(130, 140, 155)):
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.color.rgb = color
    return arr


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — What is the Silhouette Score?
# ─────────────────────────────────────────────────────────────────────────────
slide = setup_slide(
    "What is the Silhouette Score?",
    "The metric used in Stage B to choose the optimal number of prototypes k*"
)

box(slide, 0.6, 1.45, 5.6, 1.5,
    "Formula for each slot vector sᵢ:\n\n"
    "silhouette(i) = ( b(i) − a(i) ) / max(a(i), b(i))",
    LIGHT_BLUE, fs=16, bold=True)

box(slide, 6.5, 1.45, 5.7, 1.5,
    "a(i) = avg distance to other slots in the SAME cluster\n"
    "b(i) = avg distance to slots in the NEAREST other cluster",
    LIGHT_ORANGE, fs=15)

box(slide, 0.6, 3.2, 11.6, 1.0,
    "Range: −1 to +1   |   +1 = tight well-separated clusters   |   0 = overlapping   |   negative = wrong assignment",
    LIGHT_GREY, fs=15, bold=True)

box(slide, 0.6, 4.45, 2.6, 1.5, "0.71 – 1.00\nStrong structure",   ACCENT_B,   fs=16, bold=True)
box(slide, 3.4, 4.45, 2.6, 1.5, "0.51 – 0.70\nReasonable structure", ACCENT_A,   fs=16, bold=True)
box(slide, 6.2, 4.45, 2.6, 1.5, "0.26 – 0.50\nWeak structure",       ACCENT_C,   fs=16, bold=True)
box(slide, 9.0, 4.45, 2.6, 1.5, "< 0.25\nNo real structure",         ACCENT_RED, fs=16, bold=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Silhouette result in M7
# ─────────────────────────────────────────────────────────────────────────────
slide = setup_slide(
    "Silhouette Score — M7 Stage B Result",
    "Scores computed for k = 2 to 10; best k chosen automatically"
)

rows = [
    ("k=2",  "0.612", False),
    ("k=3",  "0.641", False),
    ("k=4",  "0.647", False),
    ("k=5",  "0.685", False),
    ("k=6",  "0.715 ◄ BEST", True),
    ("k=7",  "0.673", False),
    ("k=8",  "0.620", False),
    ("k=9",  "0.583", False),
    ("k=10", "0.510", False),
]

box(slide, 0.6, 1.45, 1.8, 0.55, "k", ACCENT_A, fs=16, bold=True)
box(slide, 2.6, 1.45, 3.0, 0.55, "Silhouette Score", ACCENT_A, fs=16, bold=True)

for i, (k, sc, best) in enumerate(rows):
    y = 2.1 + i * 0.52
    c = LIGHT_GREEN if best else LIGHT_GREY
    box(slide, 0.6, y, 1.8, 0.50, k,  c, fs=15, bold=best)
    box(slide, 2.6, y, 3.0, 0.50, sc, c, fs=15, bold=best)

box(slide, 6.5, 1.45, 5.7, 3.5,
    "Why 0.715 is good:\n\n"
    "- 20 slot vectors in 64-dimensional space\n"
    "- High-dimensional data is usually harder to cluster\n"
    "- Score above 0.70 confirms real physical groupings exist\n"
    "- k*=6 was chosen because it gives the highest silhouette,\n"
    "  not because we guessed it — the data itself chose it",
    LIGHT_BLUE, fs=15)

box(slide, 6.5, 5.2, 5.7, 1.0,
    "Result: 20 drop slots are genuinely organized into 6 distinct degradation modes",
    ACCENT_B, fs=16, bold=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Final Metrics Comparison Table
# ─────────────────────────────────────────────────────────────────────────────
slide = setup_slide(
    "Final Performance Comparison",
    "M6 Teacher  vs  Stage-A Student  vs  Ψ-Model (M7) — test set, original scale"
)

# Column headers
box(slide, 0.5,  1.5, 2.5, 0.55, "Metric",           ACCENT_A, fs=15, bold=True)
box(slide, 3.2,  1.5, 3.0, 0.55, "M6 Teacher",       ACCENT_A, fs=15, bold=True)
box(slide, 6.4,  1.5, 3.0, 0.55, "Stage-A Student",  ACCENT_A, fs=15, bold=True)
box(slide, 9.6,  1.5, 3.0, 0.55, "Ψ-Model (M7)",     ACCENT_A, fs=15, bold=True)

table_data = [
    ("Overall R²",  "0.9804", "0.9900", "0.9897"),
    ("KL    R²",    "0.9768", "0.9886", "0.9910"),
    ("KR    R²",    "0.9634", "0.9814", "0.9807"),
    ("KLR   R²",    "0.9690", "0.9854", "0.9860"),
    ("Params",      "56,646", "56,646", "55,770"),
]

row_colors = [LIGHT_GREY, LIGHT_BLUE, LIGHT_GREY, LIGHT_BLUE, LIGHT_GREY]

for i, (metric, m6, stA, psi) in enumerate(table_data):
    y = 2.15 + i * 0.60
    c = row_colors[i]
    box(slide, 0.5, y, 2.5, 0.58, metric, c, fs=14, bold=True)
    box(slide, 3.2, y, 3.0, 0.58, m6,     c, fs=14)
    box(slide, 6.4, y, 3.0, 0.58, stA,    c, fs=14)
    box(slide, 9.6, y, 3.0, 0.58, psi,    c, fs=14)

box(slide, 0.5, 5.25, 12.1, 1.1,
    "Ψ-Model retains ≈99% of M6 accuracy while using 1.5% fewer parameters and revealing a structured internal representation of pile degradation.",
    LIGHT_GREEN, fs=16, bold=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Key Takeaway
# ─────────────────────────────────────────────────────────────────────────────
slide = setup_slide(
    "Key Takeaway — What M7 Achieved",
    "Summary of results from the completed Ψ-NN pipeline run"
)

box(slide, 0.6, 1.5, 3.5, 1.5,
    "Structure found\nk* = 6 prototypes\nSilhouette = 0.715",
    ACCENT_B, fs=18, bold=True)

box(slide, 4.35, 1.5, 3.5, 1.5,
    "Accuracy kept\nΨ-model R² = 0.9897\nvs teacher 0.9804",
    ACCENT_A, fs=18, bold=True)

box(slide, 8.1, 1.5, 3.9, 1.5,
    "Compressed\n55,770 params\n(1.5% fewer than M6)",
    ACCENT_C, fs=18, bold=True)

box(slide, 0.6, 3.3, 11.5, 1.4,
    "The 20 degradation time-steps are not independent — they group into 6 recurring physical patterns.\n"
    "M7 discovered this structure automatically from data, without any prior assumption about how many groups exist.",
    LIGHT_BLUE, fs=16)

box(slide, 0.6, 4.95, 11.5, 1.2,
    "KL stiffness prediction improved most: R² went from 0.9768 (M6) → 0.9910 (Ψ-model)\n"
    "The structured prototype representation actually made the model more accurate on KL.",
    LIGHT_GREEN, fs=15)


prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
