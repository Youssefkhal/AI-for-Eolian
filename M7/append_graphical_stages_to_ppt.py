from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

SRC_PPT_PATH = r"c:\Users\youss\Downloads\PFE\M7\new_out_20260414_075825.pptx"
OUT_PPT_PATH = r"c:\Users\youss\Downloads\PFE\M7\new_out_20260414_075825_graphical_stages.pptx"

prs = Presentation(SRC_PPT_PATH)

# Theme colors
BG = RGBColor(244, 248, 252)
TITLE = RGBColor(21, 45, 74)
TEXT = RGBColor(35, 52, 70)
ACCENT_A = RGBColor(0, 112, 192)
ACCENT_B = RGBColor(46, 125, 50)
ACCENT_C = RGBColor(230, 126, 34)
ACCENT_D = RGBColor(120, 60, 170)
WHITE = RGBColor(255, 255, 255)


def setup_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.8))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TITLE

    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.55), Inches(0.95), Inches(12.0), Inches(0.45))
        sf = s.text_frame
        sp = sf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(15)
        sp.font.color.rgb = ACCENT_D

    return slide


def add_box(slide, x, y, w, h, text, color, fs=16, bold=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = RGBColor(200, 210, 220)

    tf = shp.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(fs)
    p.font.bold = bold
    # Auto-pick readable text color for both light and dark fill colors.
    luminance = 0.2126 * color[0] + 0.7152 * color[1] + 0.0722 * color[2]
    p.font.color.rgb = TEXT if luminance > 175 else WHITE
    return shp


def add_arrow(slide, x, y, w, h, color=RGBColor(130, 140, 155)):
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.color.rgb = color
    return arr


# Slide 1: Stage map
slide = setup_slide(
    "M7 Graphical Stage Map",
    "Three-stage Psi-NN pipeline applied to pile stiffness degradation"
)

add_box(slide, 0.7, 1.8, 3.5, 1.1, "Stage A\nDistillation", ACCENT_A, fs=18, bold=True)
add_arrow(slide, 4.35, 2.15, 1.0, 0.45)
add_box(slide, 5.45, 1.8, 3.5, 1.1, "Stage B\nStructure Discovery", ACCENT_B, fs=18, bold=True)
add_arrow(slide, 9.1, 2.15, 1.0, 0.45)
add_box(slide, 10.2, 1.8, 2.2, 1.1, "Stage C\nPsi Model", ACCENT_C, fs=18, bold=True)

add_box(
    slide, 0.8, 3.4, 11.8, 2.4,
    "Input: 8 geotechnical features per scenario\n"
    "Output: 21-step sequence of [KL, KR, KLR]\n\n"
    "Core idea: keep attention backbone, but replace 20 independent drop-slot parameters\n"
    "with k* prototype slots + relation matrix R for structured reconstruction.",
    RGBColor(255, 255, 255), fs=17
)

# Slide 2: Stage A graphical
slide = setup_slide("Stage A - Distillation (Teacher -> Student)", "Same slot-attention architecture, plus L1 sparsity on slots")

add_box(slide, 0.6, 1.7, 2.5, 1.1, "Inputs\n(8 features)", RGBColor(70, 130, 180), fs=15, bold=True)
add_arrow(slide, 3.2, 2.05, 0.7, 0.35)
add_box(slide, 3.95, 1.7, 3.0, 1.1, "Teacher M6\nSlotAttention", RGBColor(44, 92, 147), fs=15, bold=True)
add_arrow(slide, 7.05, 2.05, 0.7, 0.35)
add_box(slide, 7.8, 1.7, 2.6, 1.1, "Teacher Output\n(21 x 3)", RGBColor(88, 126, 168), fs=15)

add_box(slide, 0.6, 3.2, 2.5, 1.1, "Inputs\n(8 features)", RGBColor(70, 130, 180), fs=15, bold=True)
add_arrow(slide, 3.2, 3.55, 0.7, 0.35)
add_box(slide, 3.95, 3.2, 3.0, 1.1, "Student\nSlotAttention + L1", RGBColor(0, 112, 192), fs=15, bold=True)
add_arrow(slide, 7.05, 3.55, 0.7, 0.35)
add_box(slide, 7.8, 3.2, 2.6, 1.1, "Student Output\n(21 x 3)", RGBColor(88, 126, 168), fs=15)

add_box(
    slide, 0.7, 4.8, 11.8, 1.5,
    "Loss = MSE(student, teacher) + 0.5*MSE(student, target) + mu*L1(slots)\n"
    "Goal: imitate teacher while pushing slot representations toward sparse, discoverable structure.",
    RGBColor(255, 255, 255), fs=15
)

# Slide 3: Stage B graphical
slide = setup_slide("Stage B - Structure Discovery", "From refined student slots to k* prototypes and relation matrix R")

add_box(slide, 0.6, 1.7, 3.0, 1.05, "Refined Drop Slots\n(20 x 64 vectors)", RGBColor(56, 142, 60), fs=14, bold=True)
add_arrow(slide, 3.75, 2.03, 0.8, 0.35)
add_box(slide, 4.65, 1.7, 2.8, 1.05, "Cosine Similarity\nMatrix (20 x 20)", RGBColor(67, 160, 71), fs=14)
add_arrow(slide, 7.6, 2.03, 0.8, 0.35)
add_box(slide, 8.5, 1.7, 3.7, 1.05, "KMeans (k=2..10) +\nSilhouette Selection", RGBColor(76, 175, 80), fs=14)

add_box(slide, 0.6, 3.1, 5.6, 2.0, "Discovered clusters:\n- each cluster groups similar drop-step slots\n- cluster centroids become prototype candidates", RGBColor(255, 255, 255), fs=14)
add_box(slide, 6.4, 3.1, 5.8, 2.0, "Build R (20 x k*):\n- inverse-distance soft assignment\n- optional near one-hot sharpening\n- maps each drop slot to prototypes", RGBColor(255, 255, 255), fs=14)

add_box(slide, 0.7, 5.35, 11.6, 1.0, "Output of Stage B: k*, centroids, relation matrix R, and interpretable slot groups", RGBColor(220, 238, 222), fs=15, bold=True)

# Slide 4: Stage C graphical
slide = setup_slide("Stage C - Structured Psi Model", "Rewrite at slot parameterization level, not just MLP")

add_box(slide, 0.6, 1.6, 2.6, 1.0, "Initial Slot\n(1 x 64)", RGBColor(240, 173, 78), fs=14)
add_box(slide, 3.5, 1.6, 2.6, 1.0, "k* Prototypes\n(k* x 64)", RGBColor(230, 126, 34), fs=14, bold=True)
add_box(slide, 6.4, 1.6, 2.6, 1.0, "Relation Matrix R\n(20 x k*)", RGBColor(244, 166, 96), fs=14)
add_box(slide, 9.3, 1.6, 2.9, 1.0, "Reconstructed\nDrop Slots (20 x 64)", RGBColor(235, 152, 78), fs=14)

add_arrow(slide, 6.1, 2.0, 0.25, 0.25)
add_arrow(slide, 8.95, 2.0, 0.25, 0.25)

add_box(slide, 0.8, 3.1, 11.4, 1.2, "Concatenate [initial + reconstructed drops] -> 21 slots -> same iterative attention refinement", RGBColor(255, 255, 255), fs=15)
add_box(slide, 0.8, 4.5, 11.4, 1.2, "Heads remain: initial_proj + drop_proj + physics constraints (KL/KR negative drops, KLR positive drops)", RGBColor(255, 255, 255), fs=15)

add_box(slide, 0.8, 5.9, 11.4, 0.6, "So the major rewrite is the slot structure (prototypes + R), not only the MLP heads.", RGBColor(252, 234, 220), fs=15, bold=True)

# Slide 5: end-to-end summary table style
slide = setup_slide("What Changed vs What Stayed", "Exact architecture interpretation")

add_box(slide, 0.8, 1.7, 5.8, 0.8, "Changed in M7 Stage C", RGBColor(226, 239, 218), fs=16, bold=True)
add_box(slide, 6.9, 1.7, 5.2, 0.8, "Unchanged backbone", RGBColor(221, 235, 247), fs=16, bold=True)

add_box(slide, 0.8, 2.7, 5.8, 3.6,
        "- 20 independent drop slots -> k* prototype slots\n"
        "- Added relation matrix R (20 x k*)\n"
        "- Reconstruct drop slots from prototypes\n"
        "- Added per-slot scale factors\n"
        "- Structure discovered from clustering",
        RGBColor(255, 255, 255), fs=15)

add_box(slide, 6.9, 2.7, 5.2, 3.6,
        "- Input embedding\n"
        "- Cross-attention\n"
        "- Self-attention\n"
        "- Slot MLP refinement\n"
        "- Output heads and physics constraints",
        RGBColor(255, 255, 255), fs=15)

add_box(slide, 0.8, 6.45, 11.3, 0.6,
        "Conclusion: M7 is a structural slot rewrite with the same attention family, not an MLP-only rewrite.",
        RGBColor(248, 232, 246), fs=14, bold=True)

prs.save(OUT_PPT_PATH)
print(f"Updated presentation: {OUT_PPT_PATH}")
print(f"Total slides now: {len(prs.slides)}")
