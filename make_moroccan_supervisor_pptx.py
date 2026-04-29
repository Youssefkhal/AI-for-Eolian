"""
Présentation PPTX pour la réunion avec le superviseur marocain.
Stage à Shenzhen University, Chine.
Étudiant: Youssef KHALLOUQI — École Marocaine des Sciences de l'Ingénieur (EMSI)

Contenu (22 diapositives, en français):
  1.  Page de titre
  2.  Sommaire
  3.  Université d'accueil — Shenzhen University
  4.  Lettre d'acceptation (邀请函) & Admission doctorale
  5.  Cadre de supervision & collaboration
  6.  Sujet de recherche — Pourquoi c'est important
  7.  Problème d'ingénierie : dégradation des pieux OWT
  8.  Données & paramètres physiques
  9.  Approche IA : Jumeau Numérique Informé par la Physique
  10. Architecture du modèle — Slot Attention Transformer
  11. Évolution du modèle : M5 → M6 → M7 → M8 → M9
  12. M9 : SwiGLU Psi-NN — État de l'art actuel
  13. Pipeline en 3 étapes (Distillation → Découverte → Psi-Modèle)
  14. Contraintes physiques intégrées dans l'IA
  15. Résultats & métriques de performance
  16. Comparaison complète M6 → M9
  17. Application web interactive
  18. Contributions scientifiques & potentiel de publication
  19. Contexte du projet NSFC (面上项目)
  20. Chronologie & plan de travail
  21. Conclusion & prochaines étapes
  22. Merci / Questions
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

# --- Design tokens ---
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD    = RGBColor(0x1A, 0x24, 0x3B)
BG_ALT     = RGBColor(0x14, 0x1E, 0x33)
ACCENT1    = RGBColor(0x00, 0xD2, 0xFF)  # cyan
ACCENT2    = RGBColor(0xFF, 0x6B, 0x6B)  # corail
ACCENT3    = RGBColor(0x4E, 0xC9, 0xB0)  # emeraude
ACCENT4    = RGBColor(0xFF, 0xAA, 0x33)  # ambre
ACCENT5    = RGBColor(0xA8, 0x78, 0xFF)  # violet
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT      = RGBColor(0xCC, 0xCC, 0xDD)
DIM        = RGBColor(0x88, 0x88, 0x99)
RED        = RGBColor(0xFF, 0x44, 0x44)
GREEN      = RGBColor(0x00, 0xFF, 0x88)
GOLD       = RGBColor(0xFF, 0xD7, 0x00)
SZU_RED    = RGBColor(0xCC, 0x00, 0x00)  # Shenzhen University red
HITSZ_BLUE = RGBColor(0x00, 0x33, 0x99)  # HITSZ blue
MOROCCO_R  = RGBColor(0xC1, 0x27, 0x2D)  # rouge marocain
MOROCCO_G  = RGBColor(0x00, 0x6B, 0x3F)  # vert marocain

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, colour=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_rect(slide, left, top, width, height, fill_colour, border_colour=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    if border_colour:
        shape.line.color.rgb = border_colour
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 colour=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = colour
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=13,
                  colour=LIGHT, spacing=Pt(6), bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, tuple):
            p.text = line[0]
            p.font.color.rgb = line[1]
        else:
            p.text = line
            p.font.color.rgb = colour
        p.font.size = Pt(font_size)
        p.font.name = 'Calibri'
        p.font.bold = bold
        p.space_after = spacing
        p.alignment = alignment
    return txBox


def add_tag(slide, left, top, text, bg_colour=ACCENT1, text_colour=BG_DARK, font_size=10, width=None):
    w = width or Inches(2.0)
    h = Inches(0.32)
    rect = add_rect(slide, left, top, w, h, bg_colour)
    rect.text_frame.paragraphs[0].text = text
    rect.text_frame.paragraphs[0].font.size = Pt(font_size)
    rect.text_frame.paragraphs[0].font.bold = True
    rect.text_frame.paragraphs[0].font.color.rgb = text_colour
    rect.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    rect.text_frame.paragraphs[0].font.name = 'Calibri'
    return rect


def add_accent_bar(slide, colour=ACCENT1):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), colour)


def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(12.0), Inches(7.1), Inches(1.2), Inches(0.3),
                 f'{num}/{total}', font_size=9, colour=DIM, alignment=PP_ALIGN.RIGHT)


# ======================================================
#  DIAPOSITIVES
# ======================================================
TOTAL_SLIDES = 22


def slide_01_titre(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(4.44), Inches(0.08), MOROCCO_R)
    add_rect(slide, Inches(4.44), Inches(0), Inches(4.44), Inches(0.08), MOROCCO_G)
    add_rect(slide, Inches(8.88), Inches(0), Inches(4.45), Inches(0.08), SZU_RED)

    add_text_box(slide, Inches(1), Inches(1.3), Inches(11.3), Inches(1.2),
                 "Presentation du Stage de Recherche", font_size=38, colour=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(2.7), Inches(11.3), Inches(1.0),
                 "Jumeau Numerique base sur l'IA Informee par la Physique\n"
                 "pour la Prediction de la Degradation de Rigidite\n"
                 "des Fondations sur Pieux des Eoliennes Offshore",
                 font_size=19, colour=ACCENT1, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
                 "Stage a Shenzhen University (\u6df1\u5733\u5927\u5b66), Shenzhen, Chine",
                 font_size=16, colour=ACCENT4, alignment=PP_ALIGN.CENTER)

    add_rect(slide, Inches(3.0), Inches(5.0), Inches(7.3), Inches(1.8), BG_CARD, ACCENT1)
    add_text_box(slide, Inches(3.3), Inches(5.1), Inches(6.7), Inches(0.4),
                 'Youssef KHALLOUQI', font_size=20, colour=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.3), Inches(5.6), Inches(6.7), Inches(0.4),
                 "Ecole Marocaine des Sciences de l'Ingenieur (EMSI)",
                 font_size=14, colour=ACCENT3, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.3), Inches(6.05), Inches(6.7), Inches(0.4),
                 "Collaboration EMSI (Maroc) \u2194 Shenzhen University (Chine)",
                 font_size=12, colour=LIGHT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.3), Inches(6.4), Inches(6.7), Inches(0.3),
                 "Reunion avec le superviseur -- Avril 2026",
                 font_size=11, colour=DIM, alignment=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0), Inches(7.3), Inches(13.333), Inches(0.06), ACCENT3)


def slide_02_sommaire(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT1)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.6),
                 'Sommaire', font_size=28, colour=ACCENT1, bold=True)

    sections = [
        ('01', "Universite d'accueil & Acceptation", ACCENT1),
        ('02', 'Programme doctoral & Cadre de collaboration', ACCENT5),
        ('03', 'Sujet de recherche : Enjeux et importance', ACCENT2),
        ('04', "Probleme d'ingenierie : Degradation des pieux", ACCENT4),
        ('05', 'Approche IA : Jumeau Numerique', ACCENT3),
        ('06', 'Architecture & Evolution du modele (M5 a M9)', ACCENT1),
        ('07', 'Innovation cle : Pipeline SwiGLU Psi-NN', GREEN),
        ('08', 'Resultats, Validation & Application web', ACCENT4),
        ('09', 'Contributions scientifiques & Chronologie', ACCENT5),
        ('10', 'Conclusion & Prochaines etapes', ACCENT2),
    ]

    for i, (num, title, clr) in enumerate(sections):
        col = i // 5
        row = i % 5
        x = Inches(0.5) + col * Inches(6.4)
        y = Inches(1.2) + row * Inches(1.15)
        add_rect(slide, x, y, Inches(6.0), Inches(0.95), BG_CARD, clr)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.6), Inches(0.5),
                     num, font_size=20, colour=clr, bold=True)
        add_text_box(slide, x + Inches(0.85), y + Inches(0.2), Inches(4.8), Inches(0.5),
                     title, font_size=15, colour=WHITE)

    add_slide_number(slide, 2, TOTAL_SLIDES)


def slide_03_universite(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, SZU_RED)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Universite d'accueil -- Shenzhen University (\u6df1\u5733\u5927\u5b66)", font_size=26,
                 colour=WHITE, bold=True)

    # Colonne gauche : infos universite
    add_rect(slide, Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.7), BG_CARD, SZU_RED)
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.4), Inches(0.4),
                 "A propos de Shenzhen University", font_size=16, colour=GOLD, bold=True)

    uni_info = [
        "-- Fondee en 1983 a Shenzhen, province du Guangdong, Chine",
        "-- Shenzhen : capitale technologique de la Chine (siege de Huawei,",
        "   Tencent, BYD, DJI, ZTE...)",
        "",
        "-- Universite en croissance rapide, classee parmi les meilleures",
        "   universites chinoises (Double First-Class discipline list)",
        "-- Plus de 40 000 etudiants, 2 800+ enseignants-chercheurs",
        "",
        "-- Points forts : Ingenierie, Informatique, Science des materiaux",
        "-- Forte orientation vers l'innovation et l'industrie",
        "-- Collaborations internationales avec plus de 200 universites",
        "",
        "-- Environnement ideal : Shenzhen est un hub technologique mondial",
        "   avec un ecosysteme de recherche tres dynamique",
    ]
    add_multiline(slide, Inches(0.8), Inches(1.85), Inches(5.4), Inches(4.8),
                  uni_info, font_size=12, colour=LIGHT, spacing=Pt(4))

    # Colonne droite : labo de recherche
    add_rect(slide, Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.7), BG_CARD, ACCENT3)
    add_text_box(slide, Inches(7.1), Inches(1.3), Inches(5.4), Inches(0.4),
                 'Laboratoire de recherche', font_size=16, colour=ACCENT3, bold=True)

    lab_info = [
        "-- Departement : Genie Civil / Architecture",
        "",
        "-- Axes de recherche :",
        "   -- Fondations d'eoliennes offshore (OWT)",
        "   -- Interaction sol-structure sous chargement cyclique",
        "   -- IA / Machine Learning pour le suivi structurel",
        "   -- Reseaux de neurones informes par la physique",
        "   -- Technologie de Jumeau Numerique (Digital Twin)",
        "",
        "-- Equipements :",
        "   -- Laboratoire de tests geotechniques",
        "   -- Cluster de calcul haute performance pour l'IA",
        "",
        "-- Programmes doctoraux internationaux actifs",
        "-- Financement et bourses disponibles",
    ]
    add_multiline(slide, Inches(7.1), Inches(1.85), Inches(5.4), Inches(4.8),
                  lab_info, font_size=12, colour=LIGHT, spacing=Pt(4))

    add_slide_number(slide, 3, TOTAL_SLIDES)


def slide_04_acceptation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT4)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Lettre d'acceptation (\u9080\u8bf7\u51fd) & Admission doctorale", font_size=26,
                 colour=ACCENT4, bold=True)

    # Gauche : lettre d'acceptation
    add_rect(slide, Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.7), BG_CARD, ACCENT4)
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.4), Inches(0.4),
                 "Lettre d'invitation / Acceptation", font_size=16, colour=ACCENT4, bold=True)

    accept_info = [
        "-- Lettre officielle d'invitation (\u9080\u8bf7\u51fd) recue de",
        "   Shenzhen University, Departement de Genie Civil",
        "",
        "-- La lettre confirme :",
        "   [OK] Acceptation en tant que chercheur visiteur / doctorant",
        "   [OK] Acces aux installations et ressources informatiques",
        "   [OK] Encadrement par le groupe de recherche offshore",
        "   [OK] Sujet : Prediction par IA de la degradation de rigidite",
        "        des fondations sur pieux sous conditions de typhon",
        "",
        "-- Duree : Annee academique complete",
        "-- Statut : Officiellement accepte et confirme",
        "",
        "-- Document de reference : \u9080\u8bf7\u51fd (1).pdf",
    ]
    add_multiline(slide, Inches(0.8), Inches(1.85), Inches(5.4), Inches(4.8),
                  accept_info, font_size=12, colour=LIGHT, spacing=Pt(4))

    # Droite : admission doctorale HITSZ
    add_rect(slide, Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.7), BG_CARD, HITSZ_BLUE)
    add_text_box(slide, Inches(7.1), Inches(1.3), Inches(5.4), Inches(0.4),
                 "Doctorat -- HITSZ (Rentree 2026)", font_size=16, colour=ACCENT1, bold=True)

    phd_info = [
        "-- Accepte au Harbin Institute of Technology, Shenzhen",
        "   (HITSZ -- \u54c8\u5c14\u6ee8\u5de5\u4e1a\u5927\u5b66\u6df1\u5733)",
        "",
        ("   #3 MONDIAL en Ingenierie (US News 2025-2026)", GOLD),
        "",
        "-- College : College of Artificial Intelligence",
        "-- Specialite : Control Science and Engineering",
        "-- Programme : Ph.D. (doctorat complet)",
        "-- Langue : Anglais",
        "",
        ("-- Bourse HITSZ :", ACCENT3),
        ("   100% frais de scolarite exoneres", GREEN),
        ("   + 50 000 RMB/an (~6 400 EUR/an)", GREEN),
        "",
        "-- Inscription : 27-28 aout 2026",
        "",
        ("-- IMPORTANT : Soutenance de stage requise", ACCENT2),
        ("   debut septembre pour finaliser l'inscription", ACCENT2),
    ]
    add_multiline(slide, Inches(7.1), Inches(1.85), Inches(5.4), Inches(4.8),
                  phd_info, font_size=12, colour=LIGHT, spacing=Pt(3))

    add_slide_number(slide, 4, TOTAL_SLIDES)


def slide_05_encadrement(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT3)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 'Cadre de supervision & collaboration', font_size=26, colour=ACCENT3, bold=True)

    # Bandeau principal
    add_rect(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(2.6), BG_CARD, ACCENT3)
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.4),
                 "Modele de co-encadrement", font_size=18, colour=ACCENT3, bold=True)

    framework = [
        "-- Ce stage s'inscrit dans une collaboration bilaterale entre le Maroc et la Chine.",
        "-- L'etudiant (Youssef KHALLOUQI) effectue son stage a Shenzhen University sous encadrement chinois,",
        "   tout en maintenant un suivi academique avec le superviseur marocain a l'EMSI.",
        "-- Suite au stage, il a ete accepte en doctorat a HITSZ (Harbin Institute of Technology, Shenzhen),",
        "   classee #3 mondiale en Ingenierie (US News 2025-2026), avec une bourse complete.",
        "-- La soutenance de stage est prevue debut septembre 2026 pour permettre l'inscription a HITSZ.",
    ]
    add_multiline(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(1.8),
                  framework, font_size=12, colour=LIGHT, spacing=Pt(4))

    # Deux colonnes : Maroc / Chine
    add_rect(slide, Inches(0.5), Inches(4.0), Inches(5.9), Inches(3.0), BG_CARD, MOROCCO_R)
    add_text_box(slide, Inches(0.8), Inches(4.1), Inches(5.3), Inches(0.4),
                 "Maroc -- Etablissement d'origine", font_size=16, colour=MOROCCO_R, bold=True)
    moroccan = [
        "-- Ecole : EMSI (Ecole Marocaine des Sciences de l'Ingenieur)",
        "-- Role : Cadre academique & comite de these",
        "-- Responsabilites :",
        "   -- Validation de l'orientation de la recherche",
        "   -- Evaluation de la these & soutenance",
        "   -- Administration academique",
        "   -- Suivi regulier de l'avancement",
    ]
    add_multiline(slide, Inches(0.8), Inches(4.6), Inches(5.3), Inches(2.2),
                  moroccan, font_size=12, colour=LIGHT, spacing=Pt(3))

    add_rect(slide, Inches(6.8), Inches(4.0), Inches(6.0), Inches(3.0), BG_CARD, SZU_RED)
    add_text_box(slide, Inches(7.1), Inches(4.1), Inches(5.4), Inches(0.4),
                 "Chine -- Etablissement d'accueil", font_size=16, colour=GOLD, bold=True)
    chinese = [
        "-- Universite : Shenzhen University (\u6df1\u5733\u5927\u5b66)",
        "-- Role : Execution de la recherche & acces au labo",
        "-- Responsabilites :",
        "   -- Encadrement quotidien de la recherche",
        "   -- Acces aux donnees, au calcul, aux laboratoires",
        "   -- Co-autorat des articles de recherche",
        "   -- Formation aux methodes avancees d'IA",
    ]
    add_multiline(slide, Inches(7.1), Inches(4.6), Inches(5.4), Inches(2.2),
                  chinese, font_size=12, colour=LIGHT, spacing=Pt(3))

    add_slide_number(slide, 5, TOTAL_SLIDES)


def slide_06_sujet(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT2)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Sujet de recherche -- Pourquoi c'est important", font_size=26, colour=ACCENT2, bold=True)

    add_text_box(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.5),
                 "Prediction de la degradation de rigidite des fondations sur pieux "
                 "sous chargement cyclique induit par les typhons, a l'aide de l'IA informee par la physique",
                 font_size=14, colour=ACCENT1, bold=True)

    cols = [
        ("Contexte mondial", ACCENT4, [
            "L'eolien offshore est en pleine expansion mondiale",
            "La Chine prevoit 100+ GW d'eolien offshore d'ici 2030",
            "Le Maroc developpe son propre potentiel eolien",
            "La plupart des turbines utilisent des monopieux (5m dia.)",
            "Les typhons creent un chargement cyclique extreme",
            "L'integrite des fondations est cruciale pour la securite",
        ]),
        ("Le defi", ACCENT2, [
            "Lors des typhons, le sol autour du pieu se degrade",
            "La rigidite (resistance au deplacement) diminue",
            "Si la rigidite baisse trop : resonance puis rupture !",
            "Methodes actuelles : tests physiques couteux ou MEF lent",
            "Aucun outil rapide et precis n'existe en temps reel",
            "Les ingenieurs ont besoin d'un simulateur 'what-if'",
        ]),
        ("Notre solution", ACCENT3, [
            "Construire un 'Jumeau Numerique' IA du systeme pieu-sol",
            "L'IA apprend a partir de donnees experimentales reelles",
            "Elle encode les lois physiques fondamentales dans le reseau",
            "Predictions instantanees (millisecondes vs heures)",
            "Les ingenieurs peuvent tester des milliers de scenarios",
            "Le modele est explicable -- pas une boite noire",
        ]),
    ]

    for i, (title, clr, bullets) in enumerate(cols):
        x = Inches(0.4) + i * Inches(4.2)
        add_rect(slide, x, Inches(1.6), Inches(4.0), Inches(5.3), BG_CARD, clr)
        add_text_box(slide, x + Inches(0.2), Inches(1.7), Inches(3.6), Inches(0.4),
                     title, font_size=15, colour=clr, bold=True)
        add_multiline(slide, x + Inches(0.2), Inches(2.2), Inches(3.6), Inches(4.3),
                      [f"-- {b}" for b in bullets], font_size=11, colour=LIGHT, spacing=Pt(6))

    add_slide_number(slide, 6, TOTAL_SLIDES)


def slide_07_probleme(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT4)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Probleme d'ingenierie : Degradation de rigidite des pieux",
                 font_size=26, colour=ACCENT4, bold=True)

    add_rect(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(2.1), BG_CARD, ACCENT4)
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.4),
                 "Que se passe-t-il lors d'un typhon ?", font_size=16, colour=ACCENT4, bold=True)

    desc = [
        "1. Le vent et les vagues creent des forces cycliques laterales sur l'eolienne offshore.",
        "2. Ces forces se transmettent par la tour jusqu'au monopieu dans le fond marin.",
        "3. Le sol environnant est charge et decharge de maniere repetee --> il s'affaiblit progressivement.",
        "4. Cet affaiblissement se mesure par 3 parametres de rigidite qui evoluent dans le temps :",
    ]
    add_multiline(slide, Inches(0.8), Inches(1.65), Inches(11.5), Inches(1.3),
                  desc, font_size=12, colour=LIGHT, spacing=Pt(4))

    params = [
        ("KL -- Rigidite laterale", "Resistance au deplacement lateral", "Diminue", ACCENT2,
         "KL baisse (le sol s'affaiblit lateralement)"),
        ("KR -- Rigidite rotationnelle", "Resistance au basculement", "Diminue", ACCENT1,
         "KR baisse (le sol resiste moins au basculement)"),
        ("KLR -- Couplage croise", "Interaction laterale-rotationnelle", "Augmente (vers 0)", ACCENT3,
         "KLR est negatif, il augmente vers zero"),
    ]

    for i, (name, desc_txt, trend, clr, detail) in enumerate(params):
        x = Inches(0.5) + i * Inches(4.2)
        add_rect(slide, x, Inches(3.5), Inches(3.9), Inches(3.3), BG_CARD, clr)
        add_text_box(slide, x + Inches(0.15), Inches(3.6), Inches(3.6), Inches(0.4),
                     name, font_size=13, colour=clr, bold=True)
        add_text_box(slide, x + Inches(0.15), Inches(4.0), Inches(3.6), Inches(0.35),
                     desc_txt, font_size=11, colour=LIGHT)
        add_tag(slide, x + Inches(0.15), Inches(4.4), trend, clr, BG_DARK, width=Inches(1.8))
        add_text_box(slide, x + Inches(0.15), Inches(4.9), Inches(3.6), Inches(0.8),
                     detail, font_size=10, colour=DIM)
        add_text_box(slide, x + Inches(0.15), Inches(5.7), Inches(3.6), Inches(0.8),
                     "Code dans le modele IA\n(impossible a violer)",
                     font_size=10, colour=GREEN, bold=True)

    add_slide_number(slide, 7, TOTAL_SLIDES)


def slide_08_donnees(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT1)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Donnees & parametres physiques d'entree", font_size=26, colour=ACCENT1, bold=True)

    add_rect(slide, Inches(0.5), Inches(1.1), Inches(7.5), Inches(5.8), BG_CARD, ACCENT1)
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(7.0), Inches(0.4),
                 "8 parametres d'entree (proprietes sol & pieu)", font_size=16, colour=ACCENT1, bold=True)

    features = [
        ('PI', 'Indice de plasticite', "Classification du sol (argilosite)"),
        ('Gmax', 'Module de cisaillement max', "Rigidite initiale du sol (Pa)"),
        ('v', 'Coefficient de Poisson', "Rapport deformation lat./axiale"),
        ('Dp', 'Diametre du pieu', "Diametre du monopieu (m)"),
        ('Tp', "Epaisseur de la paroi", "Epaisseur de la paroi du pieu (m)"),
        ('Lp', 'Longueur encastree', "Longueur du pieu dans le fond marin (m)"),
        ('Ip', "Moment d'inertie", "Resistance de la section a la flexion"),
        ('Dp/Lp', 'Rapport diametre/longueur', "Rapport d'aspect du pieu"),
    ]

    for i, (sym, name, desc_txt) in enumerate(features):
        y = Inches(1.75) + i * Inches(0.58)
        bg = BG_CARD if i % 2 == 0 else BG_ALT
        add_rect(slide, Inches(0.7), y, Inches(7.0), Inches(0.5), bg)
        add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(0.8), Inches(0.35),
                     sym, font_size=12, colour=ACCENT1, bold=True)
        add_text_box(slide, Inches(1.8), y + Inches(0.05), Inches(2.5), Inches(0.35),
                     name, font_size=11, colour=WHITE)
        add_text_box(slide, Inches(4.5), y + Inches(0.05), Inches(3.0), Inches(0.35),
                     desc_txt, font_size=10, colour=DIM)

    # Droite : resume des donnees
    add_rect(slide, Inches(8.3), Inches(1.1), Inches(4.5), Inches(5.8), BG_CARD, ACCENT3)
    add_text_box(slide, Inches(8.6), Inches(1.2), Inches(4.0), Inches(0.4),
                 "Resume du jeu de donnees", font_size=16, colour=ACCENT3, bold=True)

    data_summary = [
        ('Scenarios totaux', '1 019', ACCENT1),
        ('Jeu d\'entrainement', '815 (80%)', ACCENT3),
        ('Jeu de test', '204 (20%)', ACCENT2),
        ('Pas de temps', '21 par scenario', ACCENT4),
        ('Variables de sortie', '3 (KL, KR, KLR)', WHITE),
        ('Forme de sortie', '[B, 21, 3]', DIM),
    ]

    for i, (label, value, clr) in enumerate(data_summary):
        y = Inches(1.8) + i * Inches(0.65)
        add_text_box(slide, Inches(8.6), y, Inches(2.3), Inches(0.35),
                     label, font_size=12, colour=LIGHT)
        add_text_box(slide, Inches(10.8), y, Inches(1.8), Inches(0.35),
                     value, font_size=14, colour=clr, bold=True, alignment=PP_ALIGN.RIGHT)

    add_text_box(slide, Inches(8.6), Inches(5.8), Inches(4.0), Inches(0.8),
                 "Source : Donnees experimentales reelles\nde tests cycliques sur monopieux\n"
                 "(REAL DATA.xlsx -- 44 cycles par scenario,\nsous-echantillonnes en 21 pas representatifs)",
                 font_size=10, colour=DIM)

    add_slide_number(slide, 8, TOTAL_SLIDES)


def slide_09_approche_ia(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT3)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Approche IA -- Jumeau Numerique Informe par la Physique",
                 font_size=26, colour=ACCENT3, bold=True)

    add_rect(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.5), BG_CARD, ACCENT3)
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.4),
                 "Concept central : Combiner la puissance de l'IA avec les connaissances en ingenierie",
                 font_size=15, colour=ACCENT3, bold=True)
    add_multiline(slide, Inches(0.8), Inches(1.65), Inches(11.5), Inches(0.7),
                  [
                      "Contrairement a l'IA 'boite noire', notre modele integre les lois physiques fondamentales directement dans l'architecture.",
                      "Le modele NE PEUT PAS violer ces lois -- ce sont des contraintes structurelles, pas de simples penalites.",
                  ], font_size=12, colour=LIGHT, spacing=Pt(3))

    pillars = [
        ("Slot Attention\nTransformer", ACCENT1, [
            "21 'slots' memoire apprenables",
            "Cross-attention vers les donnees",
            "Self-attention entre les slots",
            "Raffinement iteratif (3 tours)",
            "Chaque slot se specialise sur un",
            "pas de degradation temporel",
        ]),
        ("Contraintes\nphysiques", ACCENT4, [
            "Chutes de KL, KR forcees <= 0",
            "Chutes de KLR forcees >= 0",
            "Somme cumulative garantit",
            "des trajectoires monotones",
            "Echelle log-signee preserve",
            "l'info de signe (KLR < 0)",
        ]),
        ("Decouverte de\nstructure Psi-NN", ACCENT5, [
            "20 slots --> k* prototypes",
            "Clustering K-Means sur les slots",
            "Matrice de relation R relie",
            "les pas aux prototypes",
            "Reduction de 19.7% des parametres",
            "Meilleure interpretabilite",
        ]),
    ]

    for i, (title, clr, bullets) in enumerate(pillars):
        x = Inches(0.4) + i * Inches(4.2)
        add_rect(slide, x, Inches(2.9), Inches(4.0), Inches(4.0), BG_CARD, clr)
        add_text_box(slide, x + Inches(0.15), Inches(3.0), Inches(3.7), Inches(0.7),
                     title, font_size=14, colour=clr, bold=True)
        add_multiline(slide, x + Inches(0.15), Inches(3.7), Inches(3.7), Inches(2.8),
                      [f"-- {b}" for b in bullets], font_size=11, colour=LIGHT, spacing=Pt(5))

    add_slide_number(slide, 9, TOTAL_SLIDES)


def slide_10_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT1)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Architecture du modele -- Slot Attention Transformer",
                 font_size=26, colour=ACCENT1, bold=True)

    flow_steps = [
        ('ENTREE', '8 parametres : PI, Gmax, v, Dp, Tp, Lp, Ip, Dp/Lp', ACCENT2, '[B, 8]'),
        ('EMBEDDING', 'Linear(8->64) + LayerNorm + GELU', ACCENT5, '[B, 1, 64]'),
        ('21 SLOTS APPRENABLES', 'Slot 1 (initial) + Slots 2-21 (degradation)', ACCENT4, '[B, 21, 64]'),
        ('RAFFINEMENT ITERATIF x3', 'Cross-Attn -> Self-Attn -> SwiGLU MLP', ACCENT3, '[B, 21, 64]'),
        ('SEPARATION DES SLOTS', 'Slot 1 -> Tete initiale | Slots 2-21 -> Tete de chute', ACCENT1, 'Split'),
        ('CONTRAINTES PHYSIQUES', 'dKL,dKR = -|d|  |  dKLR = +|d|', ACCENT2, 'Codees en dur'),
        ('SOMME CUMULATIVE', 'K(t) = K0 + somme(chutes contraintes)', DIM, '[B, 21, 3]'),
        ('SORTIE', 'KL(t), KR(t), KLR(t) pour t=1...21', GREEN, '[B, 21, 3]'),
    ]

    for i, (name, desc_txt, clr, dim) in enumerate(flow_steps):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.5)
        y = Inches(1.1) + row * Inches(1.5)
        add_rect(slide, x, y, Inches(6.0), Inches(1.2), BG_CARD, clr)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.05), Inches(4.0), Inches(0.35),
                     name, font_size=13, colour=clr, bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.4), Inches(4.5), Inches(0.35),
                     desc_txt, font_size=10, colour=LIGHT)
        add_text_box(slide, x + Inches(4.5), y + Inches(0.05), Inches(1.3), Inches(0.35),
                     dim, font_size=9, colour=DIM, alignment=PP_ALIGN.RIGHT)

    add_slide_number(slide, 10, TOTAL_SLIDES)


def slide_11_evolution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT4)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Evolution du modele : M5 -> M6 -> M7 -> M8 -> M9",
                 font_size=26, colour=ACCENT4, bold=True)

    versions = [
        ('M5', 'LSTM + Slot Attn', '157 958', '0.9913', 'Base avec decodeur LSTM', ACCENT2,
         'Decodeur LSTM | 2 iterations | KLR non contraint'),
        ('M6', 'Enseignant simplifie', '56 646', '0.9804', 'Suppression du LSTM, ajout contrainte KLR', ACCENT1,
         'Sans LSTM (64% de params en -) | 3 iterations | 3 rigidites contraintes'),
        ('M7', 'Psi-NN de base', '55 770', '0.9897', 'Premiere tentative de compression', RGBColor(0x64, 0xB5, 0xF6),
         'Pipeline 3 etapes | Matrice R fixe | k*=6, MLP large'),
        ('M8', 'Psi-NN efficace', '45 486', '0.9882', 'MLP goulot + R apprenable', RGBColor(0xFF, 0xB7, 0x4D),
         'MLP goulot (64->48->64) | R apprenable (softmax) | Perte physique'),
        ('M9', 'SwiGLU Psi-NN', '45 502', '0.9891', 'Meilleur : mecanisme de gating SwiGLU', GREEN,
         'SwiGLU (gate x value) | Meme nb params que M8 | R2 superieur'),
    ]

    for i, (name, kind, params, r2, tagline, clr, details) in enumerate(versions):
        y = Inches(1.1) + i * Inches(1.22)
        is_best = name == 'M9'
        add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(1.1),
                 BG_CARD, GREEN if is_best else clr)

        add_tag(slide, Inches(0.65), y + Inches(0.08), name, clr, BG_DARK, width=Inches(0.8))
        add_text_box(slide, Inches(1.6), y + Inches(0.05), Inches(2.5), Inches(0.35),
                     kind, font_size=13, colour=WHITE, bold=True)
        add_text_box(slide, Inches(1.6), y + Inches(0.4), Inches(3.0), Inches(0.35),
                     tagline, font_size=10, colour=DIM)
        add_text_box(slide, Inches(4.5), y + Inches(0.1), Inches(1.5), Inches(0.35),
                     params, font_size=14, colour=clr, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(4.5), y + Inches(0.45), Inches(1.5), Inches(0.25),
                     'params', font_size=8, colour=DIM, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(6.2), y + Inches(0.1), Inches(1.2), Inches(0.35),
                     f'R2={r2}', font_size=13, colour=GREEN if is_best else ACCENT3,
                     bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(7.6), y + Inches(0.15), Inches(5.0), Inches(0.7),
                     details, font_size=9, colour=LIGHT)

    add_slide_number(slide, 11, TOTAL_SLIDES)


def slide_12_m9_detail(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, GREEN)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "M9 : SwiGLU Psi-NN -- Etat de l'art actuel", font_size=26, colour=GREEN, bold=True)

    add_rect(slide, Inches(0.5), Inches(1.1), Inches(7.0), Inches(5.8), BG_CARD, GREEN)
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(6.4), Inches(0.4),
                 "Qu'est-ce que SwiGLU ? (L'innovation cle)", font_size=16, colour=GREEN, bold=True)

    swiglu_desc = [
        "SwiGLU = Swish-Gated Linear Unit",
        "Utilise dans les meilleurs modeles d'IA : LLaMA, PaLM, Mistral",
        "",
        "Fonctionnement (dans chaque bloc feed-forward du slot) :",
        "",
        "   Entree x (vecteur 64 dimensions)",
        "      |                 |",
        "   W_gate . x       W_val . x",
        "      |                 |",
        "   SiLU(.)          (valeur brute)",
        "      |                 |",
        "   porte    x    valeur   <-- produit element par element",
        "      |",
        "   W_out . (filtre) --> Sortie (64 dims)",
        "",
        "Pourquoi c'est mieux :",
        "  -- La porte apprend QUELLES caracteristiques propager",
        "  -- Plus expressif que le simple GELU au meme nombre de params",
        "  -- 6 272 params (= M8 avec 6 256)",
        "  -- Mais R2 superieur (0.9891 vs 0.9882)",
    ]
    add_multiline(slide, Inches(0.8), Inches(1.7), Inches(6.4), Inches(5.0),
                  swiglu_desc, font_size=11, colour=LIGHT, spacing=Pt(3))

    add_rect(slide, Inches(7.8), Inches(1.1), Inches(5.0), Inches(5.8), BG_CARD, ACCENT1)
    add_text_box(slide, Inches(8.1), Inches(1.2), Inches(4.4), Inches(0.4),
                 "Performances de M9", font_size=16, colour=ACCENT1, bold=True)

    stats = [
        ('Parametres', '45 502', GREEN),
        ('Compression vs M6', '-19.7%', ACCENT4),
        ('Prototypes k*', '5', ACCENT1),
        ('Dim. cachee SwiGLU', '32', GREEN),
        ('', '', DIM),
        ('R2 global', '0.9891', GREEN),
        ('R2 KL', '0.9905', WHITE),
        ('R2 KR', '0.9796', WHITE),
        ('R2 KLR', '0.9860', WHITE),
        ('', '', DIM),
        ('MAPE courbe', '2.37%', ACCENT4),
        ('R2 courbe', '0.9829', ACCENT3),
        ('Scenarios de test', '204', ACCENT1),
    ]

    for i, (label, value, clr) in enumerate(stats):
        if not label:
            continue
        y = Inches(1.7) + i * Inches(0.38)
        add_text_box(slide, Inches(8.1), y, Inches(2.5), Inches(0.3),
                     label, font_size=11, colour=DIM)
        add_text_box(slide, Inches(10.6), y, Inches(2.0), Inches(0.3),
                     value, font_size=13, colour=clr, bold=True, alignment=PP_ALIGN.RIGHT)

    add_slide_number(slide, 12, TOTAL_SLIDES)


def slide_13_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT5)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Pipeline d'entrainement Psi-NN en 3 etapes", font_size=26, colour=ACCENT5, bold=True)

    stages = [
        ('ETAPE A', 'Distillation', ACCENT2, [
            "L'enseignant M6 (56 646 params) est gele",
            "L'eleve apprend a imiter la sortie du prof",
            "MLP SwiGLU dans l'eleve (upgrade M9)",
            "Regularisation L1 -> patterns de slots creux",
            "Perte = Distill + 0.5*Data + mu*L1",
            "2 000 epoques, patience=150",
        ]),
        ('ETAPE B', 'Decouverte de structure', ACCENT4, [
            "Extraire les 20 vecteurs de slots d'abandon",
            "Moyenne sur tous les scenarios d'entrainement",
            "Clustering K-Means (k=2..10)",
            "Score silhouette -> k* optimal = 5",
            "Construction de la matrice R (20x5)",
            "R relie chaque pas aux prototypes",
        ]),
        ('ETAPE C', 'Entrainement Psi-Modele', ACCENT3, [
            "Construction du Psi-modele avec k*=5 prototypes",
            "Initialisation depuis les centroides de l'etape B",
            "R apprenable : softmax(logits)",
            "Perte multi-objectif :",
            "  Distill + Data + Initial*5 + Forme",
            "  + 0.02*Entropie + 0.2*PhysiqueMono",
        ]),
    ]

    for i, (tag, title, clr, bullets) in enumerate(stages):
        x = Inches(0.4) + i * Inches(4.2)
        add_rect(slide, x, Inches(1.2), Inches(4.0), Inches(5.7), BG_CARD, clr)
        add_tag(slide, x + Inches(0.15), Inches(1.3), tag, clr, BG_DARK, width=Inches(1.4))
        add_text_box(slide, x + Inches(0.15), Inches(1.7), Inches(3.7), Inches(0.4),
                     title, font_size=18, colour=WHITE, bold=True)
        add_multiline(slide, x + Inches(0.15), Inches(2.3), Inches(3.7), Inches(4.2),
                      [f"-- {b}" for b in bullets], font_size=11, colour=LIGHT, spacing=Pt(7))

    for i in range(2):
        x = Inches(4.45) + i * Inches(4.2)
        add_text_box(slide, x, Inches(3.5), Inches(0.5), Inches(0.5),
                     '->', font_size=36, colour=ACCENT5, bold=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 13, TOTAL_SLIDES)


def slide_14_physique(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT2)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Contraintes physiques integrees dans le modele IA",
                 font_size=26, colour=ACCENT2, bold=True)

    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
                 "7 contraintes physiques codees en dur -- le modele NE PEUT PAS les violer structurellement.",
                 font_size=14, colour=ACCENT4, bold=True)

    constraints = [
        ("1. Application du signe", "dKL = -|d1| <= 0, dKR = -|d2| <= 0, dKLR = +|d3| >= 0",
         "Via -abs() et +abs() -- impossible de produire le mauvais signe", "Architectural"),
        ("2. Courbes monotones", "K(t) = K0 + cumsum(chutes contraintes)",
         "La somme cumulative de chutes signees garantit la monotonicite", "Architectural"),
        ("3. Separation init./chute", "Slot 1 -> tete initiale, Slots 2-21 -> tete de chute",
         "Des MLP separees refletent la distinction physique", "Architectural"),
        ("4. Echelle log-signee", "y_log = signe(y) * log(1+|y|)",
         "Preserve le fait que KLR < 0 tandis que KL, KR > 0", "Pretraitement"),
        ("5. Perte initiale ponderee", "L_initial = 5 x MSE sur le pas 1",
         "La precision du 1er pas est cruciale (base de cumsum)", "Entrainement"),
        ("6. Preservation de forme", "Perte Huber sur les 1eres differences",
         "Garantit que la forme de degradation correspond a la realite", "Entrainement"),
        ("7. Penalite physique-monotone", "relu(dKL) + relu(dKR) + relu(-dKLR) x 0.2",
         "Penalite souple supplementaire dans la perte (M8/M9)", "Entrainement"),
    ]

    for i, (name, formula, explanation, ctype) in enumerate(constraints):
        y = Inches(1.4) + i * Inches(0.82)
        bg = BG_CARD if i % 2 == 0 else BG_ALT
        add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.72), bg,
                 ACCENT2 if 'Architectural' in ctype else DIM)
        add_text_box(slide, Inches(0.7), y + Inches(0.02), Inches(2.5), Inches(0.3),
                     name, font_size=11, colour=ACCENT2 if 'Architectural' in ctype else ACCENT4, bold=True)
        add_text_box(slide, Inches(3.3), y + Inches(0.02), Inches(4.5), Inches(0.3),
                     formula, font_size=10, colour=ACCENT1)
        add_text_box(slide, Inches(3.3), y + Inches(0.35), Inches(6.0), Inches(0.3),
                     explanation, font_size=9, colour=DIM)
        add_tag(slide, Inches(11.0), y + Inches(0.18), ctype,
                ACCENT2 if 'Architectural' in ctype else ACCENT4, BG_DARK, width=Inches(1.6))

    add_slide_number(slide, 14, TOTAL_SLIDES)


def slide_15_resultats(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT3)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Resultats & Metriques de performance", font_size=26, colour=ACCENT3, bold=True)

    metrics_cards = [
        ('R2 global', '0.9891', "Le modele explique 98.9%\nde la variance", ACCENT3),
        ('MAPE courbe', '2.37%', "Erreur moy. en % abs.\nsur toutes les courbes", ACCENT4),
        ('Parametres', '45 502', "19.7% de moins que\nle prof M6 (56 646)", ACCENT1),
        ('Scenarios de test', '204', "Sur 1 019 au total\n(20% retenus)", ACCENT5),
    ]

    for i, (label, value, sub, clr) in enumerate(metrics_cards):
        x = Inches(0.5) + i * Inches(3.2)
        add_rect(slide, x, Inches(1.1), Inches(2.9), Inches(1.6), BG_CARD, clr)
        add_text_box(slide, x + Inches(0.15), Inches(1.15), Inches(2.6), Inches(0.3),
                     label, font_size=10, colour=DIM, bold=True)
        add_text_box(slide, x + Inches(0.15), Inches(1.45), Inches(2.6), Inches(0.5),
                     value, font_size=28, colour=clr, bold=True)
        add_text_box(slide, x + Inches(0.15), Inches(2.0), Inches(2.6), Inches(0.5),
                     sub, font_size=9, colour=DIM)

    add_rect(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(3.9), BG_CARD, ACCENT1)
    add_text_box(slide, Inches(0.8), Inches(3.1), Inches(5.0), Inches(0.4),
                 "Performance par variable (M9 vs M6 Enseignant)", font_size=14, colour=ACCENT1, bold=True)

    headers = ['Variable', 'R2 M6 (Prof)', 'R2 M9 (Psi-NN)', 'Amelioration', 'MAPE M9']
    for j, hdr in enumerate(headers):
        x = Inches(0.7) + j * Inches(2.4)
        add_text_box(slide, x, Inches(3.55), Inches(2.3), Inches(0.35),
                     hdr, font_size=10, colour=ACCENT1, bold=True, alignment=PP_ALIGN.CENTER)

    rows = [
        ('KL (Laterale)', '0.9768', '0.9905', '+0.0137', '2.32%'),
        ('KR (Rotationnelle)', '0.9634', '0.9796', '+0.0162', '2.42%'),
        ('KLR (Couplage)', '0.9690', '0.9860', '+0.0170', '2.38%'),
        ('Global', '0.9804', '0.9891', '+0.0087', '2.37%'),
    ]

    for r, row_data in enumerate(rows):
        y = Inches(4.0) + r * Inches(0.55)
        for c, cell in enumerate(row_data):
            x = Inches(0.7) + c * Inches(2.4)
            clr = WHITE if c == 0 else (GREEN if c == 2 else (ACCENT4 if c == 3 else LIGHT))
            is_bold = c == 0 or c == 3 or (r == 3)
            add_text_box(slide, x, y, Inches(2.3), Inches(0.4),
                         cell, font_size=11, colour=clr, bold=is_bold, alignment=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6), BG_ALT, GREEN)
    add_text_box(slide, Inches(0.8), Inches(6.25), Inches(11.7), Inches(0.45),
                 "M9 depasse le prof M6 sur TOUTES les variables avec 19.7% de parametres en moins -- "
                 "l'eleve a depasse le maitre grace a la decouverte de structure + gating SwiGLU.",
                 font_size=12, colour=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 15, TOTAL_SLIDES)


def slide_16_comparaison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT1)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Comparaison complete : M6 -> M7 -> M8 -> M9",
                 font_size=26, colour=ACCENT1, bold=True)

    headers = ['Metrique', 'M6 Prof', 'M7 Psi-NN', 'M8 Efficace', 'M9 SwiGLU']
    h_clrs = [WHITE, DIM, RGBColor(0x64, 0xB5, 0xF6), RGBColor(0xFF, 0xB7, 0x4D), GREEN]
    col_w = Inches(2.4)

    for j, (hdr, hclr) in enumerate(zip(headers, h_clrs)):
        x = Inches(0.4) + j * Inches(2.5)
        add_rect(slide, x, Inches(1.1), col_w, Inches(0.45), BG_CARD, hclr)
        add_text_box(slide, x + Inches(0.1), Inches(1.12), col_w - Inches(0.2), Inches(0.38),
                     hdr, font_size=11, colour=hclr, bold=True, alignment=PP_ALIGN.CENTER)

    rows = [
        ['Parametres', '56 646', '55 770', '45 486', '45 502'],
        ['Compression', '--', '1.5%', '19.7%', '19.7%'],
        ['MLP du slot', 'GELU 128', 'GELU 128', 'Goulot 48', 'SwiGLU 32'],
        ['k* (Prototypes)', '--', '6', '5', '5'],
        ['Matrice R', '--', 'Fixe', 'Apprenable', 'Apprenable'],
        ['Perte physique', '--', 'Non', 'Oui (l=0.2)', 'Oui (l=0.2)'],
        ['Reg. entropie', '--', 'Non', 'Oui (l=0.02)', 'Oui (l=0.02)'],
        ['R2 global', '0.9804', '0.9897', '0.9882', '0.9891'],
        ['R2 KL', '0.9768', '0.9910', '0.9888', '0.9905'],
        ['R2 KR', '0.9634', '0.9807', '0.9780', '0.9796'],
        ['R2 KLR', '0.9690', '0.9860', '0.9845', '0.9860'],
    ]

    for r, row_data in enumerate(rows):
        y = Inches(1.65) + r * Inches(0.48)
        for c, cell in enumerate(row_data):
            x = Inches(0.4) + c * Inches(2.5)
            bg = BG_CARD if r % 2 == 0 else BG_ALT
            clr = LIGHT if c == 0 else h_clrs[c]
            is_bold = c == 0

            if 'R2' in row_data[0] and c > 0:
                try:
                    val = float(cell)
                    row_vals = []
                    for cc in range(1, 5):
                        try: row_vals.append(float(row_data[cc]))
                        except: row_vals.append(0)
                    if val == max(row_vals):
                        clr = GREEN
                        is_bold = True
                except:
                    pass

            add_rect(slide, x, y, col_w, Inches(0.42), bg)
            add_text_box(slide, x + Inches(0.08), y + Inches(0.03), col_w - Inches(0.16), Inches(0.33),
                         cell, font_size=10, colour=clr, bold=is_bold,
                         alignment=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)

    add_slide_number(slide, 16, TOTAL_SLIDES)


def slide_17_webapp(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT4)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Application web interactive & Tableau de bord", font_size=26, colour=ACCENT4, bold=True)

    add_rect(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.8), BG_CARD, ACCENT4)
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.4),
                 "Tableau de bord complet construit avec Flask + Chart.js",
                 font_size=16, colour=ACCENT4, bold=True)
    webapp_desc = [
        "-- Visualisation interactive en temps reel des predictions de l'IA vs donnees experimentales reelles.",
        "-- Les ingenieurs peuvent parcourir les 204 scenarios de test, inspecter les valeurs par pas et analyser les erreurs.",
        "-- Le tableau de bord inclut la visualisation de la structure Psi-NN, l'analyse silhouette et l'importance des features.",
    ]
    add_multiline(slide, Inches(0.8), Inches(1.65), Inches(11.5), Inches(0.9),
                  webapp_desc, font_size=12, colour=LIGHT, spacing=Pt(3))

    features = [
        ("Graphiques de prediction", "Courbes KL, KR, KLR :\nCible vs Predit\navec surlignage d'erreur", ACCENT1),
        ("Tableau par pas", "Valeurs pas a pas\navec % d'erreur colore\n(vert < 5%, rouge > 20%)", ACCENT3),
        ("Navigateur de scenarios", "Parcourir les 204 cas\navec MAPE, R2, classement\net parametres d'entree", ACCENT4),
        ("Structure Psi-NN", "Clusters de prototypes,\nheatmap mat. de relation,\nanalyse silhouette", ACCENT5),
        ("Diagnostic pire cas", "Top 10 pires scenarios\nimportance des features\nanalyse de correlation", ACCENT2),
        ("Architecture du modele", "Schema complet\nComparaison SwiGLU vs GELU\ndecomposition des params", GREEN),
    ]

    for i, (title, desc_txt, clr) in enumerate(features):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + col * Inches(4.2)
        y = Inches(3.2) + row * Inches(2.0)
        add_rect(slide, x, y, Inches(3.9), Inches(1.8), BG_CARD, clr)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.08), Inches(3.6), Inches(0.35),
                     title, font_size=12, colour=clr, bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.5), Inches(3.6), Inches(1.1),
                     desc_txt, font_size=10, colour=LIGHT)

    add_slide_number(slide, 17, TOTAL_SLIDES)


def slide_18_contributions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT5)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Contributions scientifiques & Potentiel de publication",
                 font_size=26, colour=ACCENT5, bold=True)

    contribs = [
        ("1. Architecture originale : Psi-NN", ACCENT1,
         "Premiere application de reseaux a decouverte de structure en ingenierie geotechnique. "
         "Le pipeline Psi-NN trouve automatiquement le nombre optimal de composants prototypes."),
        ("2. SwiGLU pour le ML scientifique", GREEN,
         "Premiere utilisation du gating SwiGLU (issu de LLaMA/PaLM) dans des modeles informes par la physique. "
         "Demontre que les techniques NLP modernes se transferent aux domaines d'ingenierie."),
        ("3. Physique codee en dur", ACCENT4,
         "Les 3 variables de rigidite ont une monotonicite architecturalement garantie. "
         "Contrairement aux penalites souples, ces contraintes sont impossibles a violer."),
        ("4. Compression du modele", ACCENT2,
         "Reduction de 19.7% des parametres tout en ameliorant la precision (R2 0.9804 -> 0.9891). "
         "L'eleve compresse depasse le maitre grace a la decouverte de structure."),
        ("5. Jumeau Numerique interactif", ACCENT5,
         "Application web complete pour l'analyse de scenarios en temps reel, "
         "le diagnostic d'erreurs et la visualisation de l'importance des features."),
        ("6. Methodologie exhaustive", ACCENT3,
         "5 versions du modele (M5->M9) avec ablation systematique : suppression LSTM, variantes MLP, "
         "R apprenable vs fixe, perte physique -- fournissant des preuves claires pour chaque choix."),
    ]

    for i, (title, clr, desc_txt) in enumerate(contribs):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.4)
        y = Inches(1.1) + row * Inches(2.0)
        add_rect(slide, x, y, Inches(6.1), Inches(1.8), BG_CARD, clr)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.08), Inches(5.8), Inches(0.35),
                     title, font_size=13, colour=clr, bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.5), Inches(5.8), Inches(1.2),
                     desc_txt, font_size=11, colour=LIGHT)

    add_slide_number(slide, 18, TOTAL_SLIDES)


def slide_19_nsfc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, GOLD)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Contexte du projet de recherche (\u9762\u4e0a\u9879\u76ee\u7533\u8bf7\u4e66)",
                 font_size=26, colour=GOLD, bold=True)

    add_rect(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8), BG_CARD, GOLD)
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.4),
                 "Programme general NSFC -- Conception anti-typhon des fondations OWT",
                 font_size=16, colour=GOLD, bold=True)

    proposal_info = [
        "-- La recherche s'inscrit dans un projet finance par le NSFC",
        "   (Fondation Nationale des Sciences Naturelles de Chine -- Programme general).",
        "",
        "-- Titre : Jumeau Numerique pilote par l'IA pour la prediction du",
        "   decalage de frequence naturelle post-typhon des eoliennes offshore",
        "   avec fondations sur pieux degradees",
        "",
        "-- Le projet definit un cadre integrant :",
        "   1. Module de conditions d'appui (SCm) : Predit la degradation de rigidite",
        '      (KL, KR, KLR) -- c\'est ce que M6-M9 implemente',
        "   2. Module eolienne (WTm) : Methodes analytiques de Rayleigh-Ritz pour",
        "      calculer les frequences naturelles tour+pales",
        "   3. Jumeau Numerique couple : SCm alimente WTm -> predit le risque de resonance",
        "   4. IA explicable (XAI) : Propagation de pertinence couche par couche (LRP)",
        "   5. Quantification d'incertitude : Monte Carlo Dropout pour les intervalles de confiance",
        "",
        "-- Ma contribution de stage se concentre sur l'implementation et l'optimisation",
        "   du module SCm (le moteur de prediction IA), soit l'evolution M6->M9 presentee ici.",
        "",
        "-- Document de reference : \u9762\u4e0a\u9879\u76ee\u7533\u8bf7\u4e66\u64b0\u5199\u63d0\u7eb2.pdf",
    ]
    add_multiline(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.0),
                  proposal_info, font_size=12, colour=LIGHT, spacing=Pt(3))

    add_slide_number(slide, 19, TOTAL_SLIDES)


def slide_20_chronologie(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT3)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Chronologie & Plan de travail", font_size=26, colour=ACCENT3, bold=True)

    phases = [
        ('Phase 1 -- Terminee', 'Revue de litterature & analyse des donnees', ACCENT3, '[OK]', [
            "Revue de litterature sur la degradation des pieux OWT",
            "Etude des architectures slot attention & transformer",
            "Pretraitement et exploration des donnees (REAL DATA.xlsx)",
        ]),
        ('Phase 2 -- Terminee', 'Developpement du modele (M5 a M9)', GREEN, '[OK]', [
            "M5 : Base LSTM+Slot Attention (157K params)",
            "M6 : Prof simplifie (56K, sans LSTM)",
            "M7->M8->M9 : Pipeline de compression Psi-NN",
            "Upgrade SwiGLU atteignant l'etat de l'art (R2=0.99)",
        ]),
        ('Phase 3 -- Terminee', 'Application web & Validation', ACCENT1, '[OK]', [
            "Tableau de bord complet Flask+Chart.js",
            "Analyse par scenario & diagnostic d'erreurs",
            "Importance des features & identification pires cas",
        ]),
        ('Phase 4 -- URGENT', 'Soutenance de stage & Inscription HITSZ', ACCENT2, '[!!]', [
            "Soutenance de stage : debut septembre 2026",
            "Inscription HITSZ : 27-28 aout 2026",
            "Preparation du rapport de stage & presentation finale",
            "Validation des credits pour l'inscription doctorale",
        ]),
        ('Phase 5 -- Doctorat HITSZ', 'PhD en IA (Control Science & Engineering)', ACCENT4, '[..]', [
            "Rentree : automne 2026 a HITSZ (#3 mondial en Ingenierie)",
            "College of Artificial Intelligence, bourse complete",
            "Publication dans des revues internationales",
            "Extension au Jumeau Numerique couple SCm+WTm + XAI",
        ]),
    ]

    for i, (phase, title, clr, status, bullets) in enumerate(phases):
        y = Inches(1.05) + i * Inches(1.25)
        add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(1.1), BG_CARD, clr)
        add_text_box(slide, Inches(0.65), y + Inches(0.2), Inches(0.5), Inches(0.4),
                     status, font_size=12, colour=clr, bold=True)
        add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(3.2), Inches(0.35),
                     phase, font_size=11, colour=clr, bold=True)
        add_text_box(slide, Inches(1.2), y + Inches(0.38), Inches(3.2), Inches(0.35),
                     title, font_size=12, colour=WHITE, bold=True)
        bullet_str = ' | '.join(bullets)
        add_text_box(slide, Inches(4.6), y + Inches(0.15), Inches(8.0), Inches(0.7),
                     bullet_str, font_size=9, colour=LIGHT)

    add_slide_number(slide, 20, TOTAL_SLIDES)


def slide_21_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT1)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 "Conclusion & Prochaines etapes", font_size=26, colour=ACCENT1, bold=True)

    add_rect(slide, Inches(0.5), Inches(1.1), Inches(6.0), Inches(5.8), BG_CARD, ACCENT1)
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.4), Inches(0.4),
                 "Ce qui a ete accompli", font_size=16, colour=ACCENT1, bold=True)

    accomplishments = [
        "-- Developpement d'un systeme IA complet, informe par la",
        "   physique, pour la prediction de degradation de rigidite.",
        "",
        "-- Evolution a travers 5 versions (M5 -> M9),",
        "   chacune avec des ameliorations claires et justifiees.",
        "",
        "-- Le modele final M9 atteint :",
        "   -- R2 = 0.9891 (98.9% de la variance expliquee)",
        "   -- 19.7% de compression des parametres",
        "   -- Garanties physiques codees en dur",
        "   -- Gating SwiGLU a l'etat de l'art",
        "",
        "-- Construction d'un tableau de bord web interactif",
        "   pour l'analyse de scenarios en temps reel.",
        "",
        "-- Documentation exhaustive produite",
        "   (rapports LaTeX, diagrammes, code source complet).",
    ]
    add_multiline(slide, Inches(0.8), Inches(1.7), Inches(5.4), Inches(5.0),
                  accomplishments, font_size=12, colour=LIGHT, spacing=Pt(3))

    add_rect(slide, Inches(6.8), Inches(1.1), Inches(6.0), Inches(5.8), BG_CARD, ACCENT4)
    add_text_box(slide, Inches(7.1), Inches(1.2), Inches(5.4), Inches(0.4),
                 "Prochaines etapes", font_size=16, colour=ACCENT4, bold=True)

    next_steps = [
        ("ETAPE IMMEDIATE : Soutenance de stage", ACCENT2),
        "   -- Prevue debut septembre 2026",
        "   -- Requise pour valider le stage EMSI",
        "   -- Necessaire pour finaliser l'inscription a HITSZ",
        "",
        ("Inscription HITSZ : 27-28 aout 2026", GOLD),
        "",
        ("Doctorat a HITSZ (automne 2026) :", ACCENT1),
        "   -- College of Artificial Intelligence",
        "   -- Specialite : Control Science & Engineering",
        "   -- #3 mondial en Ingenierie (US News 2025-2026)",
        "   -- Bourse : 100% scolarite + 50 000 RMB/an",
        "",
        "-- Publication dans une revue internationale",
        "-- Extension au Jumeau Numerique couple SCm+WTm",
        "-- Ajout XAI (LRP) & quantification d'incertitude",
        "-- These de doctorat & soutenance",
    ]
    add_multiline(slide, Inches(7.1), Inches(1.7), Inches(5.4), Inches(5.0),
                  next_steps, font_size=12, colour=LIGHT, spacing=Pt(3))

    add_slide_number(slide, 21, TOTAL_SLIDES)


def slide_22_merci(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_rect(slide, Inches(0), Inches(0), Inches(4.44), Inches(0.08), MOROCCO_R)
    add_rect(slide, Inches(4.44), Inches(0), Inches(4.44), Inches(0.08), MOROCCO_G)
    add_rect(slide, Inches(8.88), Inches(0), Inches(4.45), Inches(0.08), SZU_RED)

    add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2),
                 'Merci', font_size=52, colour=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(0.6),
                 "\u0634\u0643\u0631\u0627 \u062c\u0632\u064a\u0644\u0627  --  \u8c22\u8c22  --  Thank you",
                 font_size=22, colour=ACCENT1, alignment=PP_ALIGN.CENTER)

    add_rect(slide, Inches(3.2), Inches(4.3), Inches(6.9), Inches(2.2), BG_CARD, ACCENT1)
    add_text_box(slide, Inches(3.5), Inches(4.4), Inches(6.3), Inches(0.4),
                 'Questions & Discussion', font_size=22, colour=ACCENT1, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.5), Inches(4.9), Inches(6.3), Inches(1.5),
                 "Youssef KHALLOUQI\n"
                 "EMSI (Ecole Marocaine des Sciences de l'Ingenieur)\n"
                 "Stage : Shenzhen University (\u6df1\u5733\u5927\u5b66) | Doctorat : HITSZ (\u54c8\u5de5\u5927\u6df1\u5733)\n"
                 "IA Informee par la Physique pour les Fondations Eoliennes Offshore",
                 font_size=12, colour=LIGHT, alignment=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0), Inches(7.3), Inches(13.333), Inches(0.06), ACCENT3)


# ======================================================
#  MAIN
# ======================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Generation des diapositives...")

    slide_01_titre(prs);        print("  [1/22] Page de titre")
    slide_02_sommaire(prs);     print("  [2/22] Sommaire")
    slide_03_universite(prs);   print("  [3/22] Universite d'accueil")
    slide_04_acceptation(prs);  print("  [4/22] Acceptation & Doctorat")
    slide_05_encadrement(prs);  print("  [5/22] Cadre de supervision")
    slide_06_sujet(prs);        print("  [6/22] Sujet de recherche")
    slide_07_probleme(prs);     print("  [7/22] Probleme d'ingenierie")
    slide_08_donnees(prs);      print("  [8/22] Donnees & parametres")
    slide_09_approche_ia(prs);  print("  [9/22] Approche IA")
    slide_10_architecture(prs); print("  [10/22] Architecture")
    slide_11_evolution(prs);    print("  [11/22] Evolution du modele")
    slide_12_m9_detail(prs);    print("  [12/22] Detail M9 SwiGLU")
    slide_13_pipeline(prs);     print("  [13/22] Pipeline 3 etapes")
    slide_14_physique(prs);     print("  [14/22] Contraintes physiques")
    slide_15_resultats(prs);    print("  [15/22] Resultats")
    slide_16_comparaison(prs);  print("  [16/22] Comparaison complete")
    slide_17_webapp(prs);       print("  [17/22] Application web")
    slide_18_contributions(prs);print("  [18/22] Contributions")
    slide_19_nsfc(prs);         print("  [19/22] Contexte NSFC")
    slide_20_chronologie(prs);  print("  [20/22] Chronologie")
    slide_21_conclusion(prs);   print("  [21/22] Conclusion")
    slide_22_merci(prs);        print("  [22/22] Merci")

    out_path = os.path.join(SCRIPT_DIR, 'Presentation_Superviseur_Marocain.pptx')
    prs.save(out_path)
    print(f"\n[OK] Presentation sauvegardee : {out_path}")
    print(f"     {TOTAL_SLIDES} diapositives generees.")
    return out_path


if __name__ == '__main__':
    main()
