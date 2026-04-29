"""
M7 vs M8 vs M9  -  Visual & Simple Explanation PPTX
====================================================
Focus: graphical architecture diagrams, simple language,
output differences, drawn layer-by-layer architectures.
"""

from pathlib import Path
import json

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

BASE_DIR = Path(__file__).resolve().parent
OUT_PATH = BASE_DIR / "M7_M8_M9_visual_comparison.pptx"

M7_JSON = BASE_DIR.parent / "M7" / "comparison.json"
M8_JSON = BASE_DIR.parent / "M8" / "comparison.json"
M9_JSON = BASE_DIR / "comparison.json"
M9_ANALYSIS_JSON = BASE_DIR / "scenario_analysis.json"

# ── Colours ─────────────────────────────────────────────────────────
C_BG     = RGBColor(250, 249, 245)
C_TITLE  = RGBColor(24, 35, 46)
C_SUB    = RGBColor(90, 105, 118)
C_M6     = RGBColor(160, 160, 160)
C_M7     = RGBColor(52, 120, 220)
C_M8     = RGBColor(231, 111, 81)
C_M9     = RGBColor(42, 157, 143)
C_ACCENT = RGBColor(244, 162, 97)
C_WHITE  = RGBColor(255, 255, 255)
C_DARK   = RGBColor(36, 46, 58)
C_LIGHT  = RGBColor(240, 240, 235)
C_GREEN  = RGBColor(39, 174, 96)
C_RED    = RGBColor(200, 60, 50)
C_GREY   = RGBColor(200, 200, 200)

def load_json(p):
    with open(p, "r", encoding="utf-8") as f:
        return json.load(f)


# ── Helpers ─────────────────────────────────────────────────────────
def bg(s, c=C_BG):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c

def title(s, t, sub=None, dark=False):
    tc = C_WHITE if dark else C_TITLE; sc = RGBColor(180,190,200) if dark else C_SUB
    b = s.shapes.add_textbox(Inches(.5), Inches(.15), Inches(12.3), Inches(.85))
    tf = b.text_frame; tf.clear(); p = tf.paragraphs[0]; r = p.add_run()
    r.text = t; r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = tc
    if sub:
        b2 = s.shapes.add_textbox(Inches(.55), Inches(.95), Inches(11.8), Inches(.5))
        tf2 = b2.text_frame; tf2.clear(); p2 = tf2.paragraphs[0]; r2 = p2.add_run()
        r2.text = sub; r2.font.size = Pt(15); r2.font.color.rgb = sc

def footer(s, t="PFE - M7 / M8 / M9 Visual Comparison"):
    b = s.shapes.add_textbox(Inches(.5), Inches(6.95), Inches(12.2), Inches(.3))
    tf = b.text_frame; tf.clear(); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = t; r.font.size = Pt(10); r.font.color.rgb = C_SUB

def box(s, x, y, w, h, fill=C_WHITE, border=C_M7, bw=2, rnd=True):
    """Rounded rectangle block."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rnd else MSO_SHAPE.RECTANGLE
    sh = s.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(bw)
    if rnd:
        # small corner radius
        sh.adjustments[0] = 0.06
    return sh

def txt(s, x, y, w, h, lines, sz=13, bold_first=False, color=C_TITLE, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line; r.font.size = Pt(sz); r.font.color.rgb = color
        r.font.bold = (i == 0 and bold_first); p.alignment = align
    return tb

def arrow_right(s, x, y, w=0.7, color=C_ACCENT):
    """Draw a right-pointing arrow shape."""
    sh = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.35))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()  # no line
    return sh

def arrow_down(s, x, y, h=0.5, color=C_ACCENT):
    sh = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(0.35), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def layer_box(s, x, y, w, h, label, sublabel, fill, border, text_color=C_WHITE):
    """A single neural-network layer box with centered label."""
    sh = box(s, x, y, w, h, fill=fill, border=border, bw=1)
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = label
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = text_color
    if sublabel:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sublabel
        r2.font.size = Pt(9); r2.font.color.rgb = text_color

def card(s, x, y, w, h, heading, bullets, color, body_sz=13):
    box(s, x, y, w, h, border=color)
    txt(s, x+.15, y+.08, w-.3, .35, [heading], sz=16, bold_first=True, color=color)
    lines = [f"  {b}" for b in bullets]
    txt(s, x+.15, y+.48, w-.3, h-.58, lines, sz=body_sz)

def styled_table(s, x, y, w, h, headers, rows, row_colors=None):
    nr = 1 + len(rows); nc = len(headers)
    table = s.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for c, hdr in enumerate(headers):
        cell = table.cell(0, c); cell.text = hdr
        cell.fill.solid(); cell.fill.fore_color.rgb = C_DARK
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = C_WHITE; r.font.bold = True; r.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri+1, ci); cell.text = str(val)
            if row_colors and ri < len(row_colors):
                cell.fill.solid(); cell.fill.fore_color.rgb = row_colors[ri]
            for p in cell.text_frame.paragraphs:
                for r in p.runs: r.font.size = Pt(11)
                p.alignment = PP_ALIGN.CENTER
    return table


# ── Draw a vertical architecture column ─────────────────────────
def draw_arch_column(s, cx, top_y, model_label, model_color, layers_info, total_params):
    """
    Draw a vertical stack of layer boxes centered at cx.
    layers_info: list of (label, sublabel, fill_color)
    """
    bw = 2.5  # box width
    bh = 0.52  # box height
    gap = 0.08  # gap between boxes
    arr_h = 0.28  # arrow height

    x = cx - bw / 2
    y = top_y

    # Model name header
    txt(s, x, y - 0.45, bw, 0.4, [model_label], sz=18, bold_first=True, color=model_color,
        align=PP_ALIGN.CENTER)

    for i, (label, sublabel, fill) in enumerate(layers_info):
        layer_box(s, x, y, bw, bh, label, sublabel, fill=fill, border=model_color)
        y += bh
        if i < len(layers_info) - 1:
            arrow_down(s, cx - 0.175, y, h=arr_h, color=model_color)
            y += arr_h + gap

    # Total params below
    txt(s, x, y + 0.1, bw, 0.3, [f"Total: {total_params:,} params"], sz=12,
        bold_first=True, color=model_color, align=PP_ALIGN.CENTER)


# ====================================================================
def build_ppt():
    m7 = load_json(M7_JSON); m8 = load_json(M8_JSON)
    m9 = load_json(M9_JSON); m9a = load_json(M9_ANALYSIS_JSON)
    m6 = m7["teacher_m6"]; m7m = m7["psi_model_m7"]
    m8m = m8["psi_model_m8"]; m9m = m9["psi_model_m9"]

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

    # ================================================================
    #  SLIDE 1 - Cover
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C_DARK)
    title(s, "M7  vs  M8  vs  M9", "A Visual & Simple Guide to the Three Model Versions", dark=True)

    # Three colored bars
    for i, (label, c, params, r2) in enumerate([
        ("M7 - Baseline", C_M7, m7m['params'], m7m['overall']['r2']),
        ("M8 - Efficient", C_M8, m8m['params'], m8m['overall']['r2']),
        ("M9 - SwiGLU (Best)", C_M9, m9m['params'], m9m['overall']['r2']),
    ]):
        bx = 1.0 + i * 3.9
        sh = box(s, bx, 2.5, 3.5, 2.0, fill=c, border=c, bw=0)
        txt(s, bx+.2, 2.6, 3.1, 1.8, [
            label,
            "",
            f"Parameters: {params:,}",
            f"Accuracy (R2): {r2:.4f}",
        ], sz=16, bold_first=True, color=C_WHITE)

    txt(s, 1.0, 5.0, 11.3, 1.5, [
        "What are these models?",
        "They all predict how a pile (foundation) degrades over time.",
        "Each version uses a smarter \"brain\" (MLP) to get more accurate with fewer parameters.",
    ], sz=15, bold_first=True, color=RGBColor(200, 210, 220))
    footer(s)

    # ================================================================
    #  SLIDE 2 - What does the model do? (simple explanation)
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "What Does the Model Do?", "A simple explanation of the input, the brain, and the output")

    # LEFT: Input
    box(s, 0.5, 1.6, 3.5, 4.5, fill=RGBColor(230,240,255), border=C_M7)
    txt(s, 0.7, 1.7, 3.1, 4.3, [
        "INPUT",
        "8 numbers describing the pile:",
        "",
        "  Soil type (PI)",
        "  Soil stiffness (Gmax)",
        "  Soil ratio (v)",
        "  Pile diameter (Dp)",
        "  Pile thickness (Tp)",
        "  Pile length (Lp)",
        "  Pile flexibility (Ip)",
        "  Shape ratio (Dp/Lp)",
        "",
        "Think of it as: describing the",
        "pile and the ground around it.",
    ], sz=13, bold_first=True, color=C_M7)

    # CENTER: Brain
    arrow_right(s, 4.1, 3.5, w=0.8, color=C_ACCENT)
    box(s, 5.0, 1.6, 3.5, 4.5, fill=RGBColor(255,248,235), border=C_ACCENT)
    txt(s, 5.2, 1.7, 3.1, 4.3, [
        "THE BRAIN (Psi-NN)",
        "",
        "1. Reads the 8 numbers",
        "",
        "2. Splits them into \"slots\"",
        "   (like memory boxes)",
        "",
        "3. Each slot learns a pattern",
        "",
        "4. Combines slots to predict",
        "   how the pile changes",
        "   over 21 loading cycles",
        "",
        "This is where M7/M8/M9 differ!",
    ], sz=13, bold_first=True, color=C_TITLE)

    # RIGHT: Output
    arrow_right(s, 8.6, 3.5, w=0.8, color=C_ACCENT)
    box(s, 9.5, 1.6, 3.5, 4.5, fill=RGBColor(230,250,240), border=C_M9)
    txt(s, 9.7, 1.7, 3.1, 4.3, [
        "OUTPUT",
        "3 curves over 21 time steps:",
        "",
        "  KL  = Lateral stiffness",
        "        (side-to-side strength)",
        "",
        "  KR  = Rocking stiffness",
        "        (tilting resistance)",
        "",
        "  KLR = Cross-coupling",
        "        (how KL and KR interact)",
        "",
        "These curves show how the pile",
        "gets weaker over time.",
    ], sz=13, bold_first=True, color=C_M9)

    footer(s)

    # ================================================================
    #  SLIDE 3 - The 5 building blocks (visual)
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "The 5 Building Blocks (All Versions Share This)",
          "Each model has the same pipeline - the difference is in block 3 (the MLP)")

    blocks = [
        (0.3, "1. Read Input", "Turns 8 numbers\ninto 64 features", RGBColor(200,215,240), C_M7),
        (2.85, "2. Slot Attention", "Splits features\ninto memory slots\n(like sorting info\ninto folders)", RGBColor(220,230,250), C_M7),
        (5.4, "3. Slot MLP", "Each slot goes\nthrough a small\n\"brain\" - THIS IS\nWHAT CHANGES!", RGBColor(255,235,220), C_M8),
        (7.95, "4. Rebuild Slots", "Uses a recipe (R)\nto reconstruct\nall 21 time steps\nfrom prototypes", RGBColor(230,250,240), C_M9),
        (10.5, "5. Output", "Produces the\n3 curves:\nKL, KR, KLR", RGBColor(255,245,220), C_ACCENT),
    ]

    for bx, lbl, desc, fill, border in blocks:
        box(s, bx, 1.7, 2.3, 3.0, fill=fill, border=border, bw=2)
        txt(s, bx+.1, 1.8, 2.1, .45, [lbl], sz=15, bold_first=True, color=border,
            align=PP_ALIGN.CENTER)
        txt(s, bx+.15, 2.35, 2.0, 2.2, desc.split("\n"), sz=12, color=C_TITLE,
            align=PP_ALIGN.CENTER)

    # Arrows between blocks
    for ax in [2.55, 5.1, 7.65, 10.2]:
        arrow_right(s, ax, 2.95, w=0.35, color=C_GREY)

    # Highlight block 3
    # Red dashed border effect (just a slightly thicker red outline)
    sh = box(s, 5.3, 1.6, 2.5, 3.2, fill=RGBColor(0,0,0), border=C_RED, bw=3, rnd=True)
    sh.fill.background()  # transparent

    txt(s, 0.3, 5.0, 12.7, 2.0, [
        "The KEY difference between M7, M8, and M9:",
        "",
        "  M7 uses a BIG brain (wide MLP) for each slot         -> lots of parameters, decent accuracy",
        "  M8 uses a SMALL brain (bottleneck MLP) for each slot  -> fewer parameters, saves space",
        "  M9 uses a SMART brain (SwiGLU gated MLP) for each slot -> same small size, BEST accuracy",
    ], sz=14, bold_first=True)
    footer(s)

    # ================================================================
    #  SLIDE 4 - M7 Architecture Drawn
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "M7 Architecture (Drawn)", "The baseline model - uses a wide MLP with 128 hidden neurons")

    # Draw vertical architecture
    m7_layers = [
        ("Input: 8 numbers", "soil + pile properties", RGBColor(180,200,230)),
        ("Embedding (8 -> 64)", "expand to 64 features", RGBColor(160,185,220)),
        ("Slot Attention (3 rounds)", "7 slots: 1 initial + 6 prototypes", RGBColor(140,170,210)),
        ("Wide MLP: 64 -> 128 -> 64", "GELU activation | 16,512 params", RGBColor(100,150,220)),
        ("Fixed Relation Matrix R", "6 protos -> 20 drop-slots (frozen)", RGBColor(120,160,215)),
        ("Output Head (64 -> 3)", "KL, KR, KLR x 21 steps", RGBColor(140,170,210)),
    ]
    draw_arch_column(s, 3.2, 1.9, "M7 - Baseline", C_M7, m7_layers, 55770)

    # Right side: simple explanation
    card(s, 6.8, 1.5, 6.0, 5.3, "How M7 Works (Simple)", [
        "1. Takes 8 input numbers about the pile and soil",
        "",
        "2. Expands them to 64 features (like zooming in)",
        "",
        "3. Slot Attention creates 7 \"memory slots\":",
        "   - 1 slot for initial state (time=0)",
        "   - 6 slots for typical degradation patterns",
        "",
        "4. Each slot goes through a WIDE MLP (128 neurons)",
        "   This is a big brain - uses lots of parameters",
        "   Think: sledgehammer to crack a nut",
        "",
        "5. A fixed recipe (R) builds 21 time steps from 6 patterns",
        "   Problem: this recipe can never improve!",
        "",
        "6. Outputs KL, KR, KLR curves (how pile degrades)",
        "",
        f"Result: R2 = {m7m['overall']['r2']:.4f}  |  Params: {m7m['params']:,}",
    ], C_M7, body_sz=12)
    footer(s)

    # ================================================================
    #  SLIDE 5 - M8 Architecture Drawn
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "M8 Architecture (Drawn)", "3 upgrades: smaller MLP, learnable recipe, physics rules")

    m8_layers = [
        ("Input: 8 numbers", "soil + pile properties", RGBColor(240,200,180)),
        ("Embedding (8 -> 64)", "expand to 64 features", RGBColor(235,185,165)),
        ("Slot Attention (3 rounds)", "6 slots: 1 initial + 5 prototypes", RGBColor(230,170,150)),
        ("Bottleneck MLP: 64->48->64", "GELU activation | 6,256 params", RGBColor(231,130,100)),
        ("Learnable Relation R", "5 protos -> 20 slots (adapts!)", RGBColor(230,150,130)),
        ("Output + Physics Check", "KL,KR must decrease / KLR increase", RGBColor(235,170,150)),
    ]
    draw_arch_column(s, 3.2, 1.9, "M8 - Efficient", C_M8, m8_layers, 45486)

    card(s, 6.8, 1.5, 6.0, 5.3, "What Changed in M8 (3 Upgrades)", [
        "UPGRADE 1: Smaller MLP (bottleneck)",
        "   Old: 64 -> 128 -> 64  (16,512 params)",
        "   New: 64 ->  48 -> 64  ( 6,256 params)  -62% smaller!",
        "   Like squeezing info through a narrow pipe",
        "   Forces the model to keep only what matters",
        "",
        "UPGRADE 2: Learnable Recipe (R)",
        "   M7: recipe frozen after discovery -> can't improve",
        "   M8: recipe keeps learning during training -> adapts",
        "   Like a chef who keeps improving a recipe vs one who",
        "   writes it down once and never changes it",
        "",
        "UPGRADE 3: Physics Rules",
        "   KL and KR MUST go down over time (pile gets weaker)",
        "   KLR MUST go up over time (coupling increases)",
        "   The model is penalized if it violates physics!",
        "",
        f"Result: R2 = {m8m['overall']['r2']:.4f}  |  Params: {m8m['params']:,}  (19.7% smaller)",
    ], C_M8, body_sz=11)
    footer(s)

    # ================================================================
    #  SLIDE 6 - M9 Architecture Drawn
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "M9 Architecture (Drawn)", "Uses SwiGLU - a smarter brain that decides what to focus on")

    m9_layers = [
        ("Input: 8 numbers", "soil + pile properties", RGBColor(180,230,220)),
        ("Embedding (8 -> 64)", "expand to 64 features", RGBColor(160,220,210)),
        ("Slot Attention (3 rounds)", "6 slots: 1 initial + 5 prototypes", RGBColor(140,210,200)),
        ("SwiGLU MLP (gated, dim=32)", "gate decides what matters | 6,272p", RGBColor(60,180,160)),
        ("Learnable Relation R", "same as M8 (adapts)", RGBColor(100,200,185)),
        ("Output + Physics Check", "same physics rules as M8", RGBColor(140,210,200)),
    ]
    draw_arch_column(s, 3.2, 1.9, "M9 - SwiGLU (Best)", C_M9, m9_layers, 45502)

    card(s, 6.8, 1.5, 6.0, 5.3, "What Changed in M9 (SwiGLU Brain)", [
        "THE KEY UPGRADE: SwiGLU replaces the plain MLP",
        "",
        "Normal MLP (M7/M8):",
        "   Data goes in -> gets transformed -> comes out",
        "   Every feature is treated the same way",
        "",
        "SwiGLU MLP (M9):",
        "   Data goes in through TWO paths:",
        "   Path 1 (Gate):  decides \"how important is this?\"",
        "   Path 2 (Value): carries the actual information",
        "   The gate CONTROLS how much of each feature passes",
        "",
        "   Like a smart filter: amplifies useful info,",
        "   suppresses noise - without adding parameters",
        "",
        "BONUS: Skips 2 training stages (reuses M8 results)",
        "   -> trains 2-3x faster!",
        "",
        f"Result: R2 = {m9m['overall']['r2']:.4f} (BEST)  |  Params: {m9m['params']:,}",
    ], C_M9, body_sz=11)
    footer(s)

    # ================================================================
    #  SLIDE 7 - Side-by-side architecture comparison
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "Side-by-Side Architecture Comparison",
          "Same pipeline, different MLP brains - spot the differences!")

    # Draw three mini-architectures
    mini_layers_m7 = [
        ("8 inputs", None, RGBColor(180,200,230)),
        ("Embed 64", None, RGBColor(160,185,220)),
        ("SlotAttn (7 slots)", None, RGBColor(140,170,210)),
        ("Wide MLP 128", "16,512 params", RGBColor(80,130,210)),
        ("Fixed R (6 proto)", None, RGBColor(120,160,215)),
        ("KL, KR, KLR", None, RGBColor(140,170,210)),
    ]
    mini_layers_m8 = [
        ("8 inputs", None, RGBColor(240,200,180)),
        ("Embed 64", None, RGBColor(235,185,165)),
        ("SlotAttn (6 slots)", None, RGBColor(230,170,150)),
        ("Bottleneck MLP 48", "6,256 params", RGBColor(231,111,81)),
        ("Learnable R (5 proto)", None, RGBColor(230,150,130)),
        ("KL,KR,KLR + physics", None, RGBColor(235,170,150)),
    ]
    mini_layers_m9 = [
        ("8 inputs", None, RGBColor(180,230,220)),
        ("Embed 64", None, RGBColor(160,220,210)),
        ("SlotAttn (6 slots)", None, RGBColor(140,210,200)),
        ("SwiGLU MLP 32", "6,272 params (gated)", RGBColor(42,157,143)),
        ("Learnable R (5 proto)", None, RGBColor(100,200,185)),
        ("KL,KR,KLR + physics", None, RGBColor(140,210,200)),
    ]

    draw_arch_column(s, 2.2, 1.8, "M7", C_M7, mini_layers_m7, 55770)
    draw_arch_column(s, 6.7, 1.8, "M8", C_M8, mini_layers_m8, 45486)
    draw_arch_column(s, 11.1, 1.8, "M9", C_M9, mini_layers_m9, 45502)

    # Highlight differences with annotations
    # Arrow between M7 and M8
    arrow_right(s, 3.7, 3.6, w=1.6, color=C_GREY)
    txt(s, 3.8, 3.2, 1.5, .3, ["Shrink MLP"], sz=10, color=C_M8, align=PP_ALIGN.CENTER)

    # Arrow between M8 and M9
    arrow_right(s, 8.2, 3.6, w=1.5, color=C_GREY)
    txt(s, 8.2, 3.2, 1.5, .3, ["Add gating"], sz=10, color=C_M9, align=PP_ALIGN.CENTER)

    footer(s)

    # ================================================================
    #  SLIDE 8 - MLP comparison (visual: data flow diagram)
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "The MLP Brain: How Data Flows Inside",
          "This is THE most important difference between M7, M8, and M9")

    # M7 MLP flow
    txt(s, 0.4, 1.5, 3.8, 0.4, ["M7: Wide MLP"], sz=18, bold_first=True, color=C_M7,
        align=PP_ALIGN.CENTER)
    m7_boxes = [
        (0.8, 2.1, "64 features", RGBColor(180,200,230)),
        (0.8, 2.9, "EXPAND to 128", RGBColor(100,150,220)),
        (0.8, 3.7, "GELU activation", RGBColor(120,160,215)),
        (0.8, 4.5, "COMPRESS to 64", RGBColor(140,170,210)),
    ]
    for bx, by, lbl, fill in m7_boxes:
        layer_box(s, bx, by, 3.0, 0.55, lbl, None, fill=fill, border=C_M7)
    for ay in [2.65, 3.45, 4.25]:
        arrow_down(s, 2.13, ay, h=0.25, color=C_M7)
    txt(s, 0.5, 5.15, 3.5, 0.7, [
        "Simple: everything goes through",
        "a wide pipe. Uses 16,512 params.",
    ], sz=11, color=C_SUB, align=PP_ALIGN.CENTER)

    # M8 MLP flow
    txt(s, 4.7, 1.5, 3.8, 0.4, ["M8: Bottleneck MLP"], sz=18, bold_first=True, color=C_M8,
        align=PP_ALIGN.CENTER)
    m8_boxes = [
        (5.1, 2.1, "64 features", RGBColor(240,200,180)),
        (5.1, 2.9, "SQUEEZE to 48", RGBColor(231,130,100)),
        (5.1, 3.7, "GELU activation", RGBColor(230,150,130)),
        (5.1, 4.5, "EXPAND to 64", RGBColor(235,170,150)),
    ]
    for bx, by, lbl, fill in m8_boxes:
        layer_box(s, bx, by, 3.0, 0.55, lbl, None, fill=fill, border=C_M8)
    for ay in [2.65, 3.45, 4.25]:
        arrow_down(s, 6.43, ay, h=0.25, color=C_M8)
    txt(s, 4.8, 5.15, 3.5, 0.7, [
        "Narrow pipe forces compression.",
        "Only 6,256 params (-62%).",
    ], sz=11, color=C_SUB, align=PP_ALIGN.CENTER)

    # M9 MLP flow (dual path!)
    txt(s, 9.0, 1.5, 4.0, 0.4, ["M9: SwiGLU (Gated)"], sz=18, bold_first=True, color=C_M9,
        align=PP_ALIGN.CENTER)
    # Input
    layer_box(s, 9.3, 2.1, 3.5, 0.55, "64 features", None,
              fill=RGBColor(180,230,220), border=C_M9)
    # Two parallel paths
    arrow_down(s, 10.2, 2.65, h=0.25, color=C_M9)
    arrow_down(s, 11.6, 2.65, h=0.25, color=C_M9)

    layer_box(s, 9.3, 3.0, 1.55, 0.55, "Gate path", "SiLU(W_g * x)",
              fill=RGBColor(60,180,160), border=C_M9)
    layer_box(s, 11.2, 3.0, 1.6, 0.55, "Value path", "W_v * x",
              fill=RGBColor(100,200,185), border=C_M9)

    # Merge
    arrow_down(s, 10.2, 3.55, h=0.25, color=C_M9)
    arrow_down(s, 11.6, 3.55, h=0.25, color=C_M9)

    layer_box(s, 9.3, 3.9, 3.5, 0.55, "MULTIPLY: gate x value", "gate controls what passes",
              fill=RGBColor(42,157,143), border=C_M9, text_color=C_WHITE)
    arrow_down(s, 10.9, 4.45, h=0.25, color=C_M9)
    layer_box(s, 9.3, 4.8, 3.5, 0.55, "Output: 64 features", None,
              fill=RGBColor(140,210,200), border=C_M9)

    txt(s, 9.1, 5.45, 3.9, 0.7, [
        "Smart filter: gate decides what",
        "matters. 6,272 params. BEST accuracy.",
    ], sz=11, color=C_SUB, align=PP_ALIGN.CENTER)

    footer(s)

    # ================================================================
    #  SLIDE 9 - Relation Matrix explained simply
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "The Relation Matrix R: A Recipe for Time Steps",
          "How do we go from 5-6 patterns to 21 time steps?")

    # Visual: prototypes -> R -> time steps
    txt(s, 0.5, 1.6, 12.3, 0.5, [
        "Think of it like this: we have 5 (or 6) typical \"patterns\" of how a pile degrades."
    ], sz=15, bold_first=True)
    txt(s, 0.5, 2.1, 12.3, 0.5, [
        "The Relation Matrix R is a RECIPE that mixes these patterns to create all 21 time steps."
    ], sz=15)

    # Left: Prototypes
    for i in range(5):
        box(s, 0.8, 2.9 + i*0.55, 2.0, 0.45, fill=RGBColor(255,235,220), border=C_M8)
        txt(s, 0.9, 2.95 + i*0.55, 1.8, 0.35, [f"Pattern {i+1}"], sz=11,
            align=PP_ALIGN.CENTER, color=C_M8)

    txt(s, 0.8, 5.7, 2.0, 0.3, ["5 prototypes"], sz=12, bold_first=True, color=C_M8,
        align=PP_ALIGN.CENTER)

    # Middle: R matrix
    arrow_right(s, 3.0, 4.0, w=1.0, color=C_ACCENT)
    box(s, 4.2, 2.9, 3.0, 3.0, fill=C_LIGHT, border=C_ACCENT)
    txt(s, 4.4, 3.0, 2.6, 2.8, [
        "RECIPE (R)",
        "",
        "Step 1 = 80% P1 + 20% P2",
        "Step 2 = 60% P1 + 40% P2",
        "Step 3 = 50% P2 + 50% P3",
        "...",
        "Step 20 = 90% P5 + 10% P4",
        "",
        "Each step is a mix",
        "of the patterns",
    ], sz=11, bold_first=True, color=C_TITLE, align=PP_ALIGN.CENTER)
    txt(s, 4.2, 5.95, 3.0, 0.3, ["Relation Matrix (20 x 5)"], sz=11,
        bold_first=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # Right: time steps
    arrow_right(s, 7.4, 4.0, w=1.0, color=C_ACCENT)
    for i in range(6):
        y = 2.9 + i * 0.48
        box(s, 8.6, y, 1.8, 0.38, fill=RGBColor(225,245,240), border=C_M9)
        txt(s, 8.7, y+0.02, 1.6, 0.3, [f"Step {i+1}"], sz=10,
            align=PP_ALIGN.CENTER, color=C_M9)
    txt(s, 8.6, 5.8, 1.8, 0.6, ["... up to\nStep 20"], sz=11,
        color=C_M9, align=PP_ALIGN.CENTER)

    # + initial slot
    box(s, 10.8, 2.9, 1.8, 0.45, fill=RGBColor(200,230,255), border=C_M7)
    txt(s, 10.9, 2.95, 1.6, 0.35, ["Step 0 (initial)"], sz=11,
        align=PP_ALIGN.CENTER, color=C_M7)
    txt(s, 10.8, 3.5, 1.8, 0.5, ["Special slot for\ntime=0 (healthy pile)"], sz=10,
        color=C_SUB, align=PP_ALIGN.CENTER)

    # Key difference
    box(s, 8.5, 4.5, 4.2, 1.8, fill=RGBColor(255,250,235), border=C_ACCENT)
    txt(s, 8.7, 4.6, 3.8, 1.6, [
        "M7: Recipe is FIXED (never changes)",
        "    -> like cooking with a rigid recipe",
        "",
        "M8/M9: Recipe LEARNS during training",
        "    -> like a chef who keeps improving",
        "    -> gets better over time!",
    ], sz=12, bold_first=True, color=C_TITLE)

    footer(s)

    # ================================================================
    #  SLIDE 10 - Physics Rules (simple visual)
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "Physics Rules: Teaching the Model Real-World Behavior",
          "M8 and M9 enforce physical laws - M7 does not")

    # Three columns: KL, KR, KLR with visual arrows
    for i, (var, desc, direction, clr) in enumerate([
        ("KL (lateral stiffness)", "Side-to-side strength", "MUST GO DOWN", C_RED),
        ("KR (rocking stiffness)", "Tilting resistance", "MUST GO DOWN", C_RED),
        ("KLR (cross-coupling)", "How KL and KR interact", "MUST GO UP", C_GREEN),
    ]):
        bx = 0.8 + i * 4.2
        box(s, bx, 1.7, 3.8, 2.5, fill=C_WHITE, border=clr)
        txt(s, bx+.2, 1.8, 3.4, .35, [var], sz=15, bold_first=True, color=clr,
            align=PP_ALIGN.CENTER)
        txt(s, bx+.2, 2.2, 3.4, .3, [desc], sz=12, color=C_SUB, align=PP_ALIGN.CENTER)

        # Visual arrow showing direction
        if "DOWN" in direction:
            # Draw declining visual
            txt(s, bx+.5, 2.6, 2.8, 1.0, [
                "Time 1:  ████████████",
                "Time 5:  █████████",
                "Time 10: ██████",
                "Time 15: ████",
                "Time 21: ██",
            ], sz=11, color=clr)
        else:
            txt(s, bx+.5, 2.6, 2.8, 1.0, [
                "Time 1:  ██",
                "Time 5:  ████",
                "Time 10: ██████",
                "Time 15: █████████",
                "Time 21: ████████████",
            ], sz=11, color=clr)

    # Explanation
    box(s, 0.8, 4.5, 5.8, 2.3, fill=RGBColor(255,245,240), border=C_M8)
    txt(s, 1.0, 4.6, 5.4, 2.1, [
        "Why does this matter?",
        "",
        "In the real world, piles ALWAYS get weaker over time.",
        "The stiffness KL and KR always decrease.",
        "The coupling KLR always increases.",
        "",
        "Without physics rules, the model might predict",
        "impossible results (like a pile getting STRONGER).",
    ], sz=13, bold_first=True)

    box(s, 7.0, 4.5, 5.8, 2.3, fill=C_WHITE, border=C_M9)
    txt(s, 7.2, 4.6, 5.4, 2.1, [
        "Which models enforce physics?",
        "",
        "   M7:  NO physics rules",
        "         (might produce impossible curves)",
        "",
        "   M8:  YES - physics penalty in loss function",
        "         (penalized if curves go wrong direction)",
        "",
        "   M9:  YES - same physics rules as M8",
    ], sz=13, bold_first=True)

    footer(s)

    # ================================================================
    #  SLIDE 11 - Training Pipeline (visual flowchart)
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "How Each Model is Trained",
          "M9 is much faster because it reuses work already done by M8")

    # Stage A
    box(s, 0.5, 1.7, 3.8, 1.8, fill=RGBColor(230,240,255), border=C_M7)
    txt(s, 0.7, 1.8, 3.4, 1.6, [
        "Stage A: Learn from Teacher",
        "",
        "A student model copies the",
        "teacher (M6) behavior.",
        "Like a student copying notes.",
        "Takes ~2,000 training rounds.",
    ], sz=12, bold_first=True, color=C_M7)

    arrow_right(s, 4.4, 2.4, w=0.6, color=C_GREY)

    # Stage B
    box(s, 5.1, 1.7, 3.5, 1.8, fill=RGBColor(255,240,225), border=C_M8)
    txt(s, 5.3, 1.8, 3.1, 1.6, [
        "Stage B: Discover Patterns",
        "",
        "Find k* typical patterns in",
        "the student's memory slots.",
        "Uses K-means clustering.",
        "Discovers R matrix recipe.",
    ], sz=12, bold_first=True, color=C_M8)

    arrow_right(s, 8.7, 2.4, w=0.6, color=C_GREY)

    # Stage C
    box(s, 9.4, 1.7, 3.5, 1.8, fill=RGBColor(225,248,240), border=C_M9)
    txt(s, 9.6, 1.8, 3.1, 1.6, [
        "Stage C: Final Training",
        "",
        "Train the compact model with",
        "k* slots and the R recipe.",
        "Takes ~2,500 training rounds.",
        "This is the final model!",
    ], sz=12, bold_first=True, color=C_M9)

    # Which model does what
    styled_table(s, 0.5, 4.0, 12.4, 2.5,
        ["", "Stage A (copy teacher)", "Stage B (find patterns)", "Stage C (final training)", "Total time"],
        [
            ["M7", "YES (2,000 rounds)", "YES (k*=6 found)", "YES (2,500 rounds)", "~4,500 rounds"],
            ["M8", "YES (2,000 rounds)", "YES (k*=5 found)", "YES (2,500 rounds)", "~4,500 rounds"],
            ["M9", "SKIP (reuses M8)", "SKIP (reuses M8)", "YES (2,500 rounds)", "~2,500 rounds!"],
        ],
        [RGBColor(230,240,255), RGBColor(255,240,225), RGBColor(210,245,230)])

    txt(s, 0.5, 6.6, 12.4, 0.5, [
        "M9 is 2-3x faster because it reuses M8's Stage A and B results - no need to redo work!",
    ], sz=14, bold_first=True, color=C_M9)
    footer(s)

    # ================================================================
    #  SLIDE 12 - Output Accuracy Comparison (charts)
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "Accuracy Comparison: Who Predicts Best?",
          "R2 = how well the model fits (1.0 = perfect)  |  RMSE = prediction error (lower = better)")

    # R2 chart
    r2d = CategoryChartData()
    r2d.categories = ["M6\n(Teacher)", "M7", "M8", "M9"]
    r2d.add_series("R2", [m6["overall"]["r2"], m7m["overall"]["r2"],
                           m8m["overall"]["r2"], m9m["overall"]["r2"]])
    r2c = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.7), Inches(6.2), Inches(4.5), r2d
    ).chart
    r2c.has_legend = False
    txt(s, 2.2, 6.2, 3.0, .3, ["R2 Score (higher = better)"], sz=12,
        bold_first=True, align=PP_ALIGN.CENTER)

    # RMSE chart
    rmse_d = CategoryChartData()
    rmse_d.categories = ["M6\n(Teacher)", "M7", "M8", "M9"]
    rmse_d.add_series("RMSE (billions)", [
        m6["overall"]["rmse"]/1e9, m7m["overall"]["rmse"]/1e9,
        m8m["overall"]["rmse"]/1e9, m9m["overall"]["rmse"]/1e9
    ])
    rmse_c = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.9), Inches(1.7), Inches(6.0), Inches(4.5), rmse_d
    ).chart
    rmse_c.has_legend = False
    txt(s, 8.5, 6.2, 3.5, .3, ["Error in billions (lower = better)"], sz=12,
        bold_first=True, align=PP_ALIGN.CENTER)

    footer(s)

    # ================================================================
    #  SLIDE 13 - Per-variable output comparison
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "Output Quality Per Variable: KL, KR, KLR",
          "How well does each model predict each stiffness curve?")

    # Bar charts for each variable
    for i, (var, clr) in enumerate([("KL", C_M7), ("KR", C_M8), ("KLR", C_M9)]):
        cd = CategoryChartData()
        cd.categories = ["M6", "M7", "M8", "M9"]
        cd.add_series(f"{var} R2", [
            m6["per_variable"][var]["r2"],
            m7m["per_variable"][var]["r2"],
            m8m["per_variable"][var]["r2"],
            m9m["per_variable"][var]["r2"],
        ])
        ch = s.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.5 + i*4.3), Inches(1.7), Inches(4.0), Inches(3.0), cd
        ).chart
        ch.has_legend = False

    txt(s, 1.5, 4.7, 2.5, .3, ["KL (side-to-side)"], sz=13, bold_first=True,
        align=PP_ALIGN.CENTER, color=C_M7)
    txt(s, 5.8, 4.7, 2.5, .3, ["KR (tilting)"], sz=13, bold_first=True,
        align=PP_ALIGN.CENTER, color=C_M8)
    txt(s, 10.1, 4.7, 2.5, .3, ["KLR (coupling)"], sz=13, bold_first=True,
        align=PP_ALIGN.CENTER, color=C_M9)

    # Simple interpretation
    box(s, 0.5, 5.2, 12.4, 1.5, fill=RGBColor(250,250,245), border=C_SUB, bw=1)
    txt(s, 0.7, 5.3, 12.0, 1.3, [
        "What this tells us:",
        f"  KL:  M7 is slightly best ({m7m['per_variable']['KL']['r2']:.4f}) but all are very close",
        f"  KR:  M9 is best ({m9m['per_variable']['KR']['r2']:.4f}) - this is the hardest variable to predict",
        f"  KLR: M9 is best ({m9m['per_variable']['KLR']['r2']:.4f}) - SwiGLU helps most on difficult variables",
        "  All student models beat the M6 teacher on every variable!",
    ], sz=13, bold_first=True)
    footer(s)

    # ================================================================
    #  SLIDE 14 - Model Size Comparison (visual)
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "Model Size: How Many Parameters?",
          "Fewer parameters = smaller model = faster to run = less memory needed")

    # Visual bars (proportional width)
    max_p = m6["params"]
    models_info = [
        ("M6 (Teacher)", m6["params"], C_M6, "0% (reference)"),
        ("M7", m7m["params"], C_M7, "1.5% smaller"),
        ("M8", m8m["params"], C_M8, "19.7% smaller"),
        ("M9", m9m["params"], C_M9, "19.7% smaller"),
    ]

    for i, (name, params, clr, comp) in enumerate(models_info):
        by = 2.0 + i * 1.2
        bar_w = (params / max_p) * 8.0

        # Bar
        sh = box(s, 2.5, by, bar_w, 0.7, fill=clr, border=clr, bw=0)

        # Label left
        txt(s, 0.3, by + 0.05, 2.0, 0.6, [name], sz=16, bold_first=True, color=clr,
            align=PP_ALIGN.RIGHT)

        # Value on bar
        txt(s, 2.7, by + 0.1, bar_w - 0.4, 0.5, [f"{params:,} params"], sz=14,
            bold_first=True, color=C_WHITE)

        # Compression label
        txt(s, 2.5 + bar_w + 0.2, by + 0.1, 3.0, 0.5, [comp], sz=13, color=clr)

    txt(s, 0.5, 6.2, 12.4, 0.5, [
        "M8 and M9 are nearly the same size (only 16 params difference!) "
        "but M9 is more accurate thanks to smarter SwiGLU gating.",
    ], sz=14, bold_first=True, color=C_TITLE)
    footer(s)

    # ================================================================
    #  SLIDE 15 - Where do parameters live? (stacked chart)
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "Where Do the Parameters Live?",
          "The MLP is where M8/M9 saved most space")

    chart_data = CategoryChartData()
    chart_data.categories = ["M6 (Teacher)", "M7", "M8", "M9"]
    chart_data.add_series("Slot MLP (the brain)", [16512, 16512, 6256, 6272])
    chart_data.add_series("Slot Attention", [20864, 20864, 20864, 20864])
    chart_data.add_series("Other (embed, init, output)", [576+1344+215, 576+448+215, 576+384+100+215, 576+384+100+215])

    ch_shape = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, Inches(0.5), Inches(1.7), Inches(7.5), Inches(5.0), chart_data
    )
    ch = ch_shape.chart
    ch.has_legend = True

    txt(s, 8.3, 1.7, 4.5, 5.0, [
        "Reading the chart:",
        "",
        "The Slot MLP (orange/top section) is where",
        "M8 and M9 saved the most space:",
        "",
        "  M6/M7: MLP = 16,512 params (huge!)",
        "  M8:    MLP =  6,256 params (-62%)",
        "  M9:    MLP =  6,272 params (-62%)",
        "",
        "Slot Attention (middle) stays the same",
        "in all models - it's already efficient.",
        "",
        "The MLP was the \"fat\" part of the model.",
        "M8 put it on a diet (bottleneck).",
        "M9 made it smarter instead (SwiGLU).",
    ], sz=13, bold_first=True)
    footer(s)

    # ================================================================
    #  SLIDE 16 - Worst-case analysis
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "Where Does M9 Struggle? (Worst Cases)",
          "Most predictions are excellent, but a few edge cases are harder")

    g = m9a["global_metrics"]
    wc = m9a["worst_scenarios"][:6]

    cat = CategoryChartData()
    cat.categories = [f"Scenario {d['scenario_id']}" for d in wc]
    cat.add_series("Average Error %", [d["overall"]["avg_curve_mape_pct"] for d in wc])
    ch = s.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.7), Inches(7.5), Inches(3.5), cat
    ).chart
    ch.has_legend = False

    card(s, 8.3, 1.7, 4.5, 3.5, "What This Means", [
        f"204 test scenarios evaluated",
        f"Average error: only {g['avg_curve_mape_pct']:.2f}%",
        f"Median error:  only {g['median_curve_mape_pct']:.2f}%",
        "",
        "Most predictions are very accurate!",
        "",
        "A few extreme pile configurations",
        "(unusual soil + pile combos) are harder.",
        "Scenario 56 is the worst (34% error).",
    ], C_ACCENT, body_sz=12)

    # Simple gauge-like indicator
    box(s, 0.5, 5.5, 12.4, 1.5, fill=RGBColor(250,250,245), border=C_SUB, bw=1)
    txt(s, 0.7, 5.6, 12.0, 1.3, [
        "Overall M9 Quality:",
        f"  Average accuracy (R2):  {g['avg_curve_r2']:.4f}  ->  Very good (close to 1.0 = perfect)",
        f"  Average error (MAPE):   {g['avg_curve_mape_pct']:.2f}%   ->  Very low (close to 0% = perfect)",
        f"  95th percentile error:  {g['p95_curve_mape_pct']:.2f}%   ->  Even worst 5% are below 6.3%",
    ], sz=14, bold_first=True)
    footer(s)

    # ================================================================
    #  SLIDE 17 - Complete differences table (simple language)
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "All Differences at a Glance",
          "A simple comparison of everything that changed between M7, M8, and M9")

    styled_table(s, 0.3, 1.5, 12.7, 5.5,
        ["What", "M7", "M8", "M9"],
        [
            ["MLP brain type", "Wide (big)", "Bottleneck (small)", "SwiGLU (smart)"],
            ["MLP brain size", "128 neurons", "48 neurons", "32 neurons (x2 paths)"],
            ["MLP parameters", "16,512", "6,256", "6,272"],
            ["Has feature gating?", "No", "No", "Yes (the key advantage)"],
            ["Recipe (R) type", "Fixed (frozen)", "Learns during training", "Learns during training"],
            ["Number of patterns", "6 prototypes", "5 prototypes", "5 (reused from M8)"],
            ["Physics rules?", "No", "Yes (penalty if wrong)", "Yes (same as M8)"],
            ["Training stages", "A + B + C (slow)", "A + B + C (slow)", "C only (fast!)"],
            ["Total parameters", "55,770", "45,486", "45,502"],
            ["Size vs teacher M6", "1.5% smaller", "19.7% smaller", "19.7% smaller"],
            [f"Accuracy (R2)", f"{m7m['overall']['r2']:.4f}", f"{m8m['overall']['r2']:.4f}",
             f"{m9m['overall']['r2']:.4f}  BEST"],
            [f"Error (RMSE x1e9)", f"{m7m['overall']['rmse']/1e9:.1f}", f"{m8m['overall']['rmse']/1e9:.1f}",
             f"{m9m['overall']['rmse']/1e9:.1f}  BEST"],
        ],
        None)
    footer(s)

    # ================================================================
    #  SLIDE 18 - Winner and Recommendation
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    title(s, "The Winner: M9", "Best accuracy, nearly same size, and fastest to train")

    # Trophy / winner visual
    sh = box(s, 4.2, 1.6, 5.0, 3.8, fill=C_M9, border=C_M9, bw=0)
    txt(s, 4.5, 1.8, 4.4, 3.4, [
        "M9 - SwiGLU Psi-NN",
        "",
        f"  Accuracy (R2):  {m9m['overall']['r2']:.4f}",
        f"  Parameters:     {m9m['params']:,}",
        f"  Compression:    19.7%",
        f"  Training:       Stage C only (fast!)",
        "",
        "  Best accuracy of all versions",
        "  Nearly same size as M8",
        "  2-3x faster to train",
        "  Enforces physics rules",
    ], sz=15, bold_first=True, color=C_WHITE)

    # Score cards
    scores = [
        ("Accuracy", "M9 wins", C_M9),
        ("Size", "M8 wins (by 16)", C_M8),
        ("Speed", "M9 wins (3x)", C_M9),
        ("Physics", "M8 = M9", C_ACCENT),
    ]
    for i, (criterion, result, clr) in enumerate(scores):
        bx = 0.5 + i * 3.2
        box(s, bx, 5.8, 2.9, 0.9, fill=C_WHITE, border=clr)
        txt(s, bx+.1, 5.85, 2.7, .4, [criterion], sz=14, bold_first=True, color=clr,
            align=PP_ALIGN.CENTER)
        txt(s, bx+.1, 6.25, 2.7, .35, [result], sz=12, color=C_TITLE,
            align=PP_ALIGN.CENTER)

    footer(s)

    # ================================================================
    #  SLIDE 19 - Executive Summary (dark)
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C_DARK)
    title(s, "Summary: The Evolution", dark=True)

    txt(s, 0.8, 1.5, 11.8, 5.5, [
        "M7 -> M8 -> M9: each version gets smarter, not bigger",
        "",
        "M7 (Baseline):",
        "   Uses a big wide brain (MLP with 128 neurons) and a fixed recipe.",
        "   Works well but wastes parameters. Only 1.5% smaller than teacher.",
        "",
        "M8 (Efficient):",
        "   Shrinks the brain (48 neurons = 62% less), makes recipe learnable,",
        "   adds physics rules. 19.7% smaller than teacher. Smart compression.",
        "",
        "M9 (SwiGLU - Best):",
        "   Replaces the brain with a gated version (SwiGLU).",
        "   The gate decides which features matter - like a smart filter.",
        "   Same small size as M8, but BETTER accuracy (0.9915 vs 0.9882).",
        "   Skips 2 training stages -> trains 2-3x faster.",
        "",
        "Bottom line: M9 is the best model. Use it.",
    ], sz=16, color=RGBColor(220, 230, 240))
    footer(s)

    # ── SAVE ──
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build_ppt()
