from pathlib import Path
import json

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

BASE_DIR = Path(__file__).resolve().parent
OUT_PATH = BASE_DIR / "M7_M8_M9_comparison.pptx"

M7_JSON = BASE_DIR.parent / "M7" / "comparison.json"
M8_JSON = BASE_DIR.parent / "M8" / "comparison.json"
M9_JSON = BASE_DIR / "comparison.json"
M9_ANALYSIS_JSON = BASE_DIR / "scenario_analysis.json"

# Color palette (warm + technical)
C_BG = RGBColor(246, 244, 238)
C_TITLE = RGBColor(24, 35, 46)
C_SUB = RGBColor(74, 90, 104)
C_M7 = RGBColor(26, 111, 223)
C_M8 = RGBColor(231, 111, 81)
C_M9 = RGBColor(42, 157, 143)
C_ACCENT = RGBColor(244, 162, 97)
C_WHITE = RGBColor(255, 255, 255)



def load_json(path: Path):
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)



def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_BG



def add_title(slide, title, subtitle=None):
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.9))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = C_TITLE

    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.55), Inches(1.0), Inches(11.8), Inches(0.6))
        stf = sbox.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.size = Pt(16)
        srun.font.color.rgb = C_SUB



def add_footer(slide, txt="PFE - Comparative model study"):
    fbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.2), Inches(0.3))
    tf = fbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(10)
    r.font.color.rgb = C_SUB



def add_bullet_card(slide, x, y, w, h, title, bullets, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    head = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.08), Inches(w - 0.3), Inches(0.4))
    htf = head.text_frame
    htf.clear()
    hp = htf.paragraphs[0]
    hr = hp.add_run()
    hr.text = title
    hr.font.bold = True
    hr.font.size = Pt(16)
    hr.font.color.rgb = color

    body = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.52), Inches(w - 0.3), Inches(h - 0.62))
    btf = body.text_frame
    btf.clear()
    for i, b in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = f"- {b}"
        p.level = 0
        p.font.size = Pt(13)
        p.font.color.rgb = C_TITLE



def build_ppt():
    m7 = load_json(M7_JSON)
    m8 = load_json(M8_JSON)
    m9 = load_json(M9_JSON)
    m9_analysis = load_json(M9_ANALYSIS_JSON)

    m7_m = m7["psi_model_m7"]
    m8_m = m8["psi_model_m8"]
    m9_m = m9["psi_model_m9"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(
        s,
        "M7 vs M8 vs M9",
        "Architecture, parameters, accuracy, and practical recommendation"
    )

    add_bullet_card(
        s, 0.8, 1.8, 3.9, 3.2, "M7 (baseline Psi-NN)",
        [
            "Wide slot MLP (64->128->64)",
            "Fixed relation matrix R",
            "k*=6 prototypes",
            "Params: 55,770"
        ], C_M7
    )
    add_bullet_card(
        s, 4.9, 1.8, 3.9, 3.2, "M8 (efficient)",
        [
            "Bottleneck MLP (64->48->64)",
            "Learnable relation matrix R",
            "Physics monotonic regularization",
            "Params: 45,486"
        ], C_M8
    )
    add_bullet_card(
        s, 9.0, 1.8, 3.3, 3.2, "M9 (SwiGLU)",
        [
            "SwiGLU slot MLP, hidden=32",
            "Same structured pipeline",
            "Reuses M8 discovery k*=5",
            "Params: 45,502"
        ], C_M9
    )
    add_footer(s)

    # Slide 2: Architecture evolution flow
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "Architecture Evolution", "What changed from M7 to M8 to M9")

    # Flow boxes
    boxes = [
        (0.7, 1.8, 3.8, 2.2, C_M7, "M7", ["Wide MLP", "Static R", "No physics monotonic loss"]),
        (4.8, 1.8, 3.8, 2.2, C_M8, "M8", ["Bottleneck MLP", "Learnable R", "Physics monotonic loss"]),
        (8.9, 1.8, 3.8, 2.2, C_M9, "M9", ["SwiGLU MLP", "Learnable R", "Physics monotonic loss"]),
    ]

    for x, y, w, h, c, title, lines in boxes:
        sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = C_WHITE
        sh.line.color.rgb = c
        sh.line.width = Pt(3)

        tb = s.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.12), Inches(w - 0.3), Inches(h - 0.2))
        tf = tb.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = title
        r0.font.bold = True
        r0.font.size = Pt(22)
        r0.font.color.rgb = c
        for line in lines:
            p = tf.add_paragraph()
            p.text = f"- {line}"
            p.font.size = Pt(13)
            p.font.color.rgb = C_TITLE

    # Simple arrows text
    a1 = s.shapes.add_textbox(Inches(4.56), Inches(2.65), Inches(0.2), Inches(0.4))
    a1.text_frame.text = ">"
    a1.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
    a1.text_frame.paragraphs[0].runs[0].font.color.rgb = C_ACCENT

    a2 = s.shapes.add_textbox(Inches(8.66), Inches(2.65), Inches(0.2), Inches(0.4))
    a2.text_frame.text = ">"
    a2.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
    a2.text_frame.paragraphs[0].runs[0].font.color.rgb = C_ACCENT

    add_footer(s)

    # Slide 3: Parameter comparison chart
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "Model Size Comparison", "M8/M9 keep strong compression vs M6 teacher")

    chart_data = CategoryChartData()
    chart_data.categories = ["M7", "M8", "M9"]
    chart_data.add_series("Params", [m7_m["params"], m8_m["params"], m9_m["params"]])

    ch = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.7), Inches(6.8), Inches(4.8), chart_data
    ).chart
    ch.has_legend = False
    ch.value_axis.has_major_gridlines = True
    ch.category_axis.has_major_gridlines = False

    # Add compression notes
    note = s.shapes.add_textbox(Inches(7.9), Inches(2.0), Inches(4.9), Inches(3.9))
    tf = note.text_frame
    tf.clear()
    lines = [
        "Compression vs M6 (56,646 params):",
        f"- M7: {(1 - m7_m['params']/56646)*100:.1f}%",
        f"- M8: {(1 - m8_m['params']/56646)*100:.1f}%",
        f"- M9: {(1 - m9_m['params']/56646)*100:.1f}%",
        "",
        "Takeaway:",
        "M8 and M9 deliver near-identical compact size,",
        "with M9 only +16 params vs M8."
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15 if i in (0, 5) else 13)
        p.font.bold = i in (0, 5)
        p.font.color.rgb = C_TITLE

    add_footer(s)

    # Slide 4: Accuracy chart
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "Accuracy Comparison (Overall)", "Higher R2 and lower RMSE indicate better fit")

    # R2 chart
    r2_data = CategoryChartData()
    r2_data.categories = ["M7", "M8", "M9"]
    r2_data.add_series("R2", [m7_m["overall"]["r2"], m8_m["overall"]["r2"], m9_m["overall"]["r2"]])
    r2_chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(1.7), Inches(6.1), Inches(4.6), r2_data
    ).chart
    r2_chart.has_legend = False

    rmse_data = CategoryChartData()
    rmse_data.categories = ["M7", "M8", "M9"]
    rmse_data.add_series(
        "RMSE (x1e9)",
        [m7_m["overall"]["rmse"] / 1e9, m8_m["overall"]["rmse"] / 1e9, m9_m["overall"]["rmse"] / 1e9]
    )
    rmse_chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.9), Inches(1.7), Inches(6.0), Inches(4.6), rmse_data
    ).chart
    rmse_chart.has_legend = False

    cap1 = s.shapes.add_textbox(Inches(2.6), Inches(6.25), Inches(2.3), Inches(0.3))
    cap1.text_frame.text = "R2 (higher is better)"
    cap1.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

    cap2 = s.shapes.add_textbox(Inches(8.8), Inches(6.25), Inches(3.0), Inches(0.3))
    cap2.text_frame.text = "RMSE in billions (lower is better)"
    cap2.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

    add_footer(s)

    # Slide 5: Variable-wise table
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "Per-variable R2", "KL, KR, KLR behavior across models")

    rows, cols = 4, 4
    table = s.shapes.add_table(rows, cols, Inches(1.0), Inches(1.8), Inches(11.2), Inches(2.8)).table
    headers = ["Model", "KL R2", "KR R2", "KLR R2"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_TITLE
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = C_WHITE
        cell.text_frame.paragraphs[0].runs[0].font.bold = True

    rows_data = [
        ("M7", m7_m["per_variable"]["KL"]["r2"], m7_m["per_variable"]["KR"]["r2"], m7_m["per_variable"]["KLR"]["r2"]),
        ("M8", m8_m["per_variable"]["KL"]["r2"], m8_m["per_variable"]["KR"]["r2"], m8_m["per_variable"]["KLR"]["r2"]),
        ("M9", m9_m["per_variable"]["KL"]["r2"], m9_m["per_variable"]["KR"]["r2"], m9_m["per_variable"]["KLR"]["r2"]),
    ]

    for r in range(1, rows):
        model_name, kl, kr, klr = rows_data[r - 1]
        vals = [model_name, f"{kl:.4f}", f"{kr:.4f}", f"{klr:.4f}"]
        for c, val in enumerate(vals):
            cell = table.cell(r, c)
            cell.text = val
            if model_name == "M7":
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(228, 240, 253)
            elif model_name == "M8":
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255, 236, 229)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(225, 245, 240)

    verdict = s.shapes.add_textbox(Inches(1.0), Inches(4.9), Inches(11.2), Inches(1.3))
    vtf = verdict.text_frame
    vtf.clear()
    vlines = [
        "Observed trend:",
        "- M9 improves overall R2 and RMSE vs both M7 and M8.",
        "- KR remains the hardest variable; this is visible in all model versions.",
    ]
    for i, line in enumerate(vlines):
        p = vtf.paragraphs[0] if i == 0 else vtf.add_paragraph()
        p.text = line
        p.font.size = Pt(14 if i == 0 else 13)
        p.font.bold = i == 0

    add_footer(s)

    # Slide 6: Stage design comparison
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "Pipeline Stage Comparison", "Practical training implications")

    add_bullet_card(
        s, 0.7, 1.6, 4.0, 4.6, "M7 Pipeline", [
            "Stage A: train student",
            "Stage B: discover k* and relation matrix R",
            "Stage C: train structured model",
            "Full process but larger final model"
        ], C_M7
    )
    add_bullet_card(
        s, 4.8, 1.6, 4.0, 4.6, "M8 Pipeline", [
            "Same 3 stages",
            "More compact model design in Stage C",
            "Introduces learnable relation matrix",
            "Good compression baseline"
        ], C_M8
    )
    add_bullet_card(
        s, 8.9, 1.6, 3.8, 4.6, "M9 Practical Pipeline", [
            "Reuses M8 Stage B artifacts",
            "Skips expensive Stage A/B retraining",
            "Runs Stage C with SwiGLU",
            "Best speed/quality tradeoff"
        ], C_M9
    )
    add_footer(s)

    # Slide 7: Worst-case diagnostics from M9
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "M9 Worst-case Diagnostics", "From scenario_analysis.json")

    g = m9_analysis["global_metrics"]
    top = m9_analysis["worst_scenarios"][0]

    summary = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(12.2), Inches(1.6))
    stf = summary.text_frame
    stf.clear()
    sum_lines = [
        f"Global curve metrics: Avg MAPE={g['avg_curve_mape_pct']:.3f}%, Avg R2={g['avg_curve_r2']:.5f}, Test scenarios={g['n_test']}",
        f"Worst scenario ID={top['scenario_id']} with Avg MAPE={top['overall']['avg_curve_mape_pct']:.2f}% and Avg R2={top['overall']['avg_curve_r2']:.4f}",
        "Interpretation: most scenarios are very accurate, but a few edge-case parameter combinations remain difficult."
    ]
    for i, line in enumerate(sum_lines):
        p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.bold = i == 0

    wc = m9_analysis["worst_scenarios"][:5]
    cat = CategoryChartData()
    cat.categories = [f"S{d['scenario_id']}" for d in wc]
    cat.add_series("Avg curve MAPE %", [d["overall"]["avg_curve_mape_pct"] for d in wc])
    ch = s.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.9), Inches(3.5), Inches(12.0), Inches(2.7), cat
    ).chart
    ch.has_legend = False

    add_footer(s)

    # Slide 8: Decision matrix
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "Decision Matrix", "Choosing the best version for deployment")

    rows, cols = 5, 5
    table = s.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(12.0), Inches(3.6)).table
    hdr = ["Criterion", "M7", "M8", "M9", "Best"]
    for i, h in enumerate(hdr):
        c = table.cell(0, i)
        c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = C_TITLE
        c.text_frame.paragraphs[0].runs[0].font.color.rgb = C_WHITE
        c.text_frame.paragraphs[0].runs[0].font.bold = True

    matrix = [
        ["Overall R2", f"{m7_m['overall']['r2']:.4f}", f"{m8_m['overall']['r2']:.4f}", f"{m9_m['overall']['r2']:.4f}", "M9"],
        ["Param count", str(m7_m['params']), str(m8_m['params']), str(m9_m['params']), "M8"],
        ["Compression vs M6", f"{(1-m7_m['params']/56646)*100:.1f}%", f"{(1-m8_m['params']/56646)*100:.1f}%", f"{(1-m9_m['params']/56646)*100:.1f}%", "M8/M9"],
        ["Training practicality", "Medium", "Medium", "High (skip A/B)", "M9"],
    ]

    for r in range(1, rows):
        for c in range(cols):
            table.cell(r, c).text = matrix[r - 1][c]

    call = s.shapes.add_textbox(Inches(0.8), Inches(5.7), Inches(12.0), Inches(0.9))
    ctf = call.text_frame
    ctf.clear()
    p = ctf.paragraphs[0]
    p.text = "Final recommendation: M9 is the best balanced choice (best accuracy, nearly same compact size as M8, and faster practical workflow)."
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_M9

    add_footer(s)

    # Slide 9: Practical runbook
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "How to Use M9", "Train, analyze worst scenarios, and open the web app")

    cmds = [
        "1) Train Stage C only:",
        "   cd M9",
        "   python train.py",
        "",
        "2) Build worst-scenario analysis:",
        "   python analyze_scenarios.py",
        "",
        "3) Start app and diagnosis page:",
        "   python webapp.py",
        "   open http://127.0.0.1:5000/diagnosis",
    ]

    cb = s.shapes.add_shape(1, Inches(1.1), Inches(1.8), Inches(11.1), Inches(4.8))
    cb.fill.solid(); cb.fill.fore_color.rgb = C_WHITE
    cb.line.color.rgb = C_ACCENT
    cb.line.width = Pt(2)

    ctb = s.shapes.add_textbox(Inches(1.4), Inches(2.1), Inches(10.5), Inches(4.2))
    tf = ctb.text_frame
    tf.clear()
    for i, line in enumerate(cmds):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16 if line and line[0].isdigit() else 14)
        p.font.bold = bool(line and line[0].isdigit())
        p.font.color.rgb = C_TITLE

    add_footer(s)

    # Slide 10: Closing summary
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_title(s, "Executive Summary", "M9 is currently the best candidate")

    add_bullet_card(
        s, 1.0, 1.8, 5.6, 4.2, "Key facts", [
            f"M9 R2: {m9_m['overall']['r2']:.4f} (best among M7/M8/M9)",
            f"M9 params: {m9_m['params']:,} (near M8 compact size)",
            "SwiGLU provides stronger feature gating",
            "Can reuse M8 Stage B artifacts to save time"
        ], C_M9
    )
    add_bullet_card(
        s, 6.9, 1.8, 5.4, 4.2, "Recommended next actions", [
            "Keep M9 as main deployment model",
            "Monitor KR-heavy edge cases from diagnosis",
            "Optional: tune prototype count around k*=5",
            "Optional: add targeted data for worst scenarios"
        ], C_ACCENT
    )

    add_footer(s, "M7-M8-M9 comparative deck")

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build_ppt()
