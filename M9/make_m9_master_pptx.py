"""
Master M7-M8-M9 comparison deck with detailed visuals.
Focus:
- detailed M9 analysis
- architecture drawings for M7, M8, M9
- Learnable relation matrix (logits -> softmax)
- SwiGLU gated MLP explanation
- parameter and error-rate comparisons
"""

from pathlib import Path
import json

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


BASE_DIR = Path(__file__).resolve().parent
M7_JSON = BASE_DIR.parent / "M7" / "comparison.json"
M8_JSON = BASE_DIR.parent / "M8" / "comparison.json"
M9_JSON = BASE_DIR / "comparison.json"
M9_ADAPTER_JSON = BASE_DIR / "comparison_cross_attn_adapter.json"
M9_SCENARIO_JSON = BASE_DIR / "scenario_analysis.json"
OUT_PATH = BASE_DIR / "M7_M8_M9_master_analysis.pptx"

# palette
C_BG = RGBColor(247, 246, 241)
C_DARK = RGBColor(23, 34, 45)
C_SUB = RGBColor(90, 106, 122)
C_WHITE = RGBColor(255, 255, 255)
C_M7 = RGBColor(45, 110, 210)
C_M8 = RGBColor(227, 120, 80)
C_M9 = RGBColor(44, 157, 139)
C_ACCENT = RGBColor(240, 165, 91)
C_MUTE = RGBColor(216, 221, 227)
C_RED = RGBColor(188, 55, 49)
C_GREEN = RGBColor(39, 167, 103)


def load_json(path: Path):
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def set_bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None, dark=False):
    tc = C_WHITE if dark else C_DARK
    sc = RGBColor(185, 196, 210) if dark else C_SUB

    t = slide.shapes.add_textbox(Inches(0.45), Inches(0.15), Inches(12.3), Inches(0.9))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = tc

    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12.0), Inches(0.45))
        stf = s.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(15)
        sr.font.color.rgb = sc


def add_footer(slide, text="PFE | M7-M8-M9 Master Analysis"):
    b = slide.shapes.add_textbox(Inches(0.45), Inches(6.95), Inches(12.3), Inches(0.28))
    tf = b.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9)
    r.font.color.rgb = C_SUB


def add_text(slide, x, y, w, h, lines, size=13, bold_first=False, color=C_DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        if i == 0 and bold_first:
            r.font.bold = True


def rect(slide, x, y, w, h, fill=C_WHITE, border=C_M7, bw=1.5, rounded=True):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(bw)
    return sh


def arrow(slide, x, y, w=0.55, color=C_ACCENT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.25))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()


def down_arrow(slide, x, y, h=0.24, color=C_ACCENT):
    sh = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(0.22), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()


def card(slide, x, y, w, h, head, bullets, color, body_size=12):
    rect(slide, x, y, w, h, fill=C_WHITE, border=color, bw=2)
    add_text(slide, x + 0.12, y + 0.06, w - 0.25, 0.35, [head], size=15, bold_first=True, color=color)
    lines = [f"- {b}" for b in bullets]
    add_text(slide, x + 0.14, y + 0.45, w - 0.28, h - 0.55, lines, size=body_size)


def layer_column(slide, cx, top, label, color, layers, total_txt):
    bw = 2.55
    bh = 0.53
    x = cx - bw / 2
    y = top

    add_text(slide, x, y - 0.42, bw, 0.34, [label], size=17, bold_first=True, color=color, align=PP_ALIGN.CENTER)

    for i, layer in enumerate(layers):
        name, sub, fill = layer
        rect(slide, x, y, bw, bh, fill=fill, border=color)
        add_text(slide, x + 0.03, y + 0.03, bw - 0.06, 0.24, [name], size=10, bold_first=True, color=C_DARK, align=PP_ALIGN.CENTER)
        if sub:
            add_text(slide, x + 0.04, y + 0.25, bw - 0.08, 0.22, [sub], size=8, color=C_SUB, align=PP_ALIGN.CENTER)
        y += bh
        if i < len(layers) - 1:
            down_arrow(slide, cx - 0.11, y, h=0.18, color=color)
            y += 0.2

    add_text(slide, x, y + 0.04, bw, 0.24, [total_txt], size=11, bold_first=True, color=color, align=PP_ALIGN.CENTER)


def styled_table(slide, x, y, w, h, headers, rows, row_colors=None):
    table = slide.shapes.add_table(1 + len(rows), len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for ci, hdr in enumerate(headers):
        c = table.cell(0, ci)
        c.text = hdr
        c.fill.solid()
        c.fill.fore_color.rgb = C_DARK
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(11)
                r.font.color.rgb = C_WHITE

    for ri, row in enumerate(rows):
        for ci, value in enumerate(row):
            c = table.cell(ri + 1, ci)
            c.text = str(value)
            if row_colors and ri < len(row_colors):
                c.fill.solid()
                c.fill.fore_color.rgb = row_colors[ri]
            for p in c.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(10)


def build():
    m7 = load_json(M7_JSON)
    m8 = load_json(M8_JSON)
    m9 = load_json(M9_JSON)
    m9a = load_json(M9_ADAPTER_JSON) if M9_ADAPTER_JSON.exists() else None
    s9 = load_json(M9_SCENARIO_JSON)

    m6 = m7["teacher_m6"]
    m7m = m7["psi_model_m7"]
    m8m = m8["psi_model_m8"]
    m9m = m9["psi_model_m9"]
    m9s = m9["student_stage_a"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # slide 1
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_DARK)
    add_title(s, "M7 vs M8 vs M9", "Master Technical Analysis with Full Architecture Drawings", dark=True)
    for i, info in enumerate([
        ("M7 Baseline", C_M7, m7m["params"], m7m["overall"]["r2"]),
        ("M8 Efficient", C_M8, m8m["params"], m8m["overall"]["r2"]),
        ("M9 SwiGLU", C_M9, m9m["params"], m9m["overall"]["r2"]),
    ]):
        lbl, color, pcount, r2 = info
        x = 0.9 + i * 4.05
        rect(s, x, 2.25, 3.7, 2.15, fill=color, border=color, bw=0)
        add_text(s, x + 0.17, 2.38, 3.35, 1.95, [
            lbl,
            "",
            f"Parameters: {pcount:,}",
            f"R2: {r2:.4f}",
        ], size=16, bold_first=True, color=C_WHITE)
    add_text(s, 0.95, 5.0, 11.9, 1.2, [
        "This deck explains model evolution, why M9 architecture matters, and how each design choice impacts accuracy, compression, and robustness.",
    ], size=14, color=RGBColor(208, 219, 231))
    add_footer(s)

    # slide 2
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Executive Snapshot", "What changed from M7 to M8 to M9")
    card(s, 0.45, 1.55, 4.1, 4.8, "M7 Core", [
        "Wide slot MLP: 64 -> 128 -> 64",
        "Fixed relation matrix",
        "k*=6 prototypes",
        f"R2={m7m['overall']['r2']:.4f}",
        f"Params={m7m['params']:,}",
    ], C_M7)
    card(s, 4.65, 1.55, 4.1, 4.8, "M8 Core", [
        "Bottleneck MLP: 64 -> 48 -> 64",
        "Learnable relation matrix",
        "Physics monotonic loss",
        "k*=5 prototypes",
        f"R2={m8m['overall']['r2']:.4f}",
        f"Params={m8m['params']:,}",
    ], C_M8)
    card(s, 8.85, 1.55, 4.0, 4.8, "M9 Core", [
        "SwiGLU gated MLP (hidden=32)",
        "Learnable relation matrix (same principle as M8)",
        "Full A/B/C retraining in latest run",
        "k*=5 prototypes",
        f"R2={m9m['overall']['r2']:.4f}",
        f"Params={m9m['params']:,}",
    ], C_M9)
    add_footer(s)

    # slide 3 outputs and objective
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Problem and Outputs", "Predict pile stiffness degradation across 21 steps")
    rect(s, 0.6, 1.7, 3.7, 4.9, fill=RGBColor(231, 242, 255), border=C_M7)
    add_text(s, 0.82, 1.82, 3.3, 4.6, [
        "Input Features (8)",
        "PI, Gmax, v, Dp, Tp, Lp, Ip, Dp/Lp",
        "",
        "Model sees one scenario and predicts full time response.",
        "",
        "Data:",
        "- 1019 scenarios",
        "- Train 815",
        "- Test 204",
    ], size=13, bold_first=True, color=C_M7)
    arrow(s, 4.45, 3.95, w=0.8)
    rect(s, 5.35, 1.7, 2.7, 4.9, fill=RGBColor(255, 247, 236), border=C_ACCENT)
    add_text(s, 5.53, 1.82, 2.35, 4.6, [
        "Psi-NN Core",
        "Slot attention",
        "Per-slot MLP",
        "Relation matrix R",
        "Output head",
        "",
        "Differences between M7/M8/M9 happen here.",
    ], size=13, bold_first=True, align=PP_ALIGN.CENTER)
    arrow(s, 8.15, 3.95, w=0.8)
    rect(s, 9.05, 1.7, 3.7, 4.9, fill=RGBColor(231, 249, 242), border=C_M9)
    add_text(s, 9.24, 1.82, 3.3, 4.6, [
        "Outputs (3 x 21)",
        "KL: lateral stiffness",
        "KR: rocking stiffness",
        "KLR: coupling stiffness",
        "",
        "Targets:",
        "- High R2",
        "- Low RMSE / MAE",
        "- Physically consistent trends",
    ], size=13, bold_first=True, color=C_M9)
    add_footer(s)

    # slide 4 m7 architecture drawing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "M7 Architecture (Drawn)", "Baseline structured model")
    m7_layers = [
        ("Input (8)", "Soil + pile features", RGBColor(198, 221, 251)),
        ("Embedding 8->64", "Feature lift", RGBColor(177, 205, 245)),
        ("Slot Attention", "1 initial + 6 prototype slots", RGBColor(157, 190, 238)),
        ("Slot MLP", "64->128->64 (GELU)", RGBColor(124, 170, 230)),
        ("Fixed Relation Matrix", "20x6, frozen", RGBColor(146, 184, 237)),
        ("Output Head", "64->3 then cumsum", RGBColor(171, 199, 242)),
    ]
    layer_column(s, 3.2, 1.85, "M7", C_M7, m7_layers, f"Total {m7m['params']:,}")
    card(s, 6.0, 1.6, 6.8, 5.1, "M7 Notes", [
        "Strong baseline accuracy",
        "High MLP parameter cost",
        "Relation matrix cannot adapt after initialization",
        "Best in this snapshot by pure R2 but weak compression",
    ], C_M7, body_size=13)
    add_footer(s)

    # slide 5 m8 architecture drawing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "M8 Architecture (Drawn)", "Efficiency upgrades + learnable structure")
    m8_layers = [
        ("Input (8)", "Same data interface", RGBColor(252, 222, 208)),
        ("Embedding 8->64", "Same", RGBColor(247, 204, 183)),
        ("Slot Attention", "1 initial + 5 prototypes", RGBColor(244, 188, 163)),
        ("Slot MLP", "64->48->64 bottleneck", RGBColor(231, 141, 104)),
        ("Learnable R", "20x5 logits->softmax", RGBColor(242, 176, 148)),
        ("Output + Physics", "KL/KR down, KLR up", RGBColor(247, 197, 175)),
    ]
    layer_column(s, 3.2, 1.85, "M8", C_M8, m8_layers, f"Total {m8m['params']:,}")
    card(s, 6.0, 1.6, 6.8, 5.1, "M8 Upgrades", [
        "MLP parameter drop by about 62% versus M7 MLP block",
        "Learnable relation matrix allows end-to-end adaptation",
        "Physics-aware loss improves physical plausibility",
        "Compression jumps to 19.7% vs M6",
    ], C_M8, body_size=13)
    add_footer(s)

    # slide 6 m9 architecture drawing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "M9 Architecture (Drawn)", "SwiGLU gating replaces plain MLP")
    m9_layers = [
        ("Input (8)", "Same", RGBColor(206, 241, 233)),
        ("Embedding 8->64", "Same", RGBColor(183, 232, 220)),
        ("Slot Attention", "1 initial + 5 prototypes", RGBColor(164, 224, 210)),
        ("SwiGLU MLP", "gate/value branches, hidden=32", RGBColor(74, 186, 165)),
        ("Learnable R", "20x5 logits->softmax", RGBColor(135, 214, 196)),
        ("Output + Physics", "same constraints", RGBColor(168, 226, 212)),
    ]
    layer_column(s, 3.2, 1.85, "M9", C_M9, m9_layers, f"Total {m9m['params']:,}")
    card(s, 6.0, 1.6, 6.8, 5.1, "M9 Intent", [
        "Keep M8 compactness while increasing representational power",
        "SwiGLU gating selects informative channels per slot",
        "Maintains learnable relation matrix and physics constraints",
        "Strong KL and KLR improvements in latest full retrain",
    ], C_M9, body_size=13)
    add_footer(s)

    # slide 7 side by side architecture
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Side-by-Side Architecture Map", "Pipeline is shared, key differences highlighted")
    layer_column(s, 2.2, 1.85, "M7", C_M7, [
        ("Slot MLP", "64->128->64", RGBColor(124, 170, 230)),
        ("R matrix", "fixed", RGBColor(146, 184, 237)),
        ("k*", "6 prototypes", RGBColor(171, 199, 242)),
    ], f"R2 {m7m['overall']['r2']:.4f}")
    layer_column(s, 6.7, 1.85, "M8", C_M8, [
        ("Slot MLP", "64->48->64", RGBColor(231, 141, 104)),
        ("R matrix", "learnable", RGBColor(242, 176, 148)),
        ("k*", "5 prototypes", RGBColor(247, 197, 175)),
    ], f"R2 {m8m['overall']['r2']:.4f}")
    layer_column(s, 11.1, 1.85, "M9", C_M9, [
        ("Slot MLP", "SwiGLU hidden=32", RGBColor(74, 186, 165)),
        ("R matrix", "learnable", RGBColor(135, 214, 196)),
        ("k*", "5 prototypes", RGBColor(168, 226, 212)),
    ], f"R2 {m9m['overall']['r2']:.4f}")
    add_footer(s)

    # slide 8 learnable relation matrix deep dive
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Learnable Relation Matrix (logits -> softmax)", "From fixed mixing to trainable prototype routing")

    rect(s, 0.55, 1.65, 5.8, 4.95, fill=C_WHITE, border=C_M8, bw=2)
    add_text(s, 0.75, 1.82, 5.4, 4.65, [
        "Definition",
        "R_logits is a trainable tensor with shape (20, k*).",
        "Each row corresponds to one reconstructed drop slot.",
        "",
        "Conversion to valid weights",
        "R_soft = softmax(R_logits, dim=1)",
        "This ensures each row sums to 1.",
        "",
        "Reconstruction",
        "drop_slots = R_soft @ prototype_slots",
        "",
        "Why this matters",
        "- Gradients update routing weights",
        "- Matrix can adapt to data and loss",
        "- Better than frozen heuristic weights",
    ], size=12, bold_first=True)

    rect(s, 6.55, 1.65, 6.25, 2.35, fill=RGBColor(255, 246, 236), border=C_ACCENT, bw=2)
    add_text(s, 6.77, 1.83, 5.8, 2.0, [
        "Visual pipeline",
        "prototype slots (k* vectors) -> R logits -> row softmax -> weighted sum -> 20 drop slots",
        "",
        "In M7: R is fixed. In M8/M9: R is learned.",
    ], size=12, bold_first=True)

    rect(s, 6.55, 4.15, 6.25, 2.45, fill=C_WHITE, border=C_M9, bw=2)
    add_text(s, 6.77, 4.33, 5.8, 2.1, [
        "Observed in M9 full retrain",
        f"- k* = {m9m['k_star']}",
        "- Learned R max-per-row average = 1.000",
        "- Sparsity (<0.05 entries): 80/100",
        "",
        "Interpretation: routing becomes highly selective.",
    ], size=12, bold_first=True)
    add_footer(s)

    # slide 9 swiglu deep dive
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "SwiGLU Gated MLP", "How M9 increases expressiveness without large parameter growth")

    rect(s, 0.55, 1.65, 6.1, 4.95, fill=C_WHITE, border=C_M9, bw=2)
    add_text(s, 0.75, 1.82, 5.7, 4.6, [
        "SwiGLU equation",
        "gate = SiLU(W_gate x)",
        "value = W_val x",
        "out = W_out (gate * value)",
        "",
        "Mechanism",
        "- gate branch controls importance per dimension",
        "- value branch carries content",
        "- element-wise multiplication acts as dynamic filter",
        "",
        "Result",
        "More expressive than plain GELU MLP at similar size",
    ], size=12, bold_first=True)

    # dual path drawing
    rect(s, 7.0, 1.8, 1.7, 0.5, fill=RGBColor(206, 241, 233), border=C_M9)
    add_text(s, 7.05, 1.93, 1.6, 0.25, ["x (64)"], size=10, align=PP_ALIGN.CENTER)
    down_arrow(s, 7.72, 2.3, h=0.2, color=C_M9)

    rect(s, 6.7, 2.55, 2.35, 0.58, fill=RGBColor(149, 222, 207), border=C_M9)
    add_text(s, 6.75, 2.72, 2.25, 0.3, ["Gate path: SiLU(W_gate x)"], size=9, align=PP_ALIGN.CENTER)

    rect(s, 9.55, 2.55, 2.35, 0.58, fill=RGBColor(181, 232, 221), border=C_M9)
    add_text(s, 9.60, 2.72, 2.25, 0.3, ["Value path: W_val x"], size=9, align=PP_ALIGN.CENTER)

    arrow(s, 9.0, 2.8, w=0.45, color=C_MUTE)
    rect(s, 8.35, 3.45, 2.1, 0.55, fill=RGBColor(81, 189, 169), border=C_M9)
    add_text(s, 8.40, 3.63, 2.0, 0.28, ["Element-wise gate * value"], size=9, align=PP_ALIGN.CENTER, color=C_WHITE)
    down_arrow(s, 9.30, 4.0, h=0.2, color=C_M9)
    rect(s, 8.15, 4.35, 2.55, 0.6, fill=RGBColor(140, 214, 198), border=C_M9)
    add_text(s, 8.20, 4.55, 2.45, 0.3, ["W_out -> output (64)"], size=10, align=PP_ALIGN.CENTER)

    rect(s, 6.7, 5.15, 5.2, 1.45, fill=C_WHITE, border=C_ACCENT, bw=2)
    add_text(s, 6.9, 5.3, 4.8, 1.2, [
        "Parameter comparison of slot MLP block",
        "M7 wide: 16,512",
        "M8 bottleneck: 6,256",
        "M9 SwiGLU: 6,272",
    ], size=11, bold_first=True)
    add_footer(s)

    # slide 10 params chart and table
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Parameter Comparison", "Compression vs representation trade-off")

    cd = CategoryChartData()
    cd.categories = ["M6", "M7", "M8", "M9"]
    cd.add_series("Total params", [m6["params"], m7m["params"], m8m["params"], m9m["params"]])
    s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.75), Inches(6.3), Inches(4.1), cd)

    rows = [
        ["M7", f"{m7m['params']:,}", "1.5%", "Wide MLP 64->128->64"],
        ["M8", f"{m8m['params']:,}", "19.7%", "Bottleneck MLP 64->48->64"],
        ["M9", f"{m9m['params']:,}", "19.7%", "SwiGLU gated MLP (hidden=32)"],
    ]
    styled_table(s, 7.05, 1.75, 5.75, 3.0, ["Model", "Params", "Compression vs M6", "MLP core"], rows,
                 [RGBColor(232, 240, 252), RGBColor(252, 237, 229), RGBColor(228, 245, 240)])

    add_text(s, 7.1, 5.0, 5.7, 1.65, [
        "Interpretation",
        "M8/M9 remove about 10.3k parameters versus M7.",
        "M9 keeps M8-size efficiency while adding gating expressiveness.",
    ], size=12, bold_first=True)
    add_footer(s)

    # slide 11 error and accuracy
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Accuracy and Error Rate", "R2 higher is better, RMSE/MAE lower is better")

    r2c = CategoryChartData()
    r2c.categories = ["M7", "M8", "M9"]
    r2c.add_series("R2", [m7m["overall"]["r2"], m8m["overall"]["r2"], m9m["overall"]["r2"]])
    s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.8), Inches(4.1), Inches(3.6), r2c)

    rmsec = CategoryChartData()
    rmsec.categories = ["M7", "M8", "M9"]
    rmsec.add_series("RMSE (B)", [
        m7m["overall"]["rmse"] / 1e9,
        m8m["overall"]["rmse"] / 1e9,
        m9m["overall"]["rmse"] / 1e9,
    ])
    s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(4.75), Inches(1.8), Inches(4.1), Inches(3.6), rmsec)

    maec = CategoryChartData()
    maec.categories = ["M7", "M8", "M9"]
    maec.add_series("MAE (B)", [
        m7m["overall"]["mae"] / 1e9,
        m8m["overall"]["mae"] / 1e9,
        m9m["overall"]["mae"] / 1e9,
    ])
    s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(9.0), Inches(1.8), Inches(3.8), Inches(3.6), maec)

    styled_table(s, 0.5, 5.55, 12.3, 1.25, ["Model", "R2", "RMSE", "MAE", "Error rate (1-R2)"], [
        ["M7", f"{m7m['overall']['r2']:.4f}", f"{m7m['overall']['rmse']:.3e}", f"{m7m['overall']['mae']:.3e}", f"{1-m7m['overall']['r2']:.4f}"],
        ["M8", f"{m8m['overall']['r2']:.4f}", f"{m8m['overall']['rmse']:.3e}", f"{m8m['overall']['mae']:.3e}", f"{1-m8m['overall']['r2']:.4f}"],
        ["M9", f"{m9m['overall']['r2']:.4f}", f"{m9m['overall']['rmse']:.3e}", f"{m9m['overall']['mae']:.3e}", f"{1-m9m['overall']['r2']:.4f}"],
    ], [RGBColor(232, 240, 252), RGBColor(252, 237, 229), RGBColor(228, 245, 240)])
    add_footer(s)

    # slide 12 per variable
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Per-Variable Performance", "KL, KR, KLR detailed view")
    styled_table(s, 0.5, 1.7, 12.3, 2.95,
                 ["Model", "KL R2", "KR R2", "KLR R2", "KL RMSE", "KR RMSE", "KLR RMSE"],
                 [
                     ["M7", f"{m7m['per_variable']['KL']['r2']:.4f}", f"{m7m['per_variable']['KR']['r2']:.4f}", f"{m7m['per_variable']['KLR']['r2']:.4f}",
                      f"{m7m['per_variable']['KL']['rmse']:.3e}", f"{m7m['per_variable']['KR']['rmse']:.3e}", f"{m7m['per_variable']['KLR']['rmse']:.3e}"],
                     ["M8", f"{m8m['per_variable']['KL']['r2']:.4f}", f"{m8m['per_variable']['KR']['r2']:.4f}", f"{m8m['per_variable']['KLR']['r2']:.4f}",
                      f"{m8m['per_variable']['KL']['rmse']:.3e}", f"{m8m['per_variable']['KR']['rmse']:.3e}", f"{m8m['per_variable']['KLR']['rmse']:.3e}"],
                     ["M9", f"{m9m['per_variable']['KL']['r2']:.4f}", f"{m9m['per_variable']['KR']['r2']:.4f}", f"{m9m['per_variable']['KLR']['r2']:.4f}",
                      f"{m9m['per_variable']['KL']['rmse']:.3e}", f"{m9m['per_variable']['KR']['rmse']:.3e}", f"{m9m['per_variable']['KLR']['rmse']:.3e}"],
                 ],
                 [RGBColor(232, 240, 252), RGBColor(252, 237, 229), RGBColor(228, 245, 240)])
    card(s, 0.5, 4.95, 6.1, 1.75, "Key observations", [
        "M9 has the best KL and KLR R2 among M7/M8/M9 in this run",
        "KR remains the hardest target and dominates RMSE scale",
    ], C_M9, body_size=12)
    card(s, 6.7, 4.95, 6.1, 1.75, "Error-rate takeaway", [
        "Absolute R2 differences are small between M7 and M9",
        "M9 advantage is efficiency + architectural scalability",
    ], C_ACCENT, body_size=12)
    add_footer(s)

    # slide 13 training pipeline and complexity
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Training Pipeline", "M9 full retraining: Stage A -> Stage B -> Stage C")
    rect(s, 0.6, 1.8, 3.85, 2.2, fill=RGBColor(233, 242, 255), border=C_M7, bw=2)
    add_text(s, 0.82, 1.98, 3.4, 1.9, [
        "Stage A: Distillation",
        "Train SwiGLU student from teacher",
        "Early stopped at epoch 1118",
        f"Student R2: {m9s['overall']['r2']:.4f}",
    ], size=12, bold_first=True)
    arrow(s, 4.58, 2.75, w=0.6)
    rect(s, 5.35, 1.8, 2.8, 2.2, fill=RGBColor(255, 242, 232), border=C_M8, bw=2)
    add_text(s, 5.50, 1.98, 2.5, 1.9, [
        "Stage B",
        "Cluster slots",
        "k*=5",
        "Silhouette=0.673",
    ], size=12, bold_first=True, align=PP_ALIGN.CENTER)
    arrow(s, 8.28, 2.75, w=0.6)
    rect(s, 9.05, 1.8, 3.85, 2.2, fill=RGBColor(229, 245, 239), border=C_M9, bw=2)
    add_text(s, 9.25, 1.98, 3.45, 1.9, [
        "Stage C: Structured M9",
        "Train final Psi model",
        "Early stopped at epoch 572",
        f"Final R2: {m9m['overall']['r2']:.4f}",
    ], size=12, bold_first=True)

    styled_table(s, 0.6, 4.35, 12.2, 2.15,
                 ["Item", "M7", "M8", "M9 (latest)"],
                 [
                     ["A/B/C pipeline", "Yes", "Yes", "Yes"],
                     ["k* discovered", "6", "5", "5"],
                     ["Relation matrix", "Fixed", "Learnable", "Learnable"],
                     ["MLP type", "Wide GELU", "Bottleneck GELU", "SwiGLU gated"],
                 ],
                 [RGBColor(240, 240, 240)] * 4)
    add_footer(s)

    # slide 14 adapter architecture extension (train new part only)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "M9 Cross-Attention Adapter: Exact Mechanism", "Code-level view of queries, attention weights, fusion, and prediction")

    rect(s, 0.55, 1.55, 6.0, 2.25, fill=RGBColor(229, 245, 239), border=C_M9, bw=2)
    add_text(s, 0.75, 1.72, 5.6, 1.95, [
        "1) Frozen encoder slots from base M9",
        "S_f in R^(B x 21 x 64) are produced by the frozen M9 core.",
        "No base M9 weights are updated in this branch.",
    ], size=11, bold_first=True)

    rect(s, 6.85, 1.55, 5.95, 2.25, fill=RGBColor(255, 246, 236), border=C_ACCENT, bw=2)
    add_text(s, 7.05, 1.72, 5.55, 1.95, [
        "2) Decoder-style cross-attention",
        "Q in R^(B x 21 x 64): learned decoder queries (trainable)",
        "K=V=S_f in R^(B x 21 x 64): frozen slot states",
        "A = softmax((Q K^T) / sqrt(64)),  CA = A V",
    ], size=11, bold_first=True)

    rect(s, 0.55, 3.95, 12.25, 1.25, fill=C_WHITE, border=C_DARK, bw=1.5)
    add_text(s, 0.80, 4.12, 11.8, 0.95, [
        "3) Fusion used in code: decoded = LN(Q + CA), gate = sigmoid(g), adapted_slots = LN(S_f + gate * decoded), then output heads predict KL/KR/KLR over 21 steps.",
    ], size=11, bold_first=True)

    # schematic arrows
    rect(s, 0.95, 5.35, 2.5, 0.72, fill=RGBColor(206, 241, 233), border=C_M9)
    add_text(s, 1.03, 5.57, 2.35, 0.34, ["Frozen slots S_f"], size=10, align=PP_ALIGN.CENTER)
    arrow(s, 3.58, 5.60, w=0.52, color=C_M9)
    rect(s, 4.15, 5.35, 2.7, 0.72, fill=RGBColor(255, 237, 220), border=C_ACCENT)
    add_text(s, 4.23, 5.57, 2.55, 0.34, ["Cross-attn: Attn(Q,K,V)"], size=10, align=PP_ALIGN.CENTER)
    arrow(s, 6.95, 5.60, w=0.52, color=C_ACCENT)
    rect(s, 7.52, 5.35, 2.1, 0.72, fill=RGBColor(242, 242, 242), border=C_DARK)
    add_text(s, 7.60, 5.57, 1.95, 0.34, ["Gated fusion + LN"], size=10, align=PP_ALIGN.CENTER)
    arrow(s, 9.72, 5.60, w=0.52, color=C_DARK)
    rect(s, 10.29, 5.35, 2.05, 0.72, fill=RGBColor(227, 245, 239), border=C_M9)
    add_text(s, 10.37, 5.57, 1.90, 0.34, ["Predictions"], size=10, align=PP_ALIGN.CENTER)

    if m9a:
        base = m9a["base_m9"]["overall"]
        adap = m9a["m9_cross_attn_adapter"]["overall"]
        row_colors = [RGBColor(232, 240, 252), RGBColor(252, 237, 229)]
        styled_table(
            s, 0.55, 6.20, 12.25, 0.75,
            ["Model", "R2", "RMSE", "MAE", "Trainable Params"],
            [
                ["Base M9", f"{base['r2']:.4f}", f"{base['rmse']:.3e}", f"{base['mae']:.3e}", "0 (frozen experiment branch)"],
                ["M9 + Cross-Attn Adapter", f"{adap['r2']:.4f}", f"{adap['rmse']:.3e}", f"{adap['mae']:.3e}", f"{m9a['m9_cross_attn_adapter']['params_trainable']:,}"],
            ],
            row_colors,
        )
    else:
        card(s, 0.55, 5.90, 12.25, 1.05, "Adapter metrics", [
            "comparison_cross_attn_adapter.json was not found at generation time.",
            "Run train_cross_attn_adapter.py to populate experiment results.",
        ], C_ACCENT, body_size=11)
    add_footer(s)

    # slide 15 robustness with scenario analysis
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "M9 Robustness Analysis", "Scenario-level diagnostics")

    g = s9["global_metrics"]
    w1 = s9["worst_scenarios"][0]

    card(s, 0.55, 1.65, 5.9, 2.55, "Global quality", [
        f"Avg curve MAPE: {g['avg_curve_mape_pct']:.3f}%",
        f"Median MAPE: {g['median_curve_mape_pct']:.3f}%",
        f"Avg curve R2: {g['avg_curve_r2']:.5f}",
        f"P95 MAPE: {g['p95_curve_mape_pct']:.3f}%",
        f"Test scenarios: {g['n_test']}",
    ], C_M9, body_size=12)

    card(s, 6.75, 1.65, 6.0, 2.55, "Worst-case signal", [
        f"Worst scenario id: {w1['scenario_id']}",
        f"Worst avg MAPE: {w1['overall']['avg_curve_mape_pct']:.2f}%",
        f"Worst avg R2: {w1['overall']['avg_curve_r2']:.4f}",
        "Use uncertainty flag for extreme edge cases",
    ], C_RED, body_size=12)

    cats = CategoryChartData()
    ws = s9["worst_scenarios"][:6]
    cats.categories = [f"S{w['scenario_id']}" for w in ws]
    cats.add_series("MAPE %", [w["overall"]["avg_curve_mape_pct"] for w in ws])
    s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.55), Inches(4.45), Inches(6.2), Inches(2.2), cats)

    add_text(s, 6.95, 4.45, 5.7, 2.15, [
        "Interpretation",
        "Most scenarios are very accurate.",
        "A small tail of hard scenarios drives peak error.",
        "This pattern is common in compact physics-informed models.",
    ], size=12, bold_first=True)
    add_footer(s)

    # slide 16 recommendation
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Recommendation and Best-Model View", "Decision using accuracy, compression, and architecture maturity")

    styled_table(s, 0.55, 1.75, 7.8, 3.0,
                 ["Criterion", "M7", "M8", "M9"],
                 [
                     ["Overall R2", f"{m7m['overall']['r2']:.4f}", f"{m8m['overall']['r2']:.4f}", f"{m9m['overall']['r2']:.4f}"],
                     ["Compression vs M6", "1.5%", "19.7%", "19.7%"],
                     ["Learnable relation matrix", "No", "Yes", "Yes"],
                     ["Gated MLP", "No", "No", "Yes"],
                     ["Scalability for future tuning", "Medium", "High", "Highest"],
                 ],
                 [RGBColor(233, 239, 249), RGBColor(251, 237, 229), RGBColor(227, 245, 239), RGBColor(240, 240, 240), RGBColor(227, 245, 239)])

    rect(s, 8.6, 1.75, 4.2, 3.0, fill=C_M9, border=C_M9, bw=0)
    add_text(s, 8.8, 1.95, 3.8, 2.65, [
        "Best practical choice",
        "M9",
        "",
        "Why:",
        "- near-top accuracy",
        "- major compression",
        "- learnable routing",
        "- SwiGLU expressiveness",
        "- strongest foundation for next tuning round",
    ], size=13, bold_first=True, color=C_WHITE)

    card(s, 0.55, 5.05, 12.25, 1.65, "Next tuning to beat current R2", [
        "Increase Stage C patience and lower LR for finer convergence",
        "Run 3-5 seeds and keep best checkpoint",
        "Augment hard scenarios from the top error tail",
    ], C_ACCENT, body_size=12)
    add_footer(s)

    # slide 17 closing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_DARK)
    add_title(s, "Closing Summary", "M9 architecture value in one page", dark=True)
    add_text(s, 0.85, 1.55, 11.8, 4.9, [
        "M7 -> M8 -> M9 evolution",
        "",
        "M7 established the baseline structured Psi-NN.",
        "M8 introduced compactness and learnable relation routing.",
        "M9 adds SwiGLU gated slot MLP on top of that structure.",
        "",
        "Learnable relation matrix (logits -> softmax)",
        "lets the model learn prototype-to-time-step assignment end-to-end.",
        "",
        "SwiGLU gated MLP",
        "improves expressive filtering at almost identical parameter budget.",
        "",
        f"Current full-retrain M9: R2={m9m['overall']['r2']:.4f}, params={m9m['params']:,}.",
        "This is a strong compact model and a solid base for final hyperparameter optimization.",
    ], size=16, color=RGBColor(220, 231, 242), bold_first=True)
    add_footer(s)

    # slide 18 cross-attention explanation talk track
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "Cross-Attention Adapter: Explanation Slide", "Simple narrative for presenting how the adapter works")

    rect(s, 0.55, 1.55, 12.25, 1.55, fill=RGBColor(229, 245, 239), border=C_M9, bw=2)
    add_text(s, 0.80, 1.76, 11.8, 1.2, [
        "1) Frozen encoder slots from base M9",
        "S_f has shape (B, 21, 64): B scenarios per batch, 21 time-step slots, 64 latent features.",
        "These slots come from base M9 and stay frozen, so the adapter learns a correction without changing base knowledge.",
    ], size=12, bold_first=True)

    rect(s, 0.55, 3.28, 12.25, 1.65, fill=RGBColor(255, 246, 236), border=C_ACCENT, bw=2)
    add_text(s, 0.80, 3.49, 11.8, 1.3, [
        "2) Decoder-style cross-attention",
        "Q is trainable with shape (B, 21, 64), while K = V = S_f are frozen context slots.",
        "Attention weights are A = softmax((Q K^T) / sqrt(64)), then retrieved context is CA = A V.",
        "Intuition: Q asks what to retrieve; A decides where to focus; CA is the retrieved information.",
    ], size=12, bold_first=True)

    rect(s, 0.55, 5.05, 12.25, 1.55, fill=C_WHITE, border=C_DARK, bw=2)
    add_text(s, 0.80, 5.25, 11.8, 1.25, [
        "3) Fusion used in code and final prediction",
        "decoded = LN(Q + CA), gate = sigmoid(g), adapted_slots = LN(S_f + gate * decoded).",
        "The output heads then predict KL, KR, and KLR across 21 steps from adapted_slots.",
        "Gate controls correction strength: low gate keeps base behavior; high gate applies stronger adapter updates.",
    ], size=12, bold_first=True)
    add_footer(s)

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
