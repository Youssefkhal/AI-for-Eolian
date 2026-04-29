"""
M9 vs M10 — Conceptual Explanation PPTX
========================================
No code. Visual pipeline diagrams explaining:
  - How the Ψ-NN pipeline works
  - What LRP is and how it traces relevance to inputs
  - How Attention Rollout aggregates information flow
  - How Token Attribution Maps reveal feature importance
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colours ──
BG       = "0F1023"
CARD     = "1A1B3A"
CARD2    = "1E1F42"
PURPLE   = "C084FC"
PURDK    = "7C3AED"
CYAN     = "00D2FF"
GREEN    = "00FF88"
ORANGE   = "FFA500"
RED      = "FF5757"
YELLOW   = "FACC15"
WHITE    = "FFFFFF"
GREY     = "AAAAAA"
LIGHT    = "E2E8F0"
PINK     = "F472B6"

H = lambda c: RGBColor(int(c[:2],16), int(c[2:4],16), int(c[4:],16))

def bg(sl):
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = H(BG)

def box(sl, l, t, w, h, fill, border=None):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = H(fill)
    if border:
        s.line.color.rgb = H(border); s.line.width = Pt(1.2)
    else:
        s.line.fill.background()
    return s

def rect(sl, l, t, w, h, fill, border=None):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = H(fill)
    if border:
        s.line.color.rgb = H(border); s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def arrow_r(sl, l, t, w, h, color):
    s = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = H(color)
    s.line.fill.background()
    return s

def arrow_d(sl, l, t, w, h, color):
    s = sl.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = H(color)
    s.line.fill.background()
    return s

def chevron(sl, l, t, w, h, color):
    s = sl.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = H(color)
    s.line.fill.background()
    return s

def txt(sl, l, t, w, h, text, sz=12, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = H(color)
    p.font.bold = bold; p.font.name = font; p.alignment = align
    return tf

def bullet(tf, text, sz=11, color=WHITE, bold=False, sp=2):
    p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = H(color)
    p.font.bold = bold; p.font.name = "Segoe UI"
    p.space_before = Pt(sp); p.space_after = Pt(0)
    return p

def stat(sl, l, t, w, label, value, col=CYAN):
    box(sl, l, t, w, 0.7, CARD, col)
    txt(sl, l+0.05, t+0.02, w-0.1, 0.22, label, 7, GREY, align=PP_ALIGN.CENTER)
    txt(sl, l+0.05, t+0.22, w-0.1, 0.42, value, 15, col, bold=True, align=PP_ALIGN.CENTER)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ═══════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)

txt(sl, 0.5, 1.0, 12.3, 0.8,
    "M9 → M10", 40, PURPLE, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 1.8, 12.3, 0.6,
    "Adding Explainable AI to the Physics-Informed Digital Twin",
    20, LIGHT, align=PP_ALIGN.CENTER)

box(sl, 2.0, 2.8, 9.3, 0.8, CARD2, PURPLE)
txt(sl, 2.2, 2.9, 8.9, 0.6,
    "How does the model decide?  Which input features matter?  "
    "Where does information flow inside the network?",
    13, GREY, align=PP_ALIGN.CENTER)

# Three pillars
pillars = [
    ("Layer-wise Relevance\nPropagation (LRP)", "Traces each prediction\nback to the 8 soil/pile inputs", GREEN),
    ("Attention Rollout", "Maps the flow of information\nthrough cross & self attention", CYAN),
    ("Token Attribution\nMaps", "Visualises which features\ndrive each degradation step", ORANGE),
]
for i, (title, desc, col) in enumerate(pillars):
    x = 1.5 + i * 3.7
    box(sl, x, 4.0, 3.3, 2.0, CARD, col)
    txt(sl, x+0.15, 4.15, 3.0, 0.6, title, 14, col, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, x+0.15, 4.85, 3.0, 0.7, desc, 10.5, LIGHT, align=PP_ALIGN.CENTER)

txt(sl, 0.5, 6.4, 12.3, 0.5,
    "Pile Stiffness Degradation under Cyclic Loading  —  SwiGLU Slot-Attention Ψ-NN",
    11, PURDK, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 6.9, 12.3, 0.35,
    "Youssef Khalil  •  Shenzhen University / EMSI  •  2026",
    9.5, GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 2 — The M9 Base Model (visual pipeline)
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)

txt(sl, 0.5, 0.25, 12.3, 0.5,
    "The Base Model: SwiGLU Ψ-NN (M9, unchanged in M10)",
    22, CYAN, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 0.75, 12.3, 0.35,
    "Understanding the architecture is the first step to explaining its predictions",
    11, GREY, align=PP_ALIGN.CENTER)

# Input box
box(sl, 0.3, 1.5, 2.0, 1.6, "1A2A4A", CYAN)
txt(sl, 0.4, 1.6, 1.8, 0.3, "8 Soil & Pile\nInputs", 13, CYAN, bold=True, align=PP_ALIGN.CENTER)
inputs = ["PI  (Plasticity Index)", "Gmax  (Shear Modulus)", "ν  (Poisson Ratio)",
          "Dp, Tp, Lp  (Geometry)", "Ip  (Moment of Inertia)", "Dp/Lp  (Slenderness)"]
tf = txt(sl, 0.4, 2.2, 1.8, 0.8, "", 1)
for inp in inputs:
    bullet(tf, f"▸ {inp}", 7.5, LIGHT, sp=1)

arrow_r(sl, 2.4, 2.0, 0.5, 0.35, CYAN)

# Embedding
box(sl, 3.0, 1.5, 1.8, 1.6, "1A2A4A", CYAN)
txt(sl, 3.1, 1.6, 1.6, 0.3, "Embedding", 12, CYAN, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 3.1, 1.95, 1.6, 0.9,
    "Transforms 8 numbers\ninto a 64-dimensional\nrepresentation the\nattention layers\ncan process", 8.5, LIGHT, align=PP_ALIGN.CENTER)

arrow_r(sl, 4.9, 2.0, 0.5, 0.35, CYAN)

# Iterative block
box(sl, 5.5, 1.3, 4.5, 2.2, "1E2840", YELLOW)
txt(sl, 5.6, 1.35, 4.3, 0.3, "Iterative Refinement  ×3 rounds", 12, YELLOW, bold=True, align=PP_ALIGN.CENTER)

# Inner boxes
box(sl, 5.7, 1.8, 1.25, 1.0, CARD2, CYAN)
txt(sl, 5.75, 1.85, 1.15, 0.3, "Cross-\nAttention", 9, CYAN, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 5.75, 2.25, 1.15, 0.5, "Slots ask:\n'What input\ninfo do I need?'", 7, LIGHT, align=PP_ALIGN.CENTER)

box(sl, 7.1, 1.8, 1.25, 1.0, CARD2, ORANGE)
txt(sl, 7.15, 1.85, 1.15, 0.3, "Self-\nAttention", 9, ORANGE, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 7.15, 2.25, 1.15, 0.5, "Slots ask:\n'What do other\nslots know?'", 7, LIGHT, align=PP_ALIGN.CENTER)

box(sl, 8.5, 1.8, 1.25, 1.0, CARD2, GREEN)
txt(sl, 8.55, 1.85, 1.15, 0.3, "SwiGLU\nMLP", 9, GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 8.55, 2.25, 1.15, 0.5, "Decides which\ninformation to\nkeep or discard", 7, LIGHT, align=PP_ALIGN.CENTER)

# Small arrows between inner boxes
arrow_r(sl, 7.0, 2.15, 0.1, 0.15, GREY)
arrow_r(sl, 8.4, 2.15, 0.1, 0.15, GREY)

arrow_r(sl, 10.1, 2.0, 0.5, 0.35, YELLOW)

# Output
box(sl, 10.7, 1.5, 2.3, 1.6, "1A2A4A", GREEN)
txt(sl, 10.8, 1.6, 2.1, 0.3, "Predictions", 13, GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 10.8, 1.95, 2.1, 1.0,
    "KL — Lateral stiffness\nKR — Rotational stiffness\nKLR — Coupled stiffness\n\n21 degradation steps\nfor each variable", 8.5, LIGHT, align=PP_ALIGN.CENTER)

# Slot explanation
box(sl, 0.3, 3.8, 12.7, 1.5, CARD, PURPLE)
txt(sl, 0.5, 3.9, 3.0, 0.3, "What are 'Slots'?", 14, PURPLE, bold=True)
txt(sl, 0.5, 4.25, 6.0, 1.0,
    "The model uses 21 slots — think of them as 21 specialised workers:\n\n"
    "▸ Slot 1 is the 'initial condition' worker — predicts starting stiffness\n"
    "▸ Slots 2–21 are 'degradation' workers — each one predicts how much\n"
    "   stiffness drops at that step of the cyclic loading process",
    9.5, LIGHT)

txt(sl, 6.8, 4.25, 6.0, 1.0,
    "These slots don't work independently — they communicate through:\n\n"
    "▸ Cross-attention: 'Let me check the soil properties again'\n"
    "▸ Self-attention: 'What did the other degradation steps learn?'\n"
    "▸ SwiGLU gating: 'I'll keep the relevant info, discard the noise'",
    9.5, LIGHT)

# M9 → M10 transition
box(sl, 0.3, 5.6, 12.7, 1.6, CARD, PURPLE)
txt(sl, 0.5, 5.7, 12.3, 0.35,
    "M9 → M10: Same model, new power — Explainability", 15, PURPLE, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 6.15, 12.3, 0.9,
    "M9 was a black box — it made accurate predictions but couldn't explain WHY.\n"
    "M10 adds an XAI layer that looks inside the attention mechanism and traces predictions back to inputs.\n"
    "No model retraining needed. Same weights, same accuracy. We simply observe and explain.",
    11, LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 3 — What is LRP? (conceptual)
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)

txt(sl, 0.5, 0.25, 12.3, 0.5,
    "Layer-wise Relevance Propagation (LRP)", 24, GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 0.75, 12.3, 0.35,
    "Tracing the model's prediction back to the 8 raw input features",
    12, GREY, align=PP_ALIGN.CENTER)

# The big idea
box(sl, 0.3, 1.3, 6.2, 2.5, CARD, GREEN)
txt(sl, 0.5, 1.4, 5.8, 0.3, "The Big Idea", 15, GREEN, bold=True)
txt(sl, 0.5, 1.8, 5.8, 1.8,
    "Imagine the model predicts KL = 5.2×10⁹ at step 10.\n\n"
    "LRP asks: 'How much did each input feature contribute\n"
    "to this specific number?'\n\n"
    "It works by computing how sensitive the output is to\n"
    "each input, then weighting by the input's actual value:\n\n"
    "   Relevance of feature j  =  (input value)  ×  (sensitivity)",
    10.5, LIGHT)

# Visual: flow diagram
box(sl, 6.8, 1.3, 6.2, 2.5, CARD, CYAN)
txt(sl, 7.0, 1.4, 5.8, 0.3, "How It Flows", 15, CYAN, bold=True)

# Input features (left)
feats = [("PI", RED), ("Gmax", ORANGE), ("ν", YELLOW), ("Dp/Lp", GREEN)]
for i, (f, c) in enumerate(feats):
    y = 1.9 + i * 0.45
    box(sl, 7.1, y, 0.8, 0.35, CARD2, c)
    txt(sl, 7.15, y+0.02, 0.7, 0.3, f, 9, c, bold=True, align=PP_ALIGN.CENTER)

# Arrow fan
for i in range(4):
    y = 2.05 + i * 0.45
    arrow_r(sl, 8.0, y, 0.5, 0.12, GREY)

# Model box
box(sl, 8.6, 1.9, 1.6, 1.7, CARD2, PURPLE)
txt(sl, 8.65, 2.1, 1.5, 0.7, "    Ψ-NN\n    Model\n\n  (black box)", 9, PURPLE, align=PP_ALIGN.CENTER)

arrow_r(sl, 10.3, 2.55, 0.5, 0.12, GREY)

# Output
box(sl, 10.9, 2.2, 1.8, 0.9, CARD2, GREEN)
txt(sl, 11.0, 2.3, 1.6, 0.7, "KL at Step 10\n= 5.2 × 10⁹", 10, GREEN, bold=True, align=PP_ALIGN.CENTER)

# Backward arrows (relevance flowing back)
txt(sl, 8.0, 3.55, 2.5, 0.3, "← Relevance flows backward", 8.5, YELLOW, bold=True)

# Result section
box(sl, 0.3, 4.0, 12.7, 1.5, CARD, ORANGE)
txt(sl, 0.5, 4.1, 5.0, 0.3, "What LRP Reveals", 15, ORANGE, bold=True)

# Example attribution bars
attrs = [
    ("Dp/Lp (Slenderness)", 73, GREEN),
    ("Gmax (Shear Modulus)", 14, CYAN),
    ("Ip (Moment of Inertia)", 4.4, PURPLE),
    ("Lp (Pile Length)", 4.0, ORANGE),
    ("PI (Plasticity Index)", 2.0, YELLOW),
    ("Tp (Wall Thickness)", 1.3, PINK),
    ("ν (Poisson Ratio)", 0.8, RED),
    ("Dp (Pile Diameter)", 0.5, GREY),
]
for i, (name, pct, col) in enumerate(attrs):
    y = 4.1 + i * 0.17
    x_start = 5.5
    bar_w = max(pct / 73 * 5.5, 0.15)
    rect(sl, x_start, y, bar_w, 0.14, col)
    txt(sl, x_start + bar_w + 0.1, y-0.02, 1.5, 0.18, f"{pct}%", 7.5, col, bold=True)
    txt(sl, 3.0, y-0.02, 2.4, 0.18, name, 7.5, LIGHT)

txt(sl, 11.5, 4.1, 1.5, 0.3, "← KL variable", 8, GREEN, bold=True)

# Key insight
box(sl, 0.3, 5.7, 12.7, 1.5, CARD, YELLOW)
txt(sl, 0.5, 5.8, 12.3, 0.3,
    "Key Insight: LRP produces a 3D attribution map", 14, YELLOW, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 6.2, 12.3, 0.8,
    "For every combination of (variable × step × feature), we get a relevance score.\n"
    "That means 3 variables × 21 steps × 8 features = 504 individual attributions per scenario.\n"
    "This tells us not just WHICH features matter overall, but HOW they change across the degradation process.",
    11, LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 4 — LRP: The Gradient × Input Method
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)

txt(sl, 0.5, 0.25, 12.3, 0.5,
    "LRP Implementation: Gradient × Input", 22, GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 0.75, 12.3, 0.35,
    "A simple yet powerful method that satisfies the conservation axiom",
    12, GREY, align=PP_ALIGN.CENTER)

# The formula
box(sl, 0.3, 1.3, 12.7, 1.5, CARD, GREEN)
txt(sl, 0.5, 1.4, 6.0, 0.3, "The Core Formula", 15, GREEN, bold=True)

txt(sl, 0.5, 1.8, 12.3, 0.4,
    "Relevance of feature j  =  xⱼ  ×  ∂ output / ∂ xⱼ",
    16, YELLOW, bold=True, align=PP_ALIGN.CENTER, font="Cambria Math")

txt(sl, 0.5, 2.25, 12.3, 0.4,
    "input value          how much the output changes when this input changes slightly",
    10, GREY, align=PP_ALIGN.CENTER)

# Two components explained
box(sl, 0.3, 3.0, 6.0, 2.3, CARD, CYAN)
txt(sl, 0.5, 3.1, 5.6, 0.3, "Component 1: The Gradient  ∂ŷ/∂xⱼ", 13, CYAN, bold=True)
txt(sl, 0.5, 3.5, 5.6, 1.6,
    "The gradient measures sensitivity — if we nudge input\n"
    "feature j by a tiny amount, how much does the\n"
    "prediction change?\n\n"
    "▸ Large gradient → model is very sensitive to this feature\n"
    "▸ Small gradient → changing this feature has little effect\n"
    "▸ Sign tells direction: positive means feature pushes\n"
    "   prediction up, negative means it pushes it down",
    9.5, LIGHT)

box(sl, 6.6, 3.0, 6.4, 2.3, CARD, ORANGE)
txt(sl, 6.8, 3.1, 6.0, 0.3, "Component 2: The Input Value  xⱼ", 13, ORANGE, bold=True)
txt(sl, 6.8, 3.5, 6.0, 1.6,
    "Multiplying by the input value is crucial because:\n\n"
    "▸ A feature might have high sensitivity (large gradient)\n"
    "   but its actual value is near zero → low real contribution\n\n"
    "▸ Conversely, a moderate gradient on a large input value\n"
    "   means significant actual contribution to the prediction\n\n"
    "▸ This product captures the TRUE relevance of each feature",
    9.5, LIGHT)

# Conservation
box(sl, 0.3, 5.5, 12.7, 1.7, CARD, PURPLE)
txt(sl, 0.5, 5.6, 12.3, 0.3,
    "The Conservation Principle", 15, PURPLE, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 6.0, 12.3, 0.35,
    "Sum of all feature relevances  ≈  the model's prediction  (first-order Taylor approximation)",
    13, YELLOW, align=PP_ALIGN.CENTER, font="Cambria Math")
txt(sl, 0.5, 6.45, 12.3, 0.6,
    "This means relevance is not artificially created or lost — the total amount of 'importance'\n"
    "distributed across the 8 features equals the model's output. Every feature gets its fair share.",
    10.5, LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 5 — Attention Rollout: Cross-Attention
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)

txt(sl, 0.5, 0.25, 12.3, 0.5,
    "Attention Rollout: Cross-Attention",
    22, CYAN, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 0.75, 12.3, 0.35,
    "How much does each slot pay attention to the input features?",
    12, GREY, align=PP_ALIGN.CENTER)

# Visual: Input → 21 slots
box(sl, 0.3, 1.3, 2.5, 1.5, CARD, CYAN)
txt(sl, 0.4, 1.4, 2.3, 0.3, "Input Embedding", 12, CYAN, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.4, 1.75, 2.3, 0.9,
    "The 8 soil/pile features\ncompressed into one\n64-dim vector\n(single token)", 9, LIGHT, align=PP_ALIGN.CENTER)

# Attention arrows of varying thickness
arrow_r(sl, 2.9, 1.5, 0.7, 0.15, CYAN)  # thick
arrow_r(sl, 2.9, 1.75, 0.6, 0.10, "0088BB")
arrow_r(sl, 2.9, 1.95, 0.5, 0.08, "006688")
arrow_r(sl, 2.9, 2.1, 0.4, 0.06, "004455")
arrow_r(sl, 2.9, 2.25, 0.3, 0.05, "003333")
arrow_r(sl, 2.9, 2.4, 0.25, 0.04, "002222")

# Slots column
box(sl, 3.7, 1.3, 2.0, 1.5, CARD, GREEN)
txt(sl, 3.8, 1.4, 1.8, 0.25, "21 Slots", 12, GREEN, bold=True, align=PP_ALIGN.CENTER)
for i in range(6):
    y = 1.7 + i * 0.17
    label = f"Slot {i+1}" if i < 5 else "..."
    c = GREEN if i == 0 else LIGHT
    txt(sl, 3.9, y, 1.6, 0.15, label, 7, c, align=PP_ALIGN.CENTER)

# Explanation
box(sl, 6.0, 1.3, 6.9, 1.5, CARD, YELLOW)
txt(sl, 6.2, 1.4, 6.5, 0.3, "What Cross-Attention Shows", 13, YELLOW, bold=True)
txt(sl, 6.2, 1.8, 6.5, 0.9,
    "Each of the 21 slots sends a 'query' to the input embedding.\n"
    "The attention weight tells us how strongly that slot requested information.\n\n"
    "▸ Thick arrow = high attention (slot heavily depends on input features)\n"
    "▸ Thin arrow = low attention (slot is more self-sufficient)\n"
    "▸ This happens 3 times (3 iterations), and we average the results",
    9.5, LIGHT)

# Three iterations diagram
box(sl, 0.3, 3.1, 12.7, 2.0, CARD, PURPLE)
txt(sl, 0.5, 3.2, 5.0, 0.3, "Rollout Across 3 Iterations", 14, PURPLE, bold=True)

iters = ["Iteration 1", "Iteration 2", "Iteration 3", "Rollout\n(Average)"]
cols = [CYAN, CYAN, CYAN, YELLOW]
for i, (label, col) in enumerate(zip(iters, cols)):
    x = 0.5 + i * 3.2
    box(sl, x, 3.6, 2.8, 1.3, CARD2, col)
    txt(sl, x+0.1, 3.65, 2.6, 0.25, label, 10, col, bold=True, align=PP_ALIGN.CENTER)
    # Mini bar chart effect
    heights = [0.4, 0.25, 0.15, 0.35, 0.2, 0.3, 0.1]
    for j, h in enumerate(heights):
        bx = x + 0.2 + j * 0.35
        by = 4.65 - h
        rect(sl, bx, by, 0.25, h, col)

    if i < 3:
        txt(sl, x + 2.85, 4.05, 0.3, 0.3, "+", 16, GREY, bold=True)

txt(sl, 0.5, 4.7, 12.3, 0.3,
    "Each iteration produces attention weights for all 21 slots → we average them → normalise so they sum to 1",
    9.5, GREY, align=PP_ALIGN.CENTER)

# Question answered
box(sl, 0.3, 5.3, 12.7, 1.9, CARD, GREEN)
txt(sl, 0.5, 5.4, 12.3, 0.3,
    "What Question Does Cross-Attention Rollout Answer?", 14, GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 5.85, 12.3, 0.9,
    "\"Which degradation steps are most influenced by the original soil/pile properties?\"\n\n"
    "A slot with high cross-attention is heavily shaped by the input features.\n"
    "A slot with low cross-attention relies more on what other slots communicated to it.\n"
    "This reveals whether early, middle, or late degradation steps are input-driven vs. context-driven.",
    11, LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 6 — Attention Rollout: Self-Attention
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)

txt(sl, 0.5, 0.25, 12.3, 0.5,
    "Attention Rollout: Self-Attention",
    22, ORANGE, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 0.75, 12.3, 0.35,
    "How do slots communicate with each other? Which degradation steps exchange information?",
    12, GREY, align=PP_ALIGN.CENTER)

# Visual: 21×21 interaction matrix concept
box(sl, 0.3, 1.3, 5.5, 3.0, CARD, ORANGE)
txt(sl, 0.5, 1.4, 5.1, 0.3, "Slot-to-Slot Communication", 14, ORANGE, bold=True)

txt(sl, 0.5, 1.8, 5.1, 0.9,
    "Self-attention creates a 21 × 21 matrix:\n"
    "each cell (i, j) tells us how much Slot i\n"
    "listens to Slot j.\n\n"
    "This reveals temporal dependencies in the\n"
    "degradation process.", 10, LIGHT)

# Mini heatmap
for i in range(7):
    for j in range(7):
        # Simulate attention pattern (diagonal-dominant)
        val = 0.8 if i == j else (0.4 if abs(i-j) == 1 else (0.15 if abs(i-j) == 2 else 0.05))
        r = int(255 * val)
        g = int(132 * val)
        b = int(252 * val)
        col_hex = f"{r:02X}{g:02X}{b:02X}"
        rect(sl, 3.5 + j * 0.3, 1.9 + i * 0.3, 0.28, 0.28, col_hex, "333355")

txt(sl, 3.5, 1.65, 2.1, 0.2, "  S1   S2   S3   S4   S5  S6  S7", 6.5, CYAN, font="Consolas")
for i, s in enumerate(["S1","S2","S3","S4","S5","S6","S7"]):
    txt(sl, 3.15, 1.93 + i * 0.3, 0.35, 0.2, s, 6.5, CYAN, font="Consolas")

txt(sl, 3.5, 4.05, 2.1, 0.2, "  Brighter = stronger attention", 7, GREY, align=PP_ALIGN.CENTER)

# Residual connection explanation
box(sl, 6.1, 1.3, 6.8, 3.0, CARD, PURPLE)
txt(sl, 6.3, 1.4, 6.4, 0.3, "Residual Mixing: The Key Innovation", 14, PURPLE, bold=True)

txt(sl, 6.3, 1.8, 6.4, 2.3,
    "In the real network, each slot keeps a copy of its own information\n"
    "(residual connection). We must account for this when rolling up\n"
    "across iterations. The formula:\n\n"
    "   Mixed Attention  =  50% × Identity  +  50% × Attention\n\n"
    "▸ The Identity part means: 'I keep half of my own information'\n"
    "▸ The Attention part means: 'I take half from what other slots say'\n\n"
    "Then we multiply across all 3 iterations:\n\n"
    "   Rollout  =  Mixed₃  ×  Mixed₂  ×  Mixed₁\n\n"
    "This matrix multiplication traces the TOTAL flow of information\n"
    "from source slots to destination slots across the entire network.", 10, LIGHT)

# Question answered
box(sl, 0.3, 4.5, 12.7, 2.7, CARD, CYAN)
txt(sl, 0.5, 4.6, 12.3, 0.3,
    "Why Separate Cross and Self Attention?", 15, CYAN, bold=True, align=PP_ALIGN.CENTER)

# Two columns
box(sl, 0.5, 5.0, 5.8, 2.0, CARD2, CYAN)
txt(sl, 0.7, 5.1, 5.4, 0.25, "Cross-Attention Rollout", 12, CYAN, bold=True)
txt(sl, 0.7, 5.4, 5.4, 1.4,
    "Shape: 21 values (one per slot)\n\n"
    "Answers: \"How strongly does each degradation\n"
    "step depend on the original soil/pile inputs?\"\n\n"
    "Interpretation: input-driven vs. context-driven behaviour",
    9.5, LIGHT)

box(sl, 6.7, 5.0, 6.1, 2.0, CARD2, ORANGE)
txt(sl, 6.9, 5.1, 5.7, 0.25, "Self-Attention Rollout", 12, ORANGE, bold=True)
txt(sl, 6.9, 5.4, 5.7, 1.4,
    "Shape: 21 × 21 matrix (slot-to-slot)\n\n"
    "Answers: \"Which degradation steps influence each\n"
    "other?  Does step 15 look at step 3 for guidance?\"\n\n"
    "Interpretation: temporal coupling in the degradation process",
    9.5, LIGHT)


# ═══════════════════════════════════════════════════════
# SLIDE 7 — Token Attribution Maps
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)

txt(sl, 0.5, 0.25, 12.3, 0.5,
    "Token Attribution Maps: The Complete Picture",
    22, ORANGE, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 0.75, 12.3, 0.35,
    "Combining LRP and Attention Rollout into an interactive visualisation",
    12, GREY, align=PP_ALIGN.CENTER)

# Heatmap concept
box(sl, 0.3, 1.3, 7.5, 3.5, CARD, ORANGE)
txt(sl, 0.5, 1.4, 7.1, 0.3, "The Attribution Heatmap", 15, ORANGE, bold=True)

txt(sl, 0.5, 1.8, 7.1, 0.4,
    "A colour-coded grid showing the relevance of each feature at each degradation step:",
    10, LIGHT)

# Column headers
txt(sl, 1.8, 2.3, 5.5, 0.2,
    "Step 1    Step 2    Step 3    Step 4    ...    Step 21", 8, CYAN, align=PP_ALIGN.CENTER)

# Feature rows with colored cells
feats_hm = [
    ("PI", [0.1, 0.08, 0.05, 0.12, 0.03, 0.07]),
    ("Gmax", [0.6, 0.3, 0.15, 0.2, 0.1, 0.25]),
    ("ν", [0.05, 0.02, 0.01, 0.03, 0.01, 0.02]),
    ("Dp", [0.08, 0.04, 0.02, 0.06, 0.02, 0.03]),
    ("Tp", [0.03, 0.01, 0.01, 0.02, 0.01, 0.02]),
    ("Lp", [0.15, 0.1, 0.05, 0.08, 0.04, 0.06]),
    ("Ip", [0.12, 0.06, 0.03, 0.05, 0.02, 0.04]),
    ("Dp/Lp", [1.0, 0.8, 0.6, 0.9, 0.5, 0.7]),
]
for i, (name, vals) in enumerate(feats_hm):
    y = 2.55 + i * 0.24
    txt(sl, 0.6, y, 0.8, 0.2, name, 8, GREEN, bold=True)
    for j, v in enumerate(vals):
        x = 1.5 + j * 0.95
        r_val = int(255 * v) if v > 0 else 0
        b_val = int(255 * (-v)) if v < 0 else 0
        c = f"{r_val:02X}{int(50*(1-v)):02X}{int(50*(1-v)):02X}" if v > 0 else "333355"
        rect(sl, x, y, 0.85, 0.2, c, "333355")

# Color legend
txt(sl, 1.5, 4.5, 5.0, 0.2,
    "■ Intense red = high positive relevance     ■ Light = low relevance", 8, GREY)

# Right side: reading the map
box(sl, 8.0, 1.3, 5.0, 3.5, CARD, GREEN)
txt(sl, 8.2, 1.4, 4.6, 0.3, "How to Read the Map", 14, GREEN, bold=True)
tf = txt(sl, 8.2, 1.8, 4.6, 0.3,
         "Each cell answers a specific question:", 10, LIGHT)
items = [
    ("Row = which feature", "PI, Gmax, ν, Dp, Tp, Lp, Ip, Dp/Lp"),
    ("Column = which step", "Step 1 (initial) through Step 21"),
    ("Colour = how important", "Red = pushes prediction up, Blue = pushes down"),
    ("Intensity = magnitude", "Darker = stronger influence"),
    ("", ""),
    ("Switchable by variable", "Show heatmap for KL, KR, or KLR"),
    ("Per-scenario", "Each test scenario generates its own attribution map"),
    ("Interactive", "Hover over any cell to see the exact relevance value"),
]
for item_title, item_desc in items:
    if item_title:
        bullet(tf, f"▸  {item_title}", 9, YELLOW, bold=True, sp=3)
        bullet(tf, f"    {item_desc}", 8, GREY, sp=0)
    else:
        bullet(tf, "", 4, GREY, sp=1)

# Bottom: the 3 dashboard panels
box(sl, 0.3, 5.0, 12.7, 2.2, CARD, PURPLE)
txt(sl, 0.5, 5.1, 12.3, 0.3,
    "The XAI Dashboard: Three Complementary Views", 14, PURPLE, bold=True, align=PP_ALIGN.CENTER)

panels = [
    ("Feature Attribution\nHeatmap", "8 features × 21 steps\nwith red/blue colour coding\nSwitchable: KL / KR / KLR", GREEN),
    ("Cross-Attention\nBar Charts", "21 bars per iteration (×3)\n+ averaged rollout chart\nTaller = more input-dependent", CYAN),
    ("Self-Attention\nHeatmaps", "21 × 21 matrix per iteration\n+ multiplicative rollout\nBrighter = stronger connection", ORANGE),
]
for i, (title, desc, col) in enumerate(panels):
    x = 0.6 + i * 4.3
    box(sl, x, 5.5, 3.8, 1.5, CARD2, col)
    txt(sl, x+0.15, 5.55, 3.5, 0.5, title, 12, col, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, x+0.15, 6.05, 3.5, 0.8, desc, 9, LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 8 — Integrated Gradients (Gold Standard)
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)

txt(sl, 0.5, 0.25, 12.3, 0.5,
    "Integrated Gradients: The Gold Standard",
    22, YELLOW, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 0.75, 12.3, 0.35,
    "A more rigorous version of LRP that guarantees perfect conservation",
    12, GREY, align=PP_ALIGN.CENTER)

# Analogy
box(sl, 0.3, 1.3, 6.2, 2.5, CARD, YELLOW)
txt(sl, 0.5, 1.4, 5.8, 0.3, "The Intuition: Walking the Path", 14, YELLOW, bold=True)
txt(sl, 0.5, 1.8, 5.8, 1.8,
    "Imagine a hiker walking from point A (all inputs = 0)\n"
    "to point B (actual input values).\n\n"
    "At each step of the walk, they measure the gradient:\n"
    "'which direction is uphill?'\n\n"
    "Integrated Gradients averages these measurements across\n"
    "the entire walk — not just at the destination.\n\n"
    "This gives a more complete picture than Gradient × Input,\n"
    "which only measures the gradient at point B.", 10, LIGHT)

# Formula
box(sl, 6.8, 1.3, 6.2, 2.5, CARD, CYAN)
txt(sl, 7.0, 1.4, 5.8, 0.3, "The Formula (Simplified)", 14, CYAN, bold=True)
txt(sl, 7.0, 1.9, 5.8, 0.5,
    "Attribution of feature j  =", 12, LIGHT, align=PP_ALIGN.CENTER)
txt(sl, 7.0, 2.3, 5.8, 0.5,
    "(xⱼ  −  0)  ×  average gradient along the path",
    14, YELLOW, bold=True, align=PP_ALIGN.CENTER, font="Cambria Math")
txt(sl, 7.0, 2.9, 5.8, 0.6,
    "We split the path into 30 equal steps, compute the\n"
    "gradient at each step, average them, then multiply\n"
    "by the input value. This guarantees:",
    10, LIGHT, align=PP_ALIGN.CENTER)
txt(sl, 7.0, 3.3, 5.8, 0.35,
    "Σ (all attributions)  =  f(x) − f(baseline)  exactly",
    12, GREEN, bold=True, align=PP_ALIGN.CENTER, font="Cambria Math")

# Comparison
box(sl, 0.3, 4.0, 12.7, 3.2, CARD, PURPLE)
txt(sl, 0.5, 4.1, 12.3, 0.3,
    "Gradient × Input vs Integrated Gradients", 14, PURPLE, bold=True, align=PP_ALIGN.CENTER)

# Visual comparison
box(sl, 0.6, 4.5, 5.7, 2.5, CARD2, GREEN)
txt(sl, 0.8, 4.6, 5.3, 0.25, "Gradient × Input  (Default)", 12, GREEN, bold=True, align=PP_ALIGN.CENTER)
tf = txt(sl, 0.8, 4.9, 5.3, 0.3, "", 1)
items_gi = [
    ("Speed", "~50 ms per scenario — instant"),
    ("Accuracy", "First-order approximation (very good)"),
    ("Conservation", "Approximate — sum ≈ prediction"),
    ("When to use", "Interactive dashboard (fast exploration)"),
    ("Gradient calls", "63 per scenario (3 vars × 21 steps)"),
]
for t, d in items_gi:
    bullet(tf, f"▸  {t}:  {d}", 9, LIGHT, sp=2)

box(sl, 6.7, 4.5, 6.1, 2.5, CARD2, YELLOW)
txt(sl, 6.9, 4.6, 5.7, 0.25, "Integrated Gradients  (Gold Standard)", 12, YELLOW, bold=True, align=PP_ALIGN.CENTER)
tf = txt(sl, 6.9, 4.9, 5.7, 0.3, "", 1)
items_ig = [
    ("Speed", "~2 seconds per scenario — slower"),
    ("Accuracy", "Exact path integral (provably correct)"),
    ("Conservation", "Exact — sum = prediction − baseline"),
    ("When to use", "Research publications, formal analysis"),
    ("Gradient calls", "1,890 per scenario (30× more)"),
]
for t, d in items_ig:
    bullet(tf, f"▸  {t}:  {d}", 9, LIGHT, sp=2)


# ═══════════════════════════════════════════════════════
# SLIDE 9 — The Complete M10 Pipeline
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)

txt(sl, 0.5, 0.25, 12.3, 0.5,
    "The Complete M10 Pipeline",
    22, PURPLE, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 0.75, 12.3, 0.35,
    "Same prediction model as M9 — with a new XAI observation layer (no retraining)",
    12, GREY, align=PP_ALIGN.CENTER)

# Top flow: M9 base
box(sl, 0.3, 1.2, 12.7, 2.0, CARD, CYAN)
txt(sl, 0.5, 1.25, 3.0, 0.3, "M9 Base (Unchanged)", 13, CYAN, bold=True)

# Pipeline steps
steps = [
    ("Soil & Pile\nParameters", "8 inputs", CYAN, 0.5),
    ("Embedding\nLayer", "8 → 64 dim", CYAN, 2.5),
    ("Cross-Attn\n(×3 iter.)", "input → slots", ORANGE, 4.5),
    ("Self-Attn\n(×3 iter.)", "slot ↔ slot", ORANGE, 6.5),
    ("SwiGLU\nMLP", "gated filtering", GREEN, 8.5),
    ("Prediction\nHeads", "slots → curves", YELLOW, 10.5),
]
for label, sub, col, x in steps:
    box(sl, x, 1.6, 1.7, 1.25, CARD2, col)
    txt(sl, x+0.05, 1.65, 1.6, 0.45, label, 9.5, col, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, x+0.05, 2.15, 1.6, 0.35, sub, 8, GREY, align=PP_ALIGN.CENTER)
    if x < 10.5:
        arrow_r(sl, x + 1.75, 2.0, 0.6, 0.12, GREY)

# Bottom: M10 XAI layer
box(sl, 0.3, 3.5, 12.7, 3.7, CARD, PURPLE)
txt(sl, 0.5, 3.55, 3.5, 0.3, "M10 XAI Layer (NEW)", 13, PURPLE, bold=True)
txt(sl, 4.0, 3.58, 5.0, 0.25,
    "— observes the model's internal state, does not modify it", 9, GREY)

# XAI methods
xai_methods = [
    ("Cross-Attention\nRollout", "Captures attention weights\nfrom all 3 iterations.\n\nAverages → normalises.\nResult: 21-value vector\n(one importance score\nper degradation slot)", CYAN, 0.5),
    ("Self-Attention\nRollout", "Captures 21×21 attention\nmatrices from each iteration.\n\nMixes with 50% identity\n(residual connection).\nMultiplies across iterations.\nResult: 21×21 flow matrix", ORANGE, 3.5),
    ("Gradient × Input\n(LRP)", "Computes the gradient\nof each output w.r.t.\neach input feature.\n\nMultiplies by input value.\nResult: 3 × 21 × 8 tensor\n(var × step × feature)", GREEN, 6.5),
    ("Feature\nImportance", "Sums absolute relevance\nacross all 21 steps\nfor each feature.\n\nNormalises to percentages.\nResult: 8 importance\nvalues per variable", YELLOW, 9.5),
]
for label, desc, col, x in xai_methods:
    box(sl, x, 3.95, 2.7, 3.0, CARD2, col)
    txt(sl, x+0.1, 4.0, 2.5, 0.4, label, 10, col, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, x+0.1, 4.45, 2.5, 2.3, desc, 8, LIGHT, align=PP_ALIGN.CENTER)

# Connection lines from M9 to XAI
for x_from, x_to in [(4.5, 1.5), (6.5, 4.3), (8.5, 7.5)]:
    arrow_d(sl, x_from + 0.5, 2.9, 0.12, 0.5, PURPLE)


# ═══════════════════════════════════════════════════════
# SLIDE 10 — Summary & Key Takeaways
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)

txt(sl, 0.5, 0.25, 12.3, 0.5,
    "Summary: M9 → M10",
    24, PURPLE, bold=True, align=PP_ALIGN.CENTER)

# Stats banner
stat(sl, 0.5, 0.9, 2.5, "MODEL RETRAINING", "None", GREEN)
stat(sl, 3.3, 0.9, 2.5, "NEW DEPENDENCIES", "Zero", GREEN)
stat(sl, 6.1, 0.9, 2.5, "MODEL ACCURACY", "Unchanged", GREEN)
stat(sl, 8.9, 0.9, 2.5, "XAI METHODS", "3 + Dashboard", PURPLE)

# Three columns
# Column 1: LRP
box(sl, 0.3, 1.9, 4.0, 3.5, CARD, GREEN)
txt(sl, 0.5, 2.0, 3.6, 0.3, "LRP / Gradient × Input", 13, GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 2.4, 3.6, 2.8,
    "What it answers:\n"
    "\"How much does each soil/pile\n"
    "parameter contribute to each\n"
    "degradation step prediction?\"\n\n"
    "Key findings:\n"
    "▸ Dp/Lp (slenderness) dominates\n"
    "   KL predictions at 73%\n"
    "▸ Gmax is 2nd most important\n"
    "▸ Attribution changes across steps\n"
    "▸ Different for KL vs KR vs KLR", 9.5, LIGHT, align=PP_ALIGN.CENTER)

# Column 2: Attention Rollout
box(sl, 4.6, 1.9, 4.2, 3.5, CARD, CYAN)
txt(sl, 4.8, 2.0, 3.8, 0.3, "Attention Rollout", 13, CYAN, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 4.8, 2.4, 3.8, 2.8,
    "What it answers:\n"
    "\"Where does information flow\n"
    "inside the model?\"\n\n"
    "Two separate views:\n"
    "▸ Cross: input → slot flow\n"
    "   (which steps use input data?)\n"
    "▸ Self: slot → slot flow\n"
    "   (which steps talk to each other?)\n"
    "▸ Residual mixing accounts for\n"
    "   skip connections in the network", 9.5, LIGHT, align=PP_ALIGN.CENTER)

# Column 3: Token Attribution Maps
box(sl, 9.1, 1.9, 3.9, 3.5, CARD, ORANGE)
txt(sl, 9.3, 2.0, 3.5, 0.3, "Token Attribution Maps", 13, ORANGE, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 9.3, 2.4, 3.5, 2.8,
    "What it provides:\n"
    "\"An interactive dashboard to\n"
    "explore explanations visually.\"\n\n"
    "Three linked panels:\n"
    "▸ Feature × Step heatmap\n"
    "   (red = positive, blue = negative)\n"
    "▸ Cross-attention bar charts\n"
    "   (one per iteration + rollout)\n"
    "▸ Self-attention 21×21 heatmaps\n"
    "▸ Per-scenario, per-variable", 9.5, LIGHT, align=PP_ALIGN.CENTER)

# Bottom message
box(sl, 0.3, 5.6, 12.7, 1.6, CARD, YELLOW)
txt(sl, 0.5, 5.7, 12.3, 0.3,
    "Why Does Explainability Matter?", 15, YELLOW, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 0.5, 6.1, 12.3, 0.9,
    "A black-box model predicts KL will drop by 15% — but can we trust it? Should engineers act on it?\n"
    "XAI transforms trust: 'KL drops because Dp/Lp = 0.25 (slender pile) and Gmax = 23 MPa (soft soil)\n"
    "dominate this scenario, and the early degradation steps are heavily input-driven.'\n"
    "This is the difference between a number and a decision support system.",
    11, LIGHT, align=PP_ALIGN.CENTER)


# ── Save ──
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "M10_XAI_Explained.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"  10 slides — no code, only concepts and visuals")
