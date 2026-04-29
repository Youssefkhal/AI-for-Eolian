"""
M10 XAI Explanation — Visual PPTX for Civil Engineering PhD Students
====================================================================
Diagram-heavy, practical presentation with minimal jargon.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, copy

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT = os.path.join(SCRIPT_DIR, "M10_XAI_Explanation_v2.pptx")

# ── Colours ──
BG       = RGBColor(0x0F, 0x10, 0x1A)
CARD     = RGBColor(0x1A, 0x1B, 0x2E)
CARD2    = RGBColor(0x22, 0x23, 0x3A)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xDD, 0xDD, 0xDD)
MUTED    = RGBColor(0x88, 0x88, 0x99)
PURPLE   = RGBColor(0xC0, 0x84, 0xFC)
CYAN     = RGBColor(0x00, 0xD2, 0xFF)
GREEN    = RGBColor(0x00, 0xFF, 0x88)
ORANGE   = RGBColor(0xFF, 0xA5, 0x00)
RED      = RGBColor(0xFF, 0x6B, 0x6B)
YELLOW   = RGBColor(0xFA, 0xCC, 0x15)
BLUE     = RGBColor(0x38, 0xBD, 0xF8)
TEAL     = RGBColor(0x2D, 0xD4, 0xBF)
PINK     = RGBColor(0xF4, 0x72, 0xB6)
DARK_PURPLE = RGBColor(0x3B, 0x1D, 0x6E)
DARK_CYAN   = RGBColor(0x0A, 0x3D, 0x50)
DARK_GREEN  = RGBColor(0x0A, 0x40, 0x2A)
DARK_ORANGE = RGBColor(0x50, 0x30, 0x0A)
DARK_RED    = RGBColor(0x50, 0x15, 0x15)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_bg(slide):
    bg = slide.background; f = bg.fill; f.solid(); f.fore_color.rgb = BG

def box(slide, l, t, w, h, fill=CARD, border=None, radius=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def rect(slide, l, t, w, h, fill=CARD, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1.2)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def circle(slide, l, t, d, fill=CARD, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, d, d)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def arrow_shape(slide, l, t, w, h, fill=CYAN):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False
    return s

def down_arrow(slide, l, t, w, h, fill=CYAN):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False
    return s

def chevron(slide, l, t, w, h, fill=CYAN):
    s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False
    return s

def txt(slide, l, t, w, h, text, sz=14, col=WHITE, bold=False,
        align=PP_ALIGN.LEFT, font="Segoe UI"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = col
    p.font.bold = bold; p.font.name = font; p.alignment = align
    return tf

def add_p(tf, text, sz=13, col=LIGHT, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    p = tf.add_paragraph(); p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = col
    p.font.bold = bold; p.font.name = font; p.alignment = align
    p.space_before = Pt(3); p.space_after = Pt(1)
    return p

def labeled_box(slide, l, t, w, h, label, sub="", fill=CARD, border=CYAN,
                label_col=WHITE, sub_col=MUTED, label_sz=12, sub_sz=9):
    """A coloured box with centred label text inside."""
    b = box(slide, l, t, w, h, fill=fill, border=border)
    tf = txt(slide, l, t + h * 0.15, w, Pt(label_sz + 4),
             label, sz=label_sz, col=label_col, bold=True, align=PP_ALIGN.CENTER)
    if sub:
        add_p(tf, sub, sz=sub_sz, col=sub_col, align=PP_ALIGN.CENTER)
    return b

def stripe(slide, color=PURPLE):
    rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill=color)

def title_block(slide, title, subtitle, color=CYAN):
    stripe(slide, color)
    txt(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5),
        title, sz=26, col=color, bold=True)
    txt(slide, Inches(0.5), Inches(0.78), Inches(12), Inches(0.35),
        subtitle, sz=13, col=MUTED)


# ════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
stripe(sl, PURPLE)

txt(sl, Inches(1), Inches(1.5), Inches(11), Inches(1),
    "M10: Explainable AI for Pile Stiffness Prediction",
    sz=34, col=WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(sl, Inches(1.5), Inches(2.6), Inches(10), Inches(0.5),
    "How does the model know which soil parameters matter?",
    sz=18, col=PURPLE, align=PP_ALIGN.CENTER)

txt(sl, Inches(2), Inches(3.3), Inches(9), Inches(0.4),
    "A practical guide for civil engineering researchers",
    sz=14, col=MUTED, align=PP_ALIGN.CENTER)

# Visual flow: Problem → Model → XAI
flow_y = Inches(4.5)
flow_h = Inches(2.0)
cards = [
    ("THE PROBLEM", "Typhoon loads degrade\nsoil stiffness around\nmonopile foundations", RED, DARK_RED),
    ("THE MODEL", "A Slot-Attention AI\npredicts KL, KR, KLR\ndegradation curves", CYAN, DARK_CYAN),
    ("THE XAI", "Three methods reveal\nWHY the model makes\neach prediction", GREEN, DARK_GREEN),
]
cw = Inches(3.3); gap = Inches(0.35)
sx = (W - 3*cw - 2*gap) // 2
for i, (title, desc, col, bg_col) in enumerate(cards):
    x = sx + i*(cw + gap)
    box(sl, x, flow_y, cw, flow_h, fill=bg_col, border=col)
    txt(sl, x + Inches(0.15), flow_y + Inches(0.2), cw - Inches(0.3), Inches(0.3),
        title, sz=14, col=col, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.15), flow_y + Inches(0.65), cw - Inches(0.3), Inches(1.1),
        desc, sz=12, col=LIGHT, align=PP_ALIGN.CENTER)
    if i < 2:
        arrow_shape(sl, x + cw + Inches(0.03), flow_y + flow_h//2 - Inches(0.15),
                    Inches(0.3), Inches(0.3), fill=MUTED)


# ════════════════════════════════════════════════════════════
# SLIDE 2 — THE BIG PICTURE: What M10 Does
# ════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
title_block(sl, "What Does M10 Do?",
            "From 8 soil/pile measurements to 21-step stiffness degradation curves + explanations")

# ── INPUT block ──
bx = Inches(0.4); by = Inches(1.5); bw = Inches(2.8); bh = Inches(5.2)
box(sl, bx, by, bw, bh, fill=DARK_CYAN, border=CYAN)
txt(sl, bx+Inches(0.1), by+Inches(0.1), bw-Inches(0.2), Inches(0.35),
    "INPUT: 8 Parameters", sz=14, col=CYAN, bold=True, align=PP_ALIGN.CENTER)

features = ["PI  (Plasticity Index)", "Gmax  (Max Shear Modulus)",
            "v  (Poisson's Ratio)", "Dp  (Pile Diameter)",
            "Tp  (Wall Thickness)", "Lp  (Pile Length)",
            "Ip  (Moment of Inertia)", "Dp/Lp  (Slenderness)"]
feat_colors = [RED, ORANGE, YELLOW, CYAN, BLUE, GREEN, TEAL, PINK]
for fi, (feat, fc) in enumerate(zip(features, feat_colors)):
    fy = by + Inches(0.55) + fi * Inches(0.56)
    box(sl, bx+Inches(0.15), fy, bw-Inches(0.3), Inches(0.45),
        fill=RGBColor(0x14, 0x14, 0x28), border=fc)
    txt(sl, bx+Inches(0.25), fy+Inches(0.07), bw-Inches(0.5), Inches(0.3),
        feat, sz=10, col=fc, bold=True, align=PP_ALIGN.LEFT)

# Arrow input → model
arrow_shape(sl, Inches(3.35), Inches(3.8), Inches(0.5), Inches(0.4), fill=CYAN)

# ── MODEL block ──
mx = Inches(4.0); my = Inches(1.5); mw = Inches(4.8); mh = Inches(5.2)
box(sl, mx, my, mw, mh, fill=DARK_PURPLE, border=PURPLE)
txt(sl, mx+Inches(0.1), my+Inches(0.1), mw-Inches(0.2), Inches(0.3),
    "M10: SwiGLU Slot-Attention Model", sz=13, col=PURPLE, bold=True, align=PP_ALIGN.CENTER)

# Model internals
model_steps = [
    ("Feature Embedding", "8 numbers  -->  64-dim vector", CYAN),
    ("21 Slots Initialised", "1 initial + 20 drop slots", ORANGE),
    ("Cross-Attention (x3)", "Slots query the input", YELLOW),
    ("Self-Attention (x3)", "Slots share information", GREEN),
    ("SwiGLU MLP (x3)", "Each slot refines itself", TEAL),
    ("Prediction Heads", "Slot --> KL, KR, KLR values", PINK),
]
for si, (step, desc, sc) in enumerate(model_steps):
    sy = my + Inches(0.5) + si * Inches(0.75)
    box(sl, mx+Inches(0.15), sy, mw-Inches(0.3), Inches(0.62),
        fill=RGBColor(0x18, 0x18, 0x30), border=sc)
    txt(sl, mx+Inches(0.3), sy+Inches(0.05), Inches(2.5), Inches(0.25),
        step, sz=11, col=sc, bold=True)
    txt(sl, mx+Inches(0.3), sy+Inches(0.3), mw-Inches(0.6), Inches(0.25),
        desc, sz=9, col=MUTED)
    if si < 5:
        down_arrow(sl, mx+mw//2-Inches(0.1), sy+Inches(0.6),
                   Inches(0.2), Inches(0.17), fill=RGBColor(0x44,0x44,0x66))

# Arrow model → output
arrow_shape(sl, Inches(8.95), Inches(3.8), Inches(0.5), Inches(0.4), fill=PURPLE)

# ── OUTPUT block ──
ox = Inches(9.6); oy = Inches(1.5); ow = Inches(3.3); oh = Inches(5.2)
box(sl, ox, oy, ow, oh, fill=DARK_GREEN, border=GREEN)
txt(sl, ox+Inches(0.1), oy+Inches(0.1), ow-Inches(0.2), Inches(0.3),
    "OUTPUT: Stiffness Curves", sz=13, col=GREEN, bold=True, align=PP_ALIGN.CENTER)

out_vars = [
    ("KL", "Lateral Stiffness", "21 values showing\nhow KL drops with\ncyclic loading", CYAN),
    ("KR", "Rotational Stiffness", "21 values showing\nhow KR drops with\ncyclic loading", ORANGE),
    ("KLR", "Coupled Stiffness", "21 values showing\nhow KLR changes\nwith cyclic loading", GREEN),
]
for vi, (vn, vlab, vdesc, vc) in enumerate(out_vars):
    vy = oy + Inches(0.5) + vi * Inches(1.5)
    box(sl, ox+Inches(0.15), vy, ow-Inches(0.3), Inches(1.35),
        fill=RGBColor(0x14, 0x14, 0x28), border=vc)
    txt(sl, ox+Inches(0.25), vy+Inches(0.08), Inches(0.6), Inches(0.3),
        vn, sz=16, col=vc, bold=True)
    txt(sl, ox+Inches(0.85), vy+Inches(0.08), Inches(1.8), Inches(0.25),
        vlab, sz=10, col=LIGHT, bold=True)
    txt(sl, ox+Inches(0.25), vy+Inches(0.45), ow-Inches(0.5), Inches(0.8),
        vdesc, sz=9, col=MUTED)


# ════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT ARE SLOTS? (Visual diagram)
# ════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
title_block(sl, "What Are Slots?",
            "The 21 internal states that represent different stages of stiffness degradation", ORANGE)

# Left — explanation
lx = Inches(0.4); ly = Inches(1.5); lw = Inches(5.5); lh = Inches(5.3)
box(sl, lx, ly, lw, lh, fill=CARD, border=RGBColor(0x33,0x33,0x55))
tf = txt(sl, lx+Inches(0.3), ly+Inches(0.2), lw-Inches(0.6), Inches(0.3),
         "Slots = Internal Representation States", sz=15, col=ORANGE, bold=True)

add_p(tf, "", sz=4);
add_p(tf, "The model uses 21 slots to represent the pile's condition:", sz=12, col=LIGHT)
add_p(tf, "", sz=4);
add_p(tf, "Slot 1  =  Initial state (before any loading)", sz=12, col=CYAN, bold=True)
add_p(tf, "     Captures the pile's original stiffness (KL0, KR0, KLR0)", sz=11, col=LIGHT)
add_p(tf, "", sz=4);
add_p(tf, "Slots 2-21  =  Drop states (progressive degradation)", sz=12, col=RED, bold=True)
add_p(tf, "     Each slot represents one step of stiffness reduction", sz=11, col=LIGHT)
add_p(tf, "     Step 2 = small load  ...  Step 21 = extreme loading", sz=11, col=LIGHT)
add_p(tf, "", sz=6);
add_p(tf, "Key Design: Physics-Consistent Constraints", sz=13, col=GREEN, bold=True)
add_p(tf, "     KL and KR can only decrease (soil gets weaker)", sz=11, col=LIGHT)
add_p(tf, "     KLR magnitude can only increase (coupling grows)", sz=11, col=LIGHT)
add_p(tf, "     This is enforced in the model architecture itself", sz=11, col=LIGHT)
add_p(tf, "", sz=6);
add_p(tf, "Prototype Compression:", sz=13, col=PURPLE, bold=True)
add_p(tf, "     20 drop slots are built from only 4 prototypes", sz=11, col=LIGHT)
add_p(tf, "     Reduces parameters while preserving accuracy", sz=11, col=LIGHT)

# Right — visual diagram of slots
rx = Inches(6.2); ry = Inches(1.5)
# Draw 21 slot circles in a structured layout
txt(sl, rx+Inches(0.2), ry+Inches(0.05), Inches(6.5), Inches(0.35),
    "21 Slots - Visual Layout", sz=14, col=YELLOW, bold=True, align=PP_ALIGN.CENTER)

# Slot 1 (Initial) - big circle
s1x = rx + Inches(2.8); s1y = ry + Inches(0.6)
circle(sl, s1x, s1y, Inches(0.9), fill=DARK_CYAN, border=CYAN)
txt(sl, s1x, s1y+Inches(0.15), Inches(0.9), Inches(0.3),
    "S1", sz=16, col=CYAN, bold=True, align=PP_ALIGN.CENTER)
txt(sl, s1x, s1y+Inches(0.45), Inches(0.9), Inches(0.3),
    "Initial", sz=9, col=LIGHT, align=PP_ALIGN.CENTER)

# Arrow down from S1
down_arrow(sl, s1x+Inches(0.35), s1y+Inches(0.9), Inches(0.2), Inches(0.3), fill=MUTED)

# Drop slots in a grid (4 rows x 5 cols)
txt(sl, rx+Inches(0.2), ry+Inches(1.9), Inches(6.5), Inches(0.25),
    "Drop Slots (Degradation Steps)", sz=11, col=RED, bold=True, align=PP_ALIGN.CENTER)

proto_colors = [RED, ORANGE, YELLOW, GREEN]
proto_names = ["Proto A", "Proto B", "Proto C", "Proto D"]
slot_idx = 2
for row in range(4):
    pc = proto_colors[row]
    dark_pc = [DARK_RED, DARK_ORANGE, RGBColor(0x45,0x40,0x0A), DARK_GREEN][row]
    for col in range(5):
        if slot_idx > 21: break
        cx = rx + Inches(0.5) + col * Inches(1.25)
        cy = ry + Inches(2.25) + row * Inches(0.85)
        circle(sl, cx, cy, Inches(0.65), fill=dark_pc, border=pc)
        txt(sl, cx, cy+Inches(0.12), Inches(0.65), Inches(0.22),
            f"S{slot_idx}", sz=11, col=pc, bold=True, align=PP_ALIGN.CENTER)
        txt(sl, cx, cy+Inches(0.35), Inches(0.65), Inches(0.2),
            f"Drop {slot_idx-1}", sz=7, col=MUTED, align=PP_ALIGN.CENTER)
        slot_idx += 1

# Prototype labels
for row, (pn, pc) in enumerate(zip(proto_names, proto_colors)):
    py = ry + Inches(2.35) + row * Inches(0.85)
    txt(sl, rx - Inches(0.05), py+Inches(0.1), Inches(0.6), Inches(0.3),
        pn, sz=8, col=pc, bold=True, align=PP_ALIGN.RIGHT)

# Legend
txt(sl, rx + Inches(0.3), ry + Inches(5.7), Inches(6), Inches(0.25),
    "Each colour = one prototype cluster.  20 drop slots are built from just 4 learned prototypes.",
    sz=10, col=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 4 — THE 3 ITERATIONS (Visual cycle)
# ════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
title_block(sl, "The 3 Iteration Cycle",
            "The model refines its predictions by repeating three operations three times", CYAN)

# Three iteration columns
iter_w = Inches(3.8); iter_gap = Inches(0.35)
iter_sx = (W - 3*iter_w - 2*iter_gap) // 2

for it in range(3):
    ix = iter_sx + it*(iter_w + iter_gap)
    iy = Inches(1.5)

    # Iteration header
    box(sl, ix, iy, iter_w, Inches(0.55),
        fill=[RGBColor(0x1A,0x2A,0x3D), RGBColor(0x20,0x2A,0x35), RGBColor(0x25,0x2A,0x30)][it],
        border=CYAN)
    txt(sl, ix+Inches(0.1), iy+Inches(0.1), iter_w-Inches(0.2), Inches(0.35),
        f"Iteration {it+1}", sz=16, col=CYAN, bold=True, align=PP_ALIGN.CENTER)

    # Three operations
    ops = [
        ("CROSS-ATTENTION", "Slots read the input data",
         "Each slot queries the 8 soil/pile\nparameters to extract relevant info",
         YELLOW, RGBColor(0x30,0x2A,0x0A)),
        ("SELF-ATTENTION", "Slots communicate with each other",
         "Slots share information to capture\ninteractions (e.g. PI affects Gmax effect)",
         GREEN, DARK_GREEN),
        ("SwiGLU MLP", "Each slot processes independently",
         "Non-linear transformation refines\neach slot's internal representation",
         PURPLE, DARK_PURPLE),
    ]
    for oi, (op_name, op_what, op_how, op_col, op_bg) in enumerate(ops):
        oy = iy + Inches(0.7) + oi * Inches(1.65)
        box(sl, ix+Inches(0.1), oy, iter_w-Inches(0.2), Inches(1.45),
            fill=op_bg, border=op_col)
        txt(sl, ix+Inches(0.2), oy+Inches(0.08), iter_w-Inches(0.4), Inches(0.25),
            op_name, sz=12, col=op_col, bold=True, align=PP_ALIGN.CENTER)
        txt(sl, ix+Inches(0.2), oy+Inches(0.38), iter_w-Inches(0.4), Inches(0.25),
            op_what, sz=11, col=WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt(sl, ix+Inches(0.2), oy+Inches(0.7), iter_w-Inches(0.4), Inches(0.6),
            op_how, sz=9, col=MUTED, align=PP_ALIGN.CENTER)

        # Down arrow between ops
        if oi < 2:
            down_arrow(sl, ix+iter_w//2-Inches(0.08), oy+Inches(1.43),
                       Inches(0.16), Inches(0.22), fill=RGBColor(0x44,0x44,0x66))

    # Arrow between iterations
    if it < 2:
        arrow_shape(sl, ix+iter_w+Inches(0.02), Inches(3.5),
                    Inches(0.32), Inches(0.3), fill=MUTED)

# Bottom note
box(sl, Inches(1), Inches(6.5), Inches(11.3), Inches(0.7),
    fill=RGBColor(0x14,0x14,0x28), border=RGBColor(0x33,0x33,0x55))
txt(sl, Inches(1.2), Inches(6.55), Inches(10.9), Inches(0.6),
    "WHY 3 ITERATIONS?  Each pass refines the slot representations.  "
    "Iteration 1 captures basic features.  Iteration 2 captures interactions.  "
    "Iteration 3 fine-tunes for accurate stiffness degradation curves.",
    sz=11, col=LIGHT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 5 — WHY XAI? (The motivation)
# ════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
title_block(sl, "Why Do We Need Explainability (XAI)?",
            "AI models are powerful but opaque -- civil engineers need to trust and understand predictions", RED)

# Two columns: Without XAI vs With XAI
col_w = Inches(5.8); col_gap = Inches(0.5)
col_sx = (W - 2*col_w - col_gap) // 2

# WITHOUT XAI
bx = col_sx; by = Inches(1.5)
box(sl, bx, by, col_w, Inches(5.3), fill=DARK_RED, border=RED)
txt(sl, bx+Inches(0.2), by+Inches(0.15), col_w-Inches(0.4), Inches(0.35),
    "WITHOUT XAI", sz=18, col=RED, bold=True, align=PP_ALIGN.CENTER)

# Black box illustration
box(sl, bx+Inches(0.6), by+Inches(0.8), Inches(1.6), Inches(0.8),
    fill=DARK_CYAN, border=CYAN)
txt(sl, bx+Inches(0.7), by+Inches(0.95), Inches(1.4), Inches(0.3),
    "8 Soil/Pile\nParameters", sz=10, col=CYAN, align=PP_ALIGN.CENTER)

arrow_shape(sl, bx+Inches(2.3), by+Inches(1.0), Inches(0.4), Inches(0.3), fill=MUTED)

box(sl, bx+Inches(2.8), by+Inches(0.7), Inches(1.0), Inches(1.0),
    fill=RGBColor(0x10,0x10,0x10), border=RGBColor(0x44,0x44,0x44))
txt(sl, bx+Inches(2.8), by+Inches(0.95), Inches(1.0), Inches(0.4),
    "  ???\nBLACK\n  BOX", sz=10, col=RGBColor(0x66,0x66,0x66), bold=True, align=PP_ALIGN.CENTER)

arrow_shape(sl, bx+Inches(3.9), by+Inches(1.0), Inches(0.4), Inches(0.3), fill=MUTED)

box(sl, bx+Inches(4.4), by+Inches(0.8), Inches(1.1), Inches(0.8),
    fill=DARK_GREEN, border=GREEN)
txt(sl, bx+Inches(4.45), by+Inches(0.95), Inches(1.0), Inches(0.3),
    "KL, KR, KLR\nCurves", sz=10, col=GREEN, align=PP_ALIGN.CENTER)

# Problems list
problems = [
    "Cannot verify if predictions are physically reasonable",
    "No way to identify which soil parameter drives the result",
    "Cannot detect when the model is extrapolating dangerously",
    "Regulators and reviewers won't accept unexplained AI",
    "Engineers cannot build intuition from black-box outputs",
]
for pi, prob in enumerate(problems):
    py = by + Inches(2.1) + pi * Inches(0.58)
    txt(sl, bx+Inches(0.3), py, col_w-Inches(0.6), Inches(0.45),
        f"X   {prob}", sz=11, col=RED)

# WITH XAI
bx2 = col_sx + col_w + col_gap
box(sl, bx2, by, col_w, Inches(5.3), fill=DARK_GREEN, border=GREEN)
txt(sl, bx2+Inches(0.2), by+Inches(0.15), col_w-Inches(0.4), Inches(0.35),
    "WITH M10 XAI", sz=18, col=GREEN, bold=True, align=PP_ALIGN.CENTER)

# Glass box illustration
box(sl, bx2+Inches(0.6), by+Inches(0.8), Inches(1.6), Inches(0.8),
    fill=DARK_CYAN, border=CYAN)
txt(sl, bx2+Inches(0.7), by+Inches(0.95), Inches(1.4), Inches(0.3),
    "8 Soil/Pile\nParameters", sz=10, col=CYAN, align=PP_ALIGN.CENTER)

arrow_shape(sl, bx2+Inches(2.3), by+Inches(1.0), Inches(0.4), Inches(0.3), fill=MUTED)

box(sl, bx2+Inches(2.8), by+Inches(0.7), Inches(1.0), Inches(1.0),
    fill=RGBColor(0x10, 0x20, 0x10), border=GREEN)
txt(sl, bx2+Inches(2.8), by+Inches(0.85), Inches(1.0), Inches(0.6),
    "GLASS\n  BOX", sz=11, col=GREEN, bold=True, align=PP_ALIGN.CENTER)

arrow_shape(sl, bx2+Inches(3.9), by+Inches(1.0), Inches(0.4), Inches(0.3), fill=MUTED)

box(sl, bx2+Inches(4.4), by+Inches(0.8), Inches(1.1), Inches(0.8),
    fill=DARK_GREEN, border=GREEN)
txt(sl, bx2+Inches(4.45), by+Inches(0.95), Inches(1.0), Inches(0.3),
    "KL, KR, KLR\n+ WHY", sz=10, col=GREEN, align=PP_ALIGN.CENTER)

benefits = [
    "See which soil parameters control each stiffness variable",
    "Track how information flows through the model's internal states",
    "Verify that the model learned real soil mechanics relationships",
    "Build confidence for peer review and regulatory acceptance",
    "Identify edge cases where the model may be unreliable",
]
for pi, ben in enumerate(benefits):
    py = by + Inches(2.1) + pi * Inches(0.58)
    txt(sl, bx2+Inches(0.3), py, col_w-Inches(0.6), Inches(0.45),
        f"+   {ben}", sz=11, col=GREEN)


# ════════════════════════════════════════════════════════════
# SLIDE 6 — THREE XAI METHODS OVERVIEW (Visual)
# ════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
title_block(sl, "M10 Uses Three Complementary XAI Methods",
            "Each method answers a different question about the model's decision process", PURPLE)

methods = [
    ("Cross-Attention Rollout",
     "Which slots are most\nactive for this scenario?",
     "Tracks how much each of the\n21 slots reads from the input\nsoil parameters, averaged\nacross all 3 iterations.",
     "Output: 21 importance scores\n(one per slot, sum = 1.0)",
     CYAN, DARK_CYAN,
     "SLOT-LEVEL\nTRANSPARENCY"),

    ("Self-Attention Rollout",
     "How do the slots\ncommunicate internally?",
     "Tracks the information flow\nbetween all 21 slots across\n3 iterations, accounting for\nresidual connections.",
     "Output: 21x21 flow matrix\nshowing slot-to-slot influence",
     GREEN, DARK_GREEN,
     "INTERACTION\nTRANSPARENCY"),

    ("Gradient x Input\n(+ Integrated Gradients)",
     "Which of the 8 input features\nmatters most for the prediction?",
     "Computes how much each soil/pile\nparameter contributes to each\npredicted stiffness value at\neach degradation step.",
     "Output: importance score for\neach of the 8 input features",
     ORANGE, DARK_ORANGE,
     "FEATURE-LEVEL\nTRANSPARENCY"),
]

mw = Inches(3.8); mh = Inches(5.0); mgap = Inches(0.35)
msx = (W - 3*mw - 2*mgap) // 2

for i, (name, question, how, output, col, bgcol, level) in enumerate(methods):
    mx = msx + i*(mw + mgap)
    my = Inches(1.5)

    # Main card
    box(sl, mx, my, mw, mh, fill=bgcol, border=col)

    # Name
    txt(sl, mx+Inches(0.15), my+Inches(0.12), mw-Inches(0.3), Inches(0.5),
        name, sz=14, col=col, bold=True, align=PP_ALIGN.CENTER)

    # Question
    box(sl, mx+Inches(0.15), my+Inches(0.7), mw-Inches(0.3), Inches(0.7),
        fill=RGBColor(0x14,0x14,0x28), border=RGBColor(0x33,0x33,0x55))
    txt(sl, mx+Inches(0.2), my+Inches(0.6), mw-Inches(0.4), Inches(0.15),
        "QUESTION:", sz=8, col=MUTED, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, mx+Inches(0.2), my+Inches(0.82), mw-Inches(0.4), Inches(0.55),
        question, sz=11, col=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # How
    txt(sl, mx+Inches(0.2), my+Inches(1.5), mw-Inches(0.4), Inches(0.2),
        "HOW:", sz=9, col=MUTED, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, mx+Inches(0.15), my+Inches(1.75), mw-Inches(0.3), Inches(1.1),
        how, sz=10, col=LIGHT, align=PP_ALIGN.CENTER)

    # Output
    box(sl, mx+Inches(0.15), my+Inches(2.95), mw-Inches(0.3), Inches(0.7),
        fill=RGBColor(0x14,0x14,0x28), border=col)
    txt(sl, mx+Inches(0.2), my+Inches(3.0), mw-Inches(0.4), Inches(0.6),
        output, sz=10, col=col, align=PP_ALIGN.CENTER)

    # Level badge
    box(sl, mx+Inches(0.5), my+Inches(3.85), mw-Inches(1.0), Inches(0.65),
        fill=bgcol, border=col)
    txt(sl, mx+Inches(0.5), my+Inches(3.9), mw-Inches(1.0), Inches(0.55),
        level, sz=11, col=col, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 7 — Cross-Attention Rollout (Detailed Visual)
# ════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
title_block(sl, "Cross-Attention Rollout: How Slots Read the Input",
            "Visualising which of the 21 slots extract the most information from the soil parameters", CYAN)

# Diagram: Input → Cross Attention → Slots
# LEFT: Input box
box(sl, Inches(0.5), Inches(1.8), Inches(1.8), Inches(3.2),
    fill=DARK_CYAN, border=CYAN)
txt(sl, Inches(0.55), Inches(1.85), Inches(1.7), Inches(0.3),
    "INPUT", sz=13, col=CYAN, bold=True, align=PP_ALIGN.CENTER)
txt(sl, Inches(0.55), Inches(2.2), Inches(1.7), Inches(0.25),
    "EMBEDDING", sz=11, col=CYAN, align=PP_ALIGN.CENTER)

mini_feats = ["PI", "Gmax", "v", "Dp", "Tp", "Lp", "Ip", "D/L"]
for fi, mf in enumerate(mini_feats):
    fy = Inches(2.55) + fi * Inches(0.28)
    box(sl, Inches(0.7), fy, Inches(1.4), Inches(0.23),
        fill=RGBColor(0x14,0x14,0x28), border=feat_colors[fi])
    txt(sl, Inches(0.73), fy+Inches(0.01), Inches(1.3), Inches(0.2),
        mf, sz=8, col=feat_colors[fi], bold=True, align=PP_ALIGN.CENTER)

# Arrows from input to attention cloud
for ai in range(5):
    ay = Inches(2.6) + ai * Inches(0.5)
    arrow_shape(sl, Inches(2.4), ay, Inches(0.6), Inches(0.18),
                fill=RGBColor(0x00, 0x80 + ai*20, 0xFF))

# MIDDLE: Cross Attention cloud
box(sl, Inches(3.2), Inches(1.5), Inches(3.5), Inches(4.2),
    fill=RGBColor(0x18, 0x18, 0x30), border=YELLOW)
txt(sl, Inches(3.3), Inches(1.6), Inches(3.3), Inches(0.3),
    "CROSS-ATTENTION", sz=14, col=YELLOW, bold=True, align=PP_ALIGN.CENTER)
txt(sl, Inches(3.3), Inches(1.95), Inches(3.3), Inches(0.25),
    "x 3 Iterations", sz=11, col=MUTED, align=PP_ALIGN.CENTER)

# Show iterations as stacked layers
iter_labels = ["Iter. 1: Initial read", "Iter. 2: Refined read", "Iter. 3: Final read"]
for ii, il in enumerate(iter_labels):
    iy = Inches(2.4) + ii * Inches(0.95)
    box(sl, Inches(3.5), iy, Inches(3.0), Inches(0.75),
        fill=[RGBColor(0x2A,0x20,0x0A), RGBColor(0x2F,0x25,0x0A), RGBColor(0x35,0x2A,0x0A)][ii],
        border=YELLOW)
    txt(sl, Inches(3.6), iy+Inches(0.1), Inches(2.8), Inches(0.22),
        il, sz=10, col=YELLOW, bold=True)
    # Mini slot weights
    txt(sl, Inches(3.6), iy+Inches(0.35), Inches(2.8), Inches(0.3),
        "S1: 0.05  S2: 0.04  S3: 0.06  ...  S21: 0.05", sz=8, col=MUTED)

# Arrows to output
for ai in range(5):
    ay = Inches(2.6) + ai * Inches(0.5)
    arrow_shape(sl, Inches(6.85), ay, Inches(0.5), Inches(0.18), fill=YELLOW)

# RIGHT: Rollout result (bar chart visual)
box(sl, Inches(7.5), Inches(1.5), Inches(5.3), Inches(4.2),
    fill=CARD, border=RGBColor(0x33,0x33,0x55))
txt(sl, Inches(7.6), Inches(1.6), Inches(5.1), Inches(0.3),
    "ROLLOUT RESULT", sz=14, col=YELLOW, bold=True, align=PP_ALIGN.CENTER)
txt(sl, Inches(7.6), Inches(1.95), Inches(5.1), Inches(0.25),
    "Average across 3 iterations, normalised to sum = 1.0",
    sz=10, col=MUTED, align=PP_ALIGN.CENTER)

# Simulated bar chart
bar_vals = [0.08, 0.06, 0.07, 0.05, 0.04, 0.05, 0.06, 0.04, 0.05, 0.05,
            0.04, 0.05, 0.04, 0.03, 0.04, 0.05, 0.04, 0.05, 0.04, 0.03, 0.04]
max_bv = max(bar_vals)
bar_base_y = Inches(4.8)
bar_top_margin = Inches(2.5)
bar_w = Inches(0.2)
bar_gap = Inches(0.03)
bar_sx = Inches(7.7)

for bi, bv in enumerate(bar_vals):
    bh_px = bv/max_bv * 1.8  # max height in inches
    bx = bar_sx + bi * (bar_w + bar_gap)
    by = bar_base_y - Inches(bh_px)
    bc = CYAN if bi == 0 else [RED, ORANGE, YELLOW, GREEN][bi % 4] if bv > 0.05 else MUTED
    rect(sl, bx, by, bar_w, Inches(bh_px), fill=bc)
    txt(sl, bx - Inches(0.02), bar_base_y + Inches(0.02), Inches(0.24), Inches(0.2),
        f"S{bi+1}", sz=6, col=MUTED, align=PP_ALIGN.CENTER)

# Bottom explanation
box(sl, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.2),
    fill=RGBColor(0x14,0x14,0x28), border=RGBColor(0x33,0x33,0x55))
tf = txt(sl, Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.25),
         "What does this tell a civil engineer?", sz=13, col=CYAN, bold=True)
add_p(tf, "- High bar = that slot is actively reading the soil parameters to make its prediction", sz=11, col=LIGHT)
add_p(tf, "- Low bar = that slot relies more on information shared from other slots (self-attention)", sz=11, col=LIGHT)
add_p(tf, "- Slot 1 (initial stiffness) typically has high cross-attention because initial values depend directly on soil properties", sz=11, col=LIGHT)


# ════════════════════════════════════════════════════════════
# SLIDE 8 — Self-Attention Rollout (Detailed Visual)
# ════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
title_block(sl, "Self-Attention Rollout: How Slots Communicate",
            "Tracking information flow between the 21 internal states across all iterations", GREEN)

# LEFT: Explanation
box(sl, Inches(0.4), Inches(1.5), Inches(5.5), Inches(5.3),
    fill=CARD, border=RGBColor(0x33,0x33,0x55))

tf = txt(sl, Inches(0.7), Inches(1.65), Inches(5.0), Inches(0.3),
         "How It Works", sz=15, col=GREEN, bold=True)

add_p(tf, "", sz=4)
add_p(tf, "In each iteration, every slot can attend to every other slot.", sz=12, col=LIGHT)
add_p(tf, "This creates a 21x21 weight matrix per iteration.", sz=12, col=LIGHT)
add_p(tf, "", sz=6)
add_p(tf, "The Residual Connection:", sz=13, col=YELLOW, bold=True)
add_p(tf, "The model adds the slot's own previous state back in:", sz=11, col=LIGHT)
add_p(tf, "   new_slot = LayerNorm( old_slot + attention_output )", sz=11, col=YELLOW, font="Consolas")
add_p(tf, "This means ~50% of information is kept from before.", sz=11, col=LIGHT)
add_p(tf, "", sz=6)
add_p(tf, "The Rollout:", sz=13, col=YELLOW, bold=True)
add_p(tf, "We mix 50% identity + 50% attention at each iteration,", sz=11, col=LIGHT)
add_p(tf, "then multiply the three matrices together:", sz=11, col=LIGHT)
add_p(tf, "   R = (0.5*I + 0.5*A3) x (0.5*I + 0.5*A2) x (0.5*I + 0.5*A1)", sz=10, col=GREEN, font="Consolas")
add_p(tf, "", sz=6)
add_p(tf, "What civil engineers learn:", sz=13, col=CYAN, bold=True)
add_p(tf, "- Which degradation steps influence which other steps", sz=11, col=LIGHT)
add_p(tf, "- Whether early loading stages affect late-stage predictions", sz=11, col=LIGHT)
add_p(tf, "- The model's internal representation of progressive failure", sz=11, col=LIGHT)

# RIGHT: 21x21 heatmap visual
hm_x = Inches(6.2); hm_y = Inches(1.5)
box(sl, hm_x, hm_y, Inches(6.7), Inches(5.3),
    fill=CARD, border=RGBColor(0x33,0x33,0x55))
txt(sl, hm_x+Inches(0.2), hm_y+Inches(0.1), Inches(6.3), Inches(0.3),
    "Self-Attention Rollout Heatmap  [21 x 21]", sz=14, col=GREEN, bold=True, align=PP_ALIGN.CENTER)

txt(sl, hm_x+Inches(0.2), hm_y+Inches(0.45), Inches(6.3), Inches(0.2),
    "Brighter = stronger information flow from column slot to row slot",
    sz=10, col=MUTED, align=PP_ALIGN.CENTER)

# Draw a simplified heatmap (10x10 showing concept)
import random
random.seed(42)
grid_n = 10  # show 10x10 for readability
cell_sz = Inches(0.44)
grid_sx = hm_x + Inches(1.2)
grid_sy = hm_y + Inches(1.0)

# Row & column labels
for gi in range(grid_n):
    slot_lbl = f"S{gi*2+1}" if gi < 10 else "..."
    txt(sl, grid_sx - Inches(0.45), grid_sy + gi*cell_sz + Inches(0.08),
        Inches(0.4), Inches(0.25), slot_lbl, sz=8, col=GREEN, align=PP_ALIGN.RIGHT)
    txt(sl, grid_sx + gi*cell_sz + Inches(0.05), grid_sy - Inches(0.25),
        Inches(0.35), Inches(0.2), slot_lbl, sz=8, col=GREEN, align=PP_ALIGN.CENTER)

# Cells
for ri in range(grid_n):
    for ci in range(grid_n):
        # Diagonal is strongest (self-retention)
        if ri == ci:
            val = 0.4 + random.random() * 0.3
        elif abs(ri-ci) <= 1:
            val = 0.15 + random.random() * 0.2
        elif ri == 0 or ci == 0:
            val = 0.08 + random.random() * 0.12
        else:
            val = random.random() * 0.1

        intensity = int(min(val * 350, 255))
        cell_col = RGBColor(intensity//4, intensity, intensity//2)
        cx = grid_sx + ci * cell_sz
        cy = grid_sy + ri * cell_sz
        rect(sl, cx, cy, cell_sz - Inches(0.02), cell_sz - Inches(0.02),
             fill=cell_col)

# Labels
txt(sl, grid_sx + Inches(1.2), grid_sy + grid_n*cell_sz + Inches(0.08),
    Inches(2), Inches(0.2), "Source Slot (Key)", sz=10, col=MUTED, align=PP_ALIGN.CENTER)

# Annotations
box(sl, hm_x + Inches(0.3), hm_y + Inches(4.0), Inches(6.1), Inches(1.1),
    fill=RGBColor(0x14,0x14,0x28), border=GREEN)
tf = txt(sl, hm_x + Inches(0.45), hm_y + Inches(4.05), Inches(5.8), Inches(0.25),
         "Reading the Heatmap:", sz=12, col=GREEN, bold=True)
add_p(tf, "Bright diagonal = each slot retains its own information (residual connection)", sz=10, col=LIGHT)
add_p(tf, "Bright off-diagonal = strong communication between those slots", sz=10, col=LIGHT)
add_p(tf, "Bright first row/column = initial slot influences or is influenced by all others", sz=10, col=LIGHT)


# ════════════════════════════════════════════════════════════
# SLIDE 9 — Feature Attribution: Gradient x Input + IG
# ════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
title_block(sl, "Feature Attribution: Which Soil Parameters Matter Most?",
            "Gradient x Input (fast) and Integrated Gradients (precise) trace predictions back to the 8 inputs", ORANGE)

# LEFT: How it works
box(sl, Inches(0.4), Inches(1.5), Inches(6.2), Inches(2.8),
    fill=CARD, border=RGBColor(0x33,0x33,0x55))
tf = txt(sl, Inches(0.7), Inches(1.65), Inches(5.6), Inches(0.3),
         "How Feature Attribution Works", sz=15, col=ORANGE, bold=True)

add_p(tf, "", sz=4)
add_p(tf, "Both methods answer: 'If I change this soil parameter slightly,", sz=12, col=LIGHT)
add_p(tf, "how much does the predicted stiffness change?'", sz=12, col=LIGHT)
add_p(tf, "", sz=6)
add_p(tf, "Gradient x Input (Default -- ~50 ms):", sz=13, col=YELLOW, bold=True)
add_p(tf, "  Takes one snapshot: multiply each feature value by its gradient", sz=11, col=LIGHT)
add_p(tf, "", sz=4)
add_p(tf, "Integrated Gradients (Precise -- ~2 sec):", sz=13, col=GREEN, bold=True)
add_p(tf, "  Traces a path from 'no soil' to 'actual soil', averaging gradients", sz=11, col=LIGHT)
add_p(tf, "  Guarantees all attributions sum exactly to the predicted value", sz=11, col=LIGHT)

# RIGHT: Visual bar chart of feature importance
box(sl, Inches(6.9), Inches(1.5), Inches(6.0), Inches(2.8),
    fill=CARD, border=RGBColor(0x33,0x33,0x55))
txt(sl, Inches(7.1), Inches(1.6), Inches(5.6), Inches(0.3),
    "Typical Feature Importance Result", sz=14, col=YELLOW, bold=True, align=PP_ALIGN.CENTER)
txt(sl, Inches(7.1), Inches(1.95), Inches(5.6), Inches(0.2),
    "(% contribution to prediction)", sz=10, col=MUTED, align=PP_ALIGN.CENTER)

# Simulated feature importance bars
fi_data = [
    ("PI", 28.5, RED),
    ("Gmax", 24.2, ORANGE),
    ("Lp", 15.8, GREEN),
    ("Dp/Lp", 11.3, TEAL),
    ("Ip", 8.1, BLUE),
    ("Dp", 5.6, CYAN),
    ("Tp", 4.2, PURPLE),
    ("v", 2.3, PINK),
]
max_fi = fi_data[0][1]
for fi, (fname, fval, fcol) in enumerate(fi_data):
    fy = Inches(2.3) + fi * Inches(0.24)
    bar_w_pct = fval / max_fi * Inches(3.5)
    txt(sl, Inches(7.2), fy, Inches(0.7), Inches(0.2),
        fname, sz=10, col=fcol, bold=True, align=PP_ALIGN.RIGHT)
    rect(sl, Inches(8.0), fy + Inches(0.02), bar_w_pct, Inches(0.17), fill=fcol)
    txt(sl, Inches(8.05) + bar_w_pct, fy, Inches(0.8), Inches(0.2),
        f"{fval}%", sz=9, col=fcol, bold=True)

# BOTTOM: Heatmap explanation (the [3, 21, 8] output)
box(sl, Inches(0.4), Inches(4.6), Inches(12.5), Inches(2.6),
    fill=CARD, border=RGBColor(0x33,0x33,0x55))
txt(sl, Inches(0.7), Inches(4.7), Inches(12), Inches(0.3),
    "Feature Attribution Heatmap: [3 variables x 21 steps x 8 features]", sz=14, col=ORANGE, bold=True)

txt(sl, Inches(0.7), Inches(5.1), Inches(3.5), Inches(0.25),
    "For each stiffness variable:", sz=12, col=LIGHT)

# Three mini-heatmaps
var_names_c = [("KL (Lateral)", CYAN), ("KR (Rotational)", ORANGE), ("KLR (Coupled)", GREEN)]
mini_sx = Inches(0.7)

for vi, (vname, vcol) in enumerate(var_names_c):
    vx = mini_sx + vi * Inches(4.2)
    vy = Inches(5.4)

    txt(sl, vx, vy - Inches(0.05), Inches(3.8), Inches(0.2),
        vname, sz=11, col=vcol, bold=True, align=PP_ALIGN.CENTER)

    # Mini heatmap (8 rows x 10 cols showing concept)
    mini_feats2 = ["PI", "Gmax", "v", "Dp", "Tp", "Lp", "Ip", "D/L"]
    cell_w2 = Inches(0.29); cell_h2 = Inches(0.17)

    for fi2, mf2 in enumerate(mini_feats2):
        txt(sl, vx - Inches(0.03), vy + Inches(0.15) + fi2*cell_h2,
            Inches(0.5), cell_h2, mf2, sz=6, col=feat_colors[fi2], align=PP_ALIGN.RIGHT)
        for si2 in range(10):
            random.seed(vi*100 + fi2*10 + si2)
            v2 = random.random()
            # Make PI and Gmax more important
            if fi2 < 2: v2 = v2 * 0.6 + 0.4
            elif fi2 < 4: v2 = v2 * 0.4 + 0.1
            else: v2 = v2 * 0.2

            is_pos = random.random() > 0.3
            if is_pos:
                r_c = int(min(v2 * 300, 255))
                cc = RGBColor(r_c, r_c//4, r_c//4)
            else:
                b_c = int(min(v2 * 250, 255))
                cc = RGBColor(b_c//5, b_c//3, b_c)

            cx2 = vx + Inches(0.5) + si2 * cell_w2
            cy2 = vy + Inches(0.15) + fi2 * cell_h2
            rect(sl, cx2, cy2, cell_w2 - Inches(0.01), cell_h2 - Inches(0.01), fill=cc)

    # Step labels
    txt(sl, vx + Inches(0.5), vy + Inches(0.15) + 8*cell_h2 + Inches(0.01),
        Inches(2.9), Inches(0.15),
        "Step 1  -->  Step 21", sz=7, col=MUTED, align=PP_ALIGN.CENTER)

# Legend
txt(sl, Inches(0.7), Inches(6.9), Inches(4), Inches(0.2),
    "RED = positive influence (pushes stiffness up)      BLUE = negative influence (pushes stiffness down)",
    sz=9, col=MUTED)
txt(sl, Inches(7), Inches(6.9), Inches(5.5), Inches(0.2),
    "BRIGHTER = STRONGER influence on the prediction",
    sz=9, col=YELLOW, bold=True)


# ════════════════════════════════════════════════════════════
# SLIDE 10 — SUMMARY TABLE
# ════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
title_block(sl, "Summary: Three Perspectives on Model Transparency",
            "Each method reveals a different aspect of how the model makes stiffness predictions", PURPLE)

# Summary cards
summaries = [
    ("Cross-Attention\nRollout", CYAN, DARK_CYAN,
     "QUESTION",
     "Which slots are most\nactive for this scenario?",
     "WHAT YOU SEE",
     "A bar chart showing the\nimportance of each of the\n21 slots in reading the input",
     "CIVIL ENG. VALUE",
     "Identifies which degradation\nstages depend most directly\non soil properties"),

    ("Self-Attention\nRollout", GREEN, DARK_GREEN,
     "QUESTION",
     "How do degradation stages\ninfluence each other?",
     "WHAT YOU SEE",
     "A 21x21 grid (heatmap)\nshowing information flow\nbetween all slot pairs",
     "CIVIL ENG. VALUE",
     "Reveals internal coupling:\ne.g. early loading stages\naffecting late-stage stiffness"),

    ("Feature\nAttribution", ORANGE, DARK_ORANGE,
     "QUESTION",
     "Which of the 8 soil/pile\nparameters drives each prediction?",
     "WHAT YOU SEE",
     "Ranked bar chart of feature\nimportance + colour-coded\nheatmap across all steps",
     "CIVIL ENG. VALUE",
     "Directly answers: 'Is PI or\nGmax more important for\nthis specific pile scenario?'"),
]

sw = Inches(3.8); sh = Inches(5.0); sgap = Inches(0.35)
ssx = (W - 3*sw - 2*sgap) // 2

for i, (name, col, bgcol, q_lab, q_txt, s_lab, s_txt, v_lab, v_txt) in enumerate(summaries):
    sx = ssx + i*(sw + sgap)
    sy = Inches(1.5)

    box(sl, sx, sy, sw, sh, fill=bgcol, border=col)

    # Name
    txt(sl, sx+Inches(0.1), sy+Inches(0.1), sw-Inches(0.2), Inches(0.55),
        name, sz=15, col=col, bold=True, align=PP_ALIGN.CENTER)

    sections = [(q_lab, q_txt, Inches(0.7)), (s_lab, s_txt, Inches(1.75)), (v_lab, v_txt, Inches(2.9))]
    for lab, content, y_off in sections:
        box(sl, sx+Inches(0.12), sy+y_off, sw-Inches(0.24), Inches(0.95),
            fill=RGBColor(0x14,0x14,0x28), border=RGBColor(0x33,0x33,0x55))
        txt(sl, sx+Inches(0.2), sy+y_off+Inches(0.03), sw-Inches(0.4), Inches(0.18),
            lab, sz=8, col=MUTED, bold=True, align=PP_ALIGN.CENTER)
        txt(sl, sx+Inches(0.2), sy+y_off+Inches(0.22), sw-Inches(0.4), Inches(0.65),
            content, sz=10, col=LIGHT, align=PP_ALIGN.CENTER)

    # Speed badge
    speeds = ["< 1 ms", "< 1 ms", "50 ms / 2 sec"]
    box(sl, sx+Inches(0.8), sy+Inches(4.05), sw-Inches(1.6), Inches(0.4),
        fill=bgcol, border=col)
    txt(sl, sx+Inches(0.8), sy+Inches(4.07), sw-Inches(1.6), Inches(0.35),
        f"Speed: {speeds[i]}", sz=10, col=col, bold=True, align=PP_ALIGN.CENTER)

# Bottom message
box(sl, Inches(1.5), Inches(6.7), Inches(10.3), Inches(0.6),
    fill=RGBColor(0x18, 0x12, 0x30), border=PURPLE)
txt(sl, Inches(1.7), Inches(6.75), Inches(9.9), Inches(0.45),
    "All three methods are available in the M10 Web Dashboard.  Select any test scenario and the XAI results are computed in real-time.",
    sz=12, col=PURPLE, align=PP_ALIGN.CENTER)


# ── SAVE ──
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"10 slides, diagram-heavy, targeting civil engineering PhD students")
