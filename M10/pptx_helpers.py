from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0x0F,0x10,0x1A); CARD=RGBColor(0x1A,0x1B,0x2E)
WHITE=RGBColor(0xFF,0xFF,0xFF); LIGHT=RGBColor(0xDD,0xDD,0xDD)
MUTED=RGBColor(0x88,0x88,0x99); CYAN=RGBColor(0x00,0xD2,0xFF)
GREEN=RGBColor(0x00,0xFF,0x88); ORANGE=RGBColor(0xFF,0xA5,0x00)
RED=RGBColor(0xFF,0x6B,0x6B); PURPLE=RGBColor(0xC0,0x84,0xFC)
YELLOW=RGBColor(0xFA,0xCC,0x15); BLUE=RGBColor(0x38,0xBD,0xF8)
TEAL=RGBColor(0x2D,0xD4,0xBF); PINK=RGBColor(0xF4,0x72,0xB6)
GOLD=RGBColor(0xFF,0xD7,0x00)
DC=RGBColor(0x0A,0x3D,0x50); DG=RGBColor(0x0A,0x40,0x2A)
DO=RGBColor(0x50,0x30,0x0A); DR=RGBColor(0x50,0x15,0x15)
DP=RGBColor(0x3B,0x1D,0x6E); DK=RGBColor(0x14,0x14,0x28)

def mk():
    p=Presentation(); p.slide_width=Inches(13.333); p.slide_height=Inches(7.5); return p

def sl(p):
    s=p.slides.add_slide(p.slide_layouts[6])
    bg=s.background; f=bg.fill; f.solid(); f.fore_color.rgb=BG
    return s

def bx(s,l,t,w,h,fill=CARD,border=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if border: sh.line.color.rgb=border; sh.line.width=Pt(1.5)
    else: sh.line.fill.background()
    sh.shadow.inherit=False; return sh

def rc(s,l,t,w,h,fill=CARD,border=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if border: sh.line.color.rgb=border; sh.line.width=Pt(1.2)
    else: sh.line.fill.background()
    sh.shadow.inherit=False; return sh

def ar(s,l,t,w,h,fill=CYAN):
    sh=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    sh.line.fill.background(); sh.shadow.inherit=False; return sh

def da(s,l,t,w,h,fill=CYAN):
    sh=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    sh.line.fill.background(); sh.shadow.inherit=False; return sh

def tx(s,l,t,w,h,text,sz=14,col=WHITE,bold=False,align=PP_ALIGN.LEFT,font="Segoe UI"):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(sz); p.font.color.rgb=col
    p.font.bold=bold; p.font.name=font; p.alignment=align; return tf

def ap(tf,text,sz=13,col=LIGHT,bold=False,align=PP_ALIGN.LEFT,font="Segoe UI"):
    p=tf.add_paragraph(); p.text=text; p.font.size=Pt(sz); p.font.color.rgb=col
    p.font.bold=bold; p.font.name=font; p.alignment=align
    p.space_before=Pt(3); p.space_after=Pt(1); return p

def stripe(s,color=PURPLE):
    rc(s,Inches(0),Inches(0),Inches(13.333),Inches(0.06),fill=color)

def title(s,t,sub,color=CYAN):
    stripe(s,color)
    tx(s,Inches(0.5),Inches(0.25),Inches(12),Inches(0.5),t,sz=26,col=color,bold=True)
    tx(s,Inches(0.5),Inches(0.78),Inches(12),Inches(0.35),sub,sz=13,col=MUTED)
